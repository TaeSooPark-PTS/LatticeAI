"""
SQLite knowledge graph for Lattice AI workspace memory.

The graph keeps raw event JSON, normalized node metadata, and edges in one
portable database so it can later migrate to Neo4j/Postgres without changing
the ingestion contract.
"""

import asyncio
import hashlib
import json
import logging
import math
import os
import platform
import re
import shutil
import sqlite3
import time
import zipfile
from collections import Counter
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple

try:
    from kg_schema import KGStoreV2, NodeType, EdgeType, _exec_script
except Exception:  # pragma: no cover - v2 schema is optional at import time
    KGStoreV2 = None  # type: ignore[assignment]
    NodeType = None   # type: ignore[assignment]
    EdgeType = None   # type: ignore[assignment]
    _exec_script = None  # type: ignore[assignment]

from latticeai.core.local_embeddings import LocalEmbeddingModel

# Default read source for the graph queries: v2 reconstruction views.
# Override with LATTICEAI_KG_READ_V2=0 to fall back to the legacy tables.
_READ_FROM_V2_DEFAULT = os.getenv("LATTICEAI_KG_READ_V2", "1") != "0"

# Bump when the v2 projection layout changes (columns, normalization rules).
# On init, a stale projection is dropped and rebuilt from the authoritative
# legacy tables — safe because nodes_v2/edges_v2 only ever hold a derived view.
# v4: summary nullable + verbatim (byte-faithful) projection of legacy values.
_PROJECTION_VERSION = 4

_llm_router_ref = None

def set_llm_router(router_instance):
    global _llm_router_ref
    _llm_router_ref = router_instance


GRAPH_SCHEMA_VERSION = 1

LOCAL_TEXT_EXTENSIONS = {".txt", ".md"}
LOCAL_CODE_EXTENSIONS = {
    ".py", ".js", ".ts", ".tsx", ".jsx", ".html", ".css", ".json",
    ".yaml", ".yml", ".xml", ".sql", ".sh", ".zsh", ".toml", ".ini",
}
LOCAL_DOCUMENT_EXTENSIONS = {".pdf", ".docx"}
LOCAL_SPREADSHEET_EXTENSIONS = {".xlsx", ".csv"}
LOCAL_SLIDE_EXTENSIONS = {".pptx"}
LOCAL_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp"}
LOCAL_SUPPORTED_EXTENSIONS = (
    LOCAL_TEXT_EXTENSIONS
    | LOCAL_CODE_EXTENSIONS
    | LOCAL_DOCUMENT_EXTENSIONS
    | LOCAL_SPREADSHEET_EXTENSIONS
    | LOCAL_SLIDE_EXTENSIONS
    | LOCAL_IMAGE_EXTENSIONS
)

LOCAL_SIZE_LIMITS = {
    "text": 4_000_000,
    "code": 4_000_000,
    "pdf": 50_000_000,
    "document": 50_000_000,
    "spreadsheet": 50_000_000,
    "slide_deck": 50_000_000,
    "image": 100_000_000,
}

COMMON_EXCLUDED_DIRS = {
    ".git", "node_modules", ".venv", "venv", "env", "__pycache__",
    ".pytest_cache", ".mypy_cache", ".ruff_cache", ".next", ".nuxt",
    ".turbo", "dist", "build", "target", "out", "coverage", ".cache",
    ".config", ".ssh", ".gnupg", ".docker", ".kube", ".aws", ".azure",
    ".npm", ".pnpm-store", ".yarn", ".bun", ".cargo", ".rustup", ".pyenv",
    ".conda", ".local", ".claude", ".codex", ".cursor", ".copilot",
    ".antigravity", ".antigravity-ide",
}

COMMON_EXCLUDED_FILE_NAMES = {
    ".env", ".env.local", ".env.production", ".env.development",
    "id_rsa", "id_ed25519", "authorized_keys", "known_hosts",
    "credentials.json", "service-account.json", "token.json", "secrets.json",
    "cookies", "login data", "history", "web data", ".ds_store", "thumbs.db",
}
COMMON_EXCLUDED_FILE_SUFFIXES = {
    ".pem", ".key", ".p12", ".pfx", ".kdbx", ".wallet", ".sqlite", ".db",
    ".exe", ".dll", ".sys", ".msi", ".dmg", ".pkg", ".app", ".zip", ".tar",
    ".gz", ".7z", ".rar", ".mp4", ".mov", ".mp3", ".wav", ".tmp", ".bak",
    ".lock",
}
SENSITIVE_PATH_KEYWORDS = {
    "secret", "secrets", "token", "password", "passwd", "credential",
    "credentials", "private", "key", "wallet", "recovery", "seed",
    "mnemonic", "cookie", "session", "auth", "oauth", "certificate",
    "cert", "api_key", "apikey",
}

MACOS_EXCLUDED_PREFIXES = (
    "/System", "/Library", "/Applications", "/private", "/tmp", "/var",
)
WINDOWS_EXCLUDED_NAMES = {
    "windows", "program files", "program files (x86)", "programdata", "appdata",
    "$recycle.bin", "system volume information", "recovery", "perflogs",
    "intel", "amd", "nvidia",
}
LINUX_EXCLUDED_PREFIXES = (
    "/bin", "/boot", "/dev", "/etc", "/lib", "/lib64", "/proc", "/root",
    "/run", "/sbin", "/sys", "/tmp", "/usr", "/var", "/snap", "/lost+found",
)


def _now() -> str:
    return datetime.now().isoformat()


def _parse_iso(raw: Optional[str]) -> Optional[datetime]:
    if not raw:
        return None
    try:
        return datetime.fromisoformat(str(raw))
    except (TypeError, ValueError):
        return None


def _recency_score(updated_at: Optional[str], *, now: Optional[datetime] = None, half_life_days: float = 14.0) -> float:
    stamp = _parse_iso(updated_at)
    if not stamp:
        return 0.0
    now = now or datetime.now()
    age_days = max(0.0, (now - stamp).total_seconds() / 86400.0)
    decay = math.log(2) / max(0.1, half_life_days)
    return math.exp(-decay * age_days)


def _json(data: Optional[Dict[str, Any]]) -> str:
    return json.dumps(data or {}, ensure_ascii=False, sort_keys=True)


def _safe_loads(raw: Optional[str]) -> Dict[str, Any]:
    """Tolerantly parse a metadata_json column — returns {} on corrupt rows."""
    if not raw:
        return {}
    try:
        value = json.loads(raw)
        return value if isinstance(value, dict) else {}
    except (json.JSONDecodeError, TypeError) as e:
        logging.warning("knowledge_graph: corrupt metadata_json (%s) — using empty dict", e)
        return {}


def _slug(text: str, max_len: int = 96) -> str:
    value = re.sub(r"\s+", " ", str(text or "")).strip().lower()
    value = re.sub(r"[^0-9a-zA-Z가-힣._:@/-]+", "-", value).strip("-")
    return (value or "untitled")[:max_len]


def _sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _sha256_text(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8", errors="replace")).hexdigest()


def _safe_iso_from_stat_mtime(mtime: float) -> str:
    try:
        return datetime.fromtimestamp(float(mtime)).isoformat()
    except (TypeError, ValueError, OSError):
        return ""


def _path_fingerprint(path: Path) -> str:
    return _sha256_text(str(path.expanduser().resolve()))[:24]


def _is_relative_to(path: Path, base: Path) -> bool:
    try:
        path.relative_to(base)
        return True
    except ValueError:
        return False


def _path_parts_lower(path: Path) -> List[str]:
    return [part.lower() for part in path.parts if part and part not in {os.sep, path.anchor}]


def _current_os_type() -> str:
    system = platform.system().lower()
    if system.startswith("darwin"):
        return "macos"
    if system.startswith("windows"):
        return "windows"
    if system.startswith("linux"):
        return "linux"
    return system or "unknown"


def _drive_id_for_path(path: Path) -> str:
    resolved = path.expanduser().resolve()
    if resolved.drive:
        return resolved.drive.upper()
    parts = resolved.parts
    if len(parts) >= 3 and parts[1] == "Volumes":
        return f"/Volumes/{parts[2]}"
    if len(parts) >= 3 and parts[1] == "media":
        return f"/media/{parts[2]}"
    if len(parts) >= 3 and parts[1] == "mnt":
        return f"/mnt/{parts[2]}"
    return resolved.anchor or "/"


def _file_category(ext: str) -> str:
    ext = (ext or "").lower()
    if ext in LOCAL_CODE_EXTENSIONS:
        return "code"
    if ext in LOCAL_TEXT_EXTENSIONS:
        return "text"
    if ext == ".pdf":
        return "pdf"
    if ext in LOCAL_DOCUMENT_EXTENSIONS:
        return "document"
    if ext in LOCAL_SPREADSHEET_EXTENSIONS:
        return "spreadsheet"
    if ext in LOCAL_SLIDE_EXTENSIONS:
        return "slide_deck"
    if ext in LOCAL_IMAGE_EXTENSIONS:
        return "image"
    return "unsupported"


def _node_type_for_category(category: str) -> str:
    return {
        "code": "CodeFile",
        "spreadsheet": "Spreadsheet",
        "slide_deck": "SlideDeck",
        "image": "Image",
        "unsupported": "File",
    }.get(category, "Document")


def _parser_type_for_category(category: str, ext: str) -> str:
    if category in {"text", "code"}:
        return "plain_text"
    if category == "spreadsheet" and ext == ".csv":
        return "csv_text"
    if category == "image":
        return "image_ocr"
    return ext.lstrip(".") or category


def _size_limit_for_category(category: str) -> int:
    return LOCAL_SIZE_LIMITS.get(category, LOCAL_SIZE_LIMITS["document"])


def _is_hidden_path(path: Path, root: Optional[Path] = None) -> bool:
    parts: Iterable[str]
    if root is not None:
        try:
            parts = path.relative_to(root).parts
        except ValueError:
            parts = path.parts
    else:
        parts = path.parts
    return any(part.startswith(".") and part not in {".", ".."} for part in parts)


def _excluded_directory_reason(path: Path, *, root: Optional[Path] = None, os_type: Optional[str] = None) -> Optional[str]:
    os_type = os_type or _current_os_type()
    name = path.name.lower()
    if name in COMMON_EXCLUDED_DIRS:
        return "excluded_folder"
    if _is_hidden_path(path, root):
        return "hidden_folder"
    parts = _path_parts_lower(path)
    if os_type == "windows" and any(part in WINDOWS_EXCLUDED_NAMES for part in parts):
        return "system_folder"
    normalized = path.as_posix()
    root_normalized = root.as_posix() if root else ""

    def _prefix_blocks(prefixes: Tuple[str, ...]) -> bool:
        for prefix in prefixes:
            path_under_prefix = normalized == prefix or normalized.startswith(f"{prefix}/")
            root_under_prefix = bool(root_normalized) and (
                root_normalized == prefix or root_normalized.startswith(f"{prefix}/")
            )
            if path_under_prefix and not root_under_prefix:
                return True
        return False

    if os_type == "macos":
        home_library = Path.home() / "Library"
        try:
            root_is_library = bool(root) and _is_relative_to(root.expanduser().resolve(), home_library.expanduser().resolve())
            if _is_relative_to(path.expanduser().resolve(), home_library.expanduser().resolve()) and not root_is_library:
                return "user_library"
        except OSError:
            pass
        if _prefix_blocks(MACOS_EXCLUDED_PREFIXES):
            return "system_folder"
    if os_type == "linux":
        if _prefix_blocks(LINUX_EXCLUDED_PREFIXES):
            return "system_folder"
    return None


def _sensitive_file_reason(path: Path, *, root: Optional[Path] = None) -> Optional[str]:
    name = path.name.lower()
    suffix = path.suffix.lower()
    if name in COMMON_EXCLUDED_FILE_NAMES or suffix in COMMON_EXCLUDED_FILE_SUFFIXES:
        return "sensitive_or_excluded_file"
    try:
        rel_text = path.relative_to(root).as_posix().lower() if root else path.as_posix().lower()
    except ValueError:
        rel_text = path.as_posix().lower()
    tokens = re.split(r"[^0-9a-zA-Z_가-힣]+", rel_text)
    if any(token in SENSITIVE_PATH_KEYWORDS for token in tokens):
        return "sensitive_name"
    return None


def _root_warning(path: Path, os_type: str) -> Optional[str]:
    resolved = path.expanduser().resolve()
    home = Path.home().expanduser().resolve()
    if os_type == "macos" and resolved == home:
        return "홈 전체에는 설정/숨김 폴더가 포함될 수 있습니다. 문서, 데스크탑, 다운로드, 프로젝트 폴더부터 추가하는 것을 권장합니다."
    if os_type == "linux" and resolved.as_posix() == "/":
        return "루트 디렉터리에는 시스템 파일이 포함되어 있습니다. 일반 사용자 폴더나 마운트된 데이터 폴더를 권장합니다."
    if os_type == "windows" and str(resolved).rstrip("\\/").upper() in {"C:", "C:\\"}:
        return "C드라이브에는 Windows 시스템 파일과 앱 설정 파일이 포함되어 있습니다. 하위 폴더를 선택하는 것을 권장합니다."
    return None


def _sample_file(path: Path, root: Path, status: str, reason: str = "") -> Dict[str, Any]:
    try:
        rel = path.relative_to(root).as_posix()
    except ValueError:
        rel = path.name
    try:
        stat = path.stat()
        size = stat.st_size if path.is_file() else None
        modified_at = _safe_iso_from_stat_mtime(stat.st_mtime)
    except OSError:
        size = None
        modified_at = ""
    return {
        "path": str(path),
        "relative_path": rel,
        "name": path.name,
        "extension": path.suffix.lower(),
        "status": status,
        "reason": reason,
        "size_bytes": size,
        "modified_at": modified_at,
    }


def _clean_text(text: str) -> str:
    return re.sub(r"\s+", " ", str(text or "")).strip()


def _chunks(text: str, size: int = 1200, overlap: int = 160) -> List[str]:
    cleaned = str(text or "").strip()
    if not cleaned:
        return []
    chunks: List[str] = []
    start = 0
    while start < len(cleaned):
        end = min(len(cleaned), start + size)
        chunks.append(cleaned[start:end])
        if end >= len(cleaned):
            break
        start = max(0, end - overlap)
    return chunks


_LLM_EXTRACT_CONCEPT_PROMPT = """Extract the key concepts from the following text.
Return ONLY a JSON array of objects, each with "concept" (string) and "importance" (float 0-1).
Extract up to {limit} concepts. Focus on named entities, technical terms, and domain-specific nouns.
Do NOT include common words, stop words, or generic terms.

Text:
{text}

JSON:"""

_LLM_EXTRACT_TRIPLE_PROMPT = """Extract relationship triples from the following text.
Return ONLY a JSON array of objects, each with:
- "subject": source concept (string)
- "relation": relationship verb (string, Korean or English)
- "object": target concept (string)
- "evidence": the sentence supporting this triple (string, max 240 chars)
- "confidence": how confident you are (float 0-1)

Extract up to {limit} triples. Focus on meaningful semantic relationships.

Text:
{text}

Concepts already identified: {concepts}

JSON:"""

ENABLE_LLM_EXTRACTION = os.getenv("LATTICEAI_LLM_EXTRACTION", "true").lower() in ("1", "true", "yes")


def _llm_extract_concepts(text: str, limit: int = 12) -> Optional[List[str]]:
    if not ENABLE_LLM_EXTRACTION or not _llm_router_ref:
        return None
    if not _llm_router_ref.current_model_id:
        return None
    prompt = _LLM_EXTRACT_CONCEPT_PROMPT.format(text=text[:3000], limit=limit)
    try:
        loop = asyncio.get_event_loop()
        if loop.is_running():
            import concurrent.futures
            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
                future = pool.submit(asyncio.run, _llm_router_ref.generate(prompt, max_tokens=1024, temperature=0.1))
                raw = future.result(timeout=30)
        else:
            raw = asyncio.run(_llm_router_ref.generate(prompt, max_tokens=1024, temperature=0.1))
        raw = raw.strip()
        if raw.startswith("```"):
            raw = re.sub(r"^```(?:json)?\s*", "", raw)
            raw = re.sub(r"\s*```$", "", raw)
        parsed = json.loads(raw)
        if isinstance(parsed, list):
            concepts = []
            for item in parsed[:limit]:
                if isinstance(item, dict) and "concept" in item:
                    concepts.append(item["concept"])
                elif isinstance(item, str):
                    concepts.append(item)
            return concepts if concepts else None
    except Exception as e:
        logging.debug("LLM concept extraction failed (falling back to rules): %s", e)
    return None


def _llm_extract_triples(text: str, concepts: List[str], limit: int = 20) -> Optional[List[Dict[str, str]]]:
    if not ENABLE_LLM_EXTRACTION or not _llm_router_ref:
        return None
    if not _llm_router_ref.current_model_id:
        return None
    prompt = _LLM_EXTRACT_TRIPLE_PROMPT.format(
        text=text[:3000], limit=limit,
        concepts=", ".join(concepts[:15]),
    )
    try:
        loop = asyncio.get_event_loop()
        if loop.is_running():
            import concurrent.futures
            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
                future = pool.submit(asyncio.run, _llm_router_ref.generate(prompt, max_tokens=2048, temperature=0.1))
                raw = future.result(timeout=30)
        else:
            raw = asyncio.run(_llm_router_ref.generate(prompt, max_tokens=2048, temperature=0.1))
        raw = raw.strip()
        if raw.startswith("```"):
            raw = re.sub(r"^```(?:json)?\s*", "", raw)
            raw = re.sub(r"\s*```$", "", raw)
        parsed = json.loads(raw)
        if isinstance(parsed, list):
            triples = []
            for item in parsed[:limit]:
                if isinstance(item, dict) and "subject" in item and "object" in item:
                    triples.append({
                        "subject": str(item["subject"]),
                        "relation": str(item.get("relation", "관련됨")),
                        "object": str(item["object"]),
                        "context": str(item.get("evidence", ""))[:240],
                        "confidence": float(item.get("confidence", 0.8)),
                    })
            return triples if triples else None
    except Exception as e:
        logging.debug("LLM triple extraction failed (falling back to rules): %s", e)
    return None


_CONCEPT_STOP: set = {
    # English stop words
    "the", "and", "for", "with", "this", "that", "from", "into", "which",
    "are", "was", "were", "has", "have", "had", "can", "will", "would",
    "could", "should", "may", "might", "must", "shall", "being", "been",
    "also", "just", "then", "than", "when", "where", "what", "how", "why",
    "its", "their", "your", "our", "you", "they", "them", "these", "those",
    "use", "used", "using", "based", "like", "such", "via", "per", "let",
    "yes", "not", "but", "are", "all", "any", "out", "new", "get", "set",
    # Korean stop words
    "사용자", "내용", "파일", "채팅", "답변", "입니다", "그리고", "처럼",
    "있어", "없어", "이야", "이다", "한다", "하다", "되다", "됩니다",
    "경우", "방법", "부분", "상태", "정도", "결과", "이후", "이전",
    "그것", "이것", "저것", "여기", "거기", "저기", "우리", "저희",
    "기능", "서버", "모델", "설정", "설명", "버전", "지원", "사용", "실행",
    "todo", "fixme", "note", "참고", "주의", "warning",
}


def _extract_concepts(text: str, limit: int = 12) -> List[str]:
    """LLM-first concept extraction with rule-based fallback."""
    llm_result = _llm_extract_concepts(text, limit)
    if llm_result:
        return llm_result
    return _extract_concepts_rules(text, limit)


def _extract_concepts_rules(text: str, limit: int = 12) -> List[str]:
    """Extract meaningful named concepts from text (rule-based).

    Priority order:
    1. Backtick / quoted terms (explicitly technical)
    2. Multi-word proper nouns (Lattice AI, GPT-4o, Claude Sonnet)
    3. Single capitalized proper nouns not at sentence start (Claude, Python, FastAPI)
    4. Korean compound technical terms (멀티모달, 에이전트, 그래프RAG)
    5. Hyphenated / versioned identifiers (gpt-4o, mlx-vlm, gemma-4)
    """
    text = str(text or "")
    seen: dict = {}  # concept_lower → original form

    def _add(term: str) -> None:
        key = term.strip().lower()
        if (
            key
            and key not in _CONCEPT_STOP
            and not key.isdigit()
            and len(key) >= 2
        ):
            seen.setdefault(key, term.strip())

    # 1. Backtick-quoted code/term (highest confidence)
    for m in re.findall(r'`([^`]{2,40})`', text):
        if not re.search(r'[\(\)\[\]{}]', m):  # skip code expressions
            _add(m)

    # 2. Double/single quoted terms
    for m in re.findall(r'"([^"]{2,40})"', text):
        _add(m)

    # 3. Multi-word English proper nouns (Title Case or ALL-CAPS first word, 2–4 words).
    #    Pattern A: Mixed-case first word — "Lattice AI", "Tool Use", "Graph RAG"
    for m in re.findall(
        r'([A-Z][a-z]{1,20}(?:\s+(?:[A-Z]{2,10}|[A-Z][a-z0-9]{1,20}|\d[\w.]{0,6})){1,3})',
        text,
    ):
        _add(m)
    #    Pattern B: ALL-CAPS first word — "VS Code", "MCP Server", "GPT-4o Mini"
    for m in re.findall(
        r'([A-Z]{2,6}(?:\s+(?:[A-Z]{2,10}|[A-Z][a-z0-9]{1,20})){1,2})',
        text,
    ):
        _add(m)

    # 4. Single capitalized proper noun.
    #    Use ASCII-boundary lookaround instead of \b so Korean particles
    #    (와, 의, 는 …) after an English word don't block the match.
    all_caps_words = re.findall(r'(?<![A-Za-z0-9])([A-Z][A-Za-z0-9]{2,24})(?![A-Za-z0-9])', text)
    freq: Dict[str, int] = {}
    for w in all_caps_words:
        freq[w] = freq.get(w, 0) + 1
    sentence_starts = set(re.findall(r'(?:^|(?<=[.!?])\s+)([A-Z][a-z]+)', text))
    for m, cnt in freq.items():
        if m.lower() in _CONCEPT_STOP:
            continue
        if cnt >= 2 or m not in sentence_starts:
            _add(m)

    # 5. Korean technical compound nouns (3–12 chars, no common particles)
    for m in re.findall(r'[가-힣]{2,12}(?:AI|LLM|API|UI|RAG|bot|Bot|기능|모델|서버|에이전트|파이프라인|워크플로)', text):
        _add(m)
    # Korean standalone terms that appear after topic markers (은/는/이/가 앞)
    for m in re.findall(r'([가-힣]{2,12})(?:은|는|이|가|을|를|의|에서|으로|와|과)', text):
        if m.lower() not in _CONCEPT_STOP and len(m) >= 2:
            # Only add if it's non-trivial (has 3+ chars or appears multiple times)
            cnt = text.count(m)
            if len(m) >= 3 or cnt >= 2:
                _add(m)

    # 6. Hyphenated / versioned identifiers (gpt-4o, gemma-4, mlx-vlm)
    for m in re.findall(r'\b([a-zA-Z][a-zA-Z0-9]*(?:-[a-zA-Z0-9.]+)+)\b', text):
        if len(m) >= 4:
            _add(m)

    # De-duplicate: remove shorter if ALL its occurrences in the source text
    # are followed immediately by the suffix that forms the longer concept.
    # "Lattice" → dropped when every occurrence is "Lattice AI"
    # "Claude"  → kept  because it appears as just "Claude" too.
    values = list(seen.values())
    values_lower = [v.lower() for v in values]
    keep = set(range(len(values)))
    for i, v in enumerate(values):
        vl = v.lower()
        for j, wl in enumerate(values_lower):
            if i == j or j not in keep:
                continue
            # Check if vl is a word-prefix of wl
            suffix = wl[len(vl):]
            if not (wl.startswith(vl) and re.match(r'^[\s\-]', suffix)):
                continue
            # Count occurrences of v NOT followed by the suffix
            suffix_stripped = suffix.lstrip(" -")
            # Escape for regex
            pattern_with_suffix = re.escape(v) + r'[\s\-]+' + re.escape(suffix_stripped)
            pattern_alone = re.escape(v) + r'(?![\s\-]*' + re.escape(suffix_stripped) + r')'
            alone_count = len(re.findall(pattern_alone, text, re.IGNORECASE))
            if alone_count == 0:
                # Shorter term never appears alone → safe to remove
                keep.discard(i)
                break

    final = [values[i] for i in range(len(values)) if i in keep]
    return final[:limit]


# ──────────────────────────────────────────────────────────────────────────────
# Node type taxonomy  (점 = 명사)
# ──────────────────────────────────────────────────────────────────────────────
# Chat      — 대화 세션
# Document  — 파일 (PDF·PPT·Word·Excel·이미지 등)
# Concept   — 개념·아이디어·기술 용어
# Person    — 사람 (사용자, 언급된 인물)
# Error     — 오류·버그·예외
# Code      — 코드 스니펫·함수·클래스
# Feature   — 소프트웨어 기능
# Task      — 할 일·액션 아이템
# Decision  — 결정 사항

# Edge type vocabulary  (선 = 동사 — 과거형 서술어)
EDGE_VERB = {
    "언급함":   r"언급|mention|refer|cited",
    "포함함":   r"포함|include|consist|구성|탑재|contains",
    "해결함":   r"해결|resolv|fix|수정|고쳤|closed",
    "의존함":   r"의존|depend|require|필요|based on",
    "설명함":   r"설명|explain|describe|정의|란|이란|means",
    "비교함":   r"비교|versus|vs\.?|차이|다르|compare",
    "사용함":   r"사용|use|활용|이용|apply",
    "연결함":   r"연결|connect|통합|integrate|연동|link",
    "확장함":   r"확장|extend|플러그인|plugin|addon",
    "생성함":   r"생성|만들|create|generate|build|produced",
    "대체함":   r"대체|replace|instead|alternative",
    "지원함":   r"지원|support|제공|provide|offer",
    "발생함":   r"발생|occur|throw|raise|triggered",
    "관련됨":   r"관련|related|associated|연관",
}


def _infer_edge(sentence: str) -> str:
    """Return the best-matching verb-form edge label for a sentence."""
    s = sentence.lower()
    for label, pattern in EDGE_VERB.items():
        if re.search(pattern, s):
            return label
    return "관련됨"


# Technical words that cannot be person names
_NOT_PERSON_WORDS: set = {
    "use", "api", "rag", "sdk", "ide", "cli", "llm", "mcp", "ui", "ux",
    "new", "old", "get", "set", "run", "add", "fix", "tool", "code",
    "base", "core", "data", "file", "test", "type", "mode", "view",
}


def _classify_node_type(concept: str, text: str) -> str:
    """Classify a concept into the node taxonomy.

    Term-level signals take priority; then a tight ±60-char window is used
    so distant keywords don't cause mis-classification.
    """
    term = concept.lower()

    # ── Term-level signals (highest confidence) ───────────────────────────
    if re.search(r'(?:error|exception|traceback|오류|에러|버그)$', term, re.I):
        return "Error"
    if re.search(r'error|exception|err\b', term, re.I) and len(concept) < 30:
        return "Error"
    if re.search(r'\(\)|\.py$|\.js$|\.ts$|\.go$|::\w', term):
        return "Code"

    # Person: "First Last" pattern, neither word is a known technical term
    if re.match(r'^[A-Z][a-z]{1,15} [A-Z][a-z]{1,15}$', concept):
        words = term.split()
        if not any(w in _NOT_PERSON_WORDS for w in words):
            return "Person"

    # ── Windowed context (±60 chars) — NOT used for Error to avoid false positives
    idx = text.lower().find(term)
    if idx >= 0:
        win = text[max(0, idx - 60): idx + len(concept) + 60].lower()
        if re.search(r'def |class |function|함수|클래스|메서드|import', win):
            return "Code"
        # Feature: concept appears DIRECTLY adjacent to 기능/feature keyword
        if (
            len(concept) <= 12
            and re.search(
                rf'{re.escape(term)}.{{0,8}}(?:기능|feature)|(?:기능|feature).{{0,8}}{re.escape(term)}',
                win,
            )
        ):
            return "Feature"

    return "Concept"


def _extract_triples(
    text: str,
    concepts: List[str],
    limit: int = 20,
) -> List[Dict[str, str]]:
    """LLM-first triple extraction with rule-based fallback."""
    llm_result = _llm_extract_triples(text, concepts, limit)
    if llm_result:
        return llm_result
    return _extract_triples_rules(text, concepts, limit)


def _extract_triples_rules(
    text: str,
    concepts: List[str],
    limit: int = 20,
) -> List[Dict[str, str]]:
    """Extract (subject, verb-edge, object, context) triples from text (rule-based).

    For each sentence containing ≥2 concepts, infer the verb-form edge label
    from surrounding context and create a directed triple.
    """
    if len(concepts) < 2:
        return []

    concept_lower = {c.lower(): c for c in concepts}
    triples: List[Dict[str, str]] = []
    seen_pairs: set = set()

    # Split on sentence boundaries
    sentences = re.split(r'(?<=[.!?\n])\s+|\n{2,}', text)
    for sent in sentences:
        sent = sent.strip()
        if len(sent) < 8:
            continue
        sent_lower = sent.lower()

        present = [concept_lower[k] for k in concept_lower if k in sent_lower]
        if len(present) < 2:
            continue

        edge = _infer_edge(sent)

        for i in range(len(present) - 1):
            subj, obj = present[i], present[i + 1]
            # Deduplicate by (subj, obj) regardless of direction for same edge
            pair_key = tuple(sorted([subj.lower(), obj.lower()])) + (edge,)
            if pair_key in seen_pairs:
                continue
            seen_pairs.add(pair_key)
            triples.append({
                "subject": subj,
                "relation": edge,          # verb form (동사)
                "object": obj,
                "context": sent[:240],
            })
            if len(triples) >= limit:
                return triples

    return triples


def _semantic_items(text: str) -> List[Dict[str, str]]:
    """Extract explicit decision / task items from text."""
    items: List[Dict[str, str]] = []
    for raw_line in str(text or "").splitlines():
        line = _clean_text(raw_line)
        if len(line) < 6:
            continue
        lowered = line.lower()
        if re.search(r"(결정|확정|하기로|decided|decision)", lowered):
            items.append({"type": "Decision", "title": line[:120], "summary": line[:500]})
        if re.search(r"(todo|해야|하자|진행|구현|수정|확인|next|task|\[ \])", lowered):
            items.append({"type": "Task", "title": line[:120], "summary": line[:500]})
    return items[:8]


def _topic_candidates(text: str, limit: int = 8) -> List[str]:
    """Return compact keyword candidates for fallback graph search."""
    candidates = _extract_concepts(text, limit=limit)
    if candidates:
        return candidates[:limit]
    seen: Dict[str, str] = {}
    for token in re.findall(r"[A-Za-z][A-Za-z0-9_.:-]{2,}|[가-힣]{2,12}", str(text or "")):
        key = token.lower()
        if key in _CONCEPT_STOP or key.isdigit():
            continue
        seen.setdefault(key, token)
        if len(seen) >= limit:
            break
    return list(seen.values())[:limit]


class KnowledgeGraphStore:
    def __init__(self, db_path: Path, blob_dir: Path):
        self.db_path = Path(db_path)
        self.blob_dir = Path(blob_dir)
        self.db_path.parent.mkdir(parents=True, exist_ok=True)
        self.blob_dir.mkdir(parents=True, exist_ok=True)
        self._embedding_model = LocalEmbeddingModel()
        self._init_db()
        # Read graph queries from the v2 projection (kgv2_* views) when available.
        # Toggle off (e.g. in tests) to compare against the legacy tables.
        self._read_from_v2 = KGStoreV2 is not None and _READ_FROM_V2_DEFAULT

    def _read_tables(self) -> tuple:
        """Return (nodes_table, edges_table) for read queries.

        Same read code runs against the legacy tables or the v2 reconstruction
        views, so the two paths are equivalent by construction.
        """
        if self._read_from_v2:
            return ("kgv2_nodes", "kgv2_edges")
        return ("nodes", "edges")

    def _connect(self) -> sqlite3.Connection:
        conn = sqlite3.connect(str(self.db_path))
        conn.row_factory = sqlite3.Row
        conn.execute("PRAGMA journal_mode=WAL")
        conn.execute("PRAGMA foreign_keys=ON")
        return conn

    def _init_db(self) -> None:
        with self._connect() as conn:
            conn.executescript(
                """
                CREATE TABLE IF NOT EXISTS graph_meta (
                  key TEXT PRIMARY KEY,
                  value TEXT NOT NULL
                );
                CREATE TABLE IF NOT EXISTS nodes (
                  id TEXT PRIMARY KEY,
                  type TEXT NOT NULL,
                  title TEXT NOT NULL,
                  summary TEXT,
                  metadata_json TEXT NOT NULL CHECK (json_valid(metadata_json)),
                  raw_json TEXT NOT NULL CHECK (json_valid(raw_json)),
                  created_at TEXT NOT NULL,
                  updated_at TEXT NOT NULL
                );
                CREATE TABLE IF NOT EXISTS edges (
                  id TEXT PRIMARY KEY,
                  from_node TEXT NOT NULL,
                  to_node TEXT NOT NULL,
                  type TEXT NOT NULL,
                  weight REAL NOT NULL DEFAULT 1.0,
                  metadata_json TEXT NOT NULL CHECK (json_valid(metadata_json)),
                  created_at TEXT NOT NULL,
                  UNIQUE(from_node, to_node, type),
                  FOREIGN KEY(from_node) REFERENCES nodes(id) ON DELETE CASCADE,
                  FOREIGN KEY(to_node) REFERENCES nodes(id) ON DELETE CASCADE
                );
                CREATE TABLE IF NOT EXISTS chunks (
                  id TEXT PRIMARY KEY,
                  source_node TEXT NOT NULL,
                  text TEXT NOT NULL,
                  metadata_json TEXT NOT NULL CHECK (json_valid(metadata_json)),
                  created_at TEXT NOT NULL,
                  FOREIGN KEY(source_node) REFERENCES nodes(id) ON DELETE CASCADE
                );
                CREATE TABLE IF NOT EXISTS knowledge_sources (
                  id TEXT PRIMARY KEY,
                  root_path TEXT NOT NULL UNIQUE,
                  os_type TEXT NOT NULL,
                  drive_id TEXT,
                  label TEXT,
                  status TEXT NOT NULL,
                  include_ocr INTEGER NOT NULL DEFAULT 0,
                  watch_enabled INTEGER NOT NULL DEFAULT 0,
                  consent_json TEXT NOT NULL CHECK (json_valid(consent_json)),
                  created_at TEXT NOT NULL,
                  updated_at TEXT NOT NULL,
                  last_scanned_at TEXT
                );
                CREATE TABLE IF NOT EXISTS local_file_index (
                  id TEXT PRIMARY KEY,
                  source_id TEXT NOT NULL,
                  os_type TEXT NOT NULL,
                  drive_id TEXT,
                  root_path TEXT NOT NULL,
                  file_path TEXT NOT NULL,
                  relative_path TEXT NOT NULL,
                  file_name TEXT NOT NULL,
                  extension TEXT NOT NULL,
                  size_bytes INTEGER,
                  modified_at TEXT,
                  sha256 TEXT,
                  last_scanned_at TEXT,
                  last_indexed_at TEXT,
                  parser_type TEXT,
                  status TEXT NOT NULL,
                  error_message TEXT,
                  graph_node_id TEXT,
                  deleted INTEGER NOT NULL DEFAULT 0,
                  metadata_json TEXT NOT NULL CHECK (json_valid(metadata_json)),
                  UNIQUE(source_id, relative_path),
                  FOREIGN KEY(source_id) REFERENCES knowledge_sources(id) ON DELETE CASCADE
                );
                CREATE TABLE IF NOT EXISTS vector_embeddings (
                  item_id TEXT PRIMARY KEY,
                  item_type TEXT NOT NULL,
                  source_node TEXT NOT NULL,
                  text_hash TEXT NOT NULL,
                  embedding BLOB NOT NULL,
                  embedding_dim INTEGER NOT NULL,
                  embedding_model TEXT NOT NULL,
                  metadata_json TEXT NOT NULL CHECK (json_valid(metadata_json)),
                  indexed_at TEXT NOT NULL,
                  FOREIGN KEY(source_node) REFERENCES nodes(id) ON DELETE CASCADE
                );
                CREATE TABLE IF NOT EXISTS vector_index_operations (
                  id TEXT PRIMARY KEY,
                  operation TEXT NOT NULL,
                  status TEXT NOT NULL,
                  requested_at TEXT NOT NULL,
                  started_at TEXT,
                  completed_at TEXT,
                  items_total INTEGER NOT NULL DEFAULT 0,
                  items_indexed INTEGER NOT NULL DEFAULT 0,
                  items_skipped INTEGER NOT NULL DEFAULT 0,
                  error_message TEXT,
                  metadata_json TEXT NOT NULL CHECK (json_valid(metadata_json))
                );
                CREATE INDEX IF NOT EXISTS idx_nodes_type ON nodes(type);
                CREATE INDEX IF NOT EXISTS idx_edges_from ON edges(from_node);
                CREATE INDEX IF NOT EXISTS idx_edges_to ON edges(to_node);
                CREATE INDEX IF NOT EXISTS idx_chunks_source ON chunks(source_node);
                CREATE INDEX IF NOT EXISTS idx_knowledge_sources_root ON knowledge_sources(root_path);
                CREATE INDEX IF NOT EXISTS idx_local_file_index_source ON local_file_index(source_id);
                CREATE INDEX IF NOT EXISTS idx_local_file_index_status ON local_file_index(status);
                CREATE INDEX IF NOT EXISTS idx_local_file_index_graph_node ON local_file_index(graph_node_id);
                CREATE INDEX IF NOT EXISTS idx_vector_embeddings_type ON vector_embeddings(item_type);
                CREATE INDEX IF NOT EXISTS idx_vector_embeddings_source ON vector_embeddings(source_node);
                CREATE INDEX IF NOT EXISTS idx_vector_embeddings_model ON vector_embeddings(embedding_model);
                CREATE INDEX IF NOT EXISTS idx_vector_index_operations_requested ON vector_index_operations(requested_at);
                """
            )
            conn.execute(
                "INSERT OR REPLACE INTO graph_meta(key, value) VALUES (?, ?)",
                ("schema_version", str(GRAPH_SCHEMA_VERSION)),
            )
        self._init_v2_schema()

    # SQL views that reconstruct the *exact* legacy row shape on top of the
    # normalized v2 tables, so the read methods run unchanged against either
    # source. The projection stores the raw legacy type string in ``legacy_type``
    # and promotes summary + metadata to first-class columns (no more
    # ``attrs._kg`` passthrough / ``evidence`` abuse), so these views are
    # byte-faithful to the legacy nodes/edges tables.
    _V2_VIEWS_SQL = """
    CREATE VIEW IF NOT EXISTS kgv2_nodes AS
      SELECT id,
             COALESCE(legacy_type, type) AS type,
             label AS title,
             summary,
             attrs AS metadata_json,
             created_at, updated_at
      FROM nodes_v2;
    CREATE VIEW IF NOT EXISTS kgv2_edges AS
      SELECT id, source AS from_node, target AS to_node,
             COALESCE(legacy_type, type) AS type,
             weight,
             metadata AS metadata_json,
             created_at
      FROM edges_v2;
    """

    def _init_v2_schema(self) -> None:
        """Initialize the normalized v2 tables + reconstruction views, migrating
        the projection layout when it is stale — **atomically**.

        The entire DROP → CREATE → VIEWS → BACKFILL → version-stamp sequence runs
        in a single transaction on one connection: on any failure it rolls back,
        leaving the prior projection untouched and the version unchanged, so the
        next startup simply retries. The migration only ever touches the v2
        tables/views and the ``projection_version`` key — never the authoritative
        legacy ``nodes``/``edges`` — so legacy data cannot be corrupted even if
        the rebuild fails midway.
        """
        if KGStoreV2 is None or _exec_script is None:
            return
        try:
            with self._connect() as conn:
                conn.execute("BEGIN")
                stale = self._projection_version(conn) != _PROJECTION_VERSION
                if stale:
                    # The projection is non-authoritative; drop it so init_schema
                    # recreates the tables with the current normalized columns.
                    for stmt in (
                        "DROP VIEW IF EXISTS kgv2_edges",
                        "DROP VIEW IF EXISTS kgv2_nodes",
                        "DROP TABLE IF EXISTS edges_v2",
                        "DROP TABLE IF EXISTS nodes_v2",
                    ):
                        conn.execute(stmt)
                # init_schema(conn=...) joins this transaction (no implicit commit)
                KGStoreV2(self.db_path).init_schema(conn=conn)
                _exec_script(conn, self._V2_VIEWS_SQL)
                self._backfill_v2_on(conn, force=stale)
                # version stamp commits together with the backfill — never stranded
                conn.execute(
                    "INSERT OR REPLACE INTO kg_meta(key, value) VALUES ('projection_version', ?)",
                    (str(_PROJECTION_VERSION),),
                )
        except Exception as e:
            logging.warning("knowledge_graph: v2 schema init/backfill skipped: %s", e)

    def _projection_version(self, conn: sqlite3.Connection) -> int:
        """Return the stored v2 projection layout version (0 if unknown).

        A fresh DB (kg_meta absent) raises ``sqlite3.OperationalError`` here and
        is correctly treated as version 0 → rebuild. Only sqlite errors are
        swallowed so a real bug doesn't masquerade as a stale projection.
        """
        try:
            row = conn.execute(
                "SELECT value FROM kg_meta WHERE key='projection_version'"
            ).fetchone()
            return int(row["value"]) if row and row["value"] is not None else 0
        except sqlite3.Error:
            return 0

    def _backfill_v2_if_needed(self, *, force: bool = False) -> None:
        """Project legacy nodes/edges into v2 on a fresh transaction.

        Thin wrapper around :meth:`_backfill_v2_on` for callers (tests, ad-hoc
        re-sync) that aren't already inside the migration transaction.
        """
        try:
            with self._connect() as conn:
                self._backfill_v2_on(conn, force=force)
        except Exception as ex:
            logging.warning("knowledge_graph: v2 backfill skipped: %s", ex)

    def _backfill_v2_on(self, conn: sqlite3.Connection, *, force: bool = False) -> None:
        """Project legacy nodes/edges into the normalized v2 tables on ``conn``.

        Non-destructive to legacy. ``force`` rebuilds unconditionally (used after
        a layout migration); otherwise it only projects when v2 is empty. The v2
        graph is a derived projection, so clearing + rebuilding it is always safe.
        Idempotent: no-ops once v2 carries the current projection. Copies the
        legacy column values **verbatim** so the kgv2_* views are byte-faithful.
        """
        legacy_nodes = conn.execute("SELECT COUNT(*) FROM nodes").fetchone()[0]
        if legacy_nodes == 0:
            return
        v2_nodes = conn.execute("SELECT COUNT(*) FROM nodes_v2").fetchone()[0]
        if v2_nodes > 0 and not force:
            return  # current projection already present
        # (re)project: clear v2 graph (not authoritative) and rebuild
        conn.execute("DELETE FROM edges_v2")
        conn.execute("DELETE FROM nodes_v2")
        n = e = 0
        for r in conn.execute(
            "SELECT id, type, title, summary, metadata_json, created_at, updated_at FROM nodes"
        ).fetchall():
            self._v2_project_node(
                conn, r["id"], r["type"], r["title"], r["summary"], r["metadata_json"],
                created_at=r["created_at"], updated_at=r["updated_at"],
            )
            n += 1
        for r in conn.execute(
            "SELECT id, from_node, to_node, type, weight, metadata_json, created_at FROM edges"
        ).fetchall():
            self._v2_project_edge(
                conn, r["from_node"], r["to_node"], r["type"], float(r["weight"] or 1.0),
                r["metadata_json"], edge_id=r["id"], created_at=r["created_at"],
            )
            e += 1
        logging.info("knowledge_graph: projected legacy → v2 (%d nodes, %d edges)", n, e)

    # ── v2 dual-write projection (normalized type, byte-faithful legacy values) ──
    # The projection stores the legacy ``title``/``summary``/``metadata_json``
    # values it is handed VERBATIM (no truncation or JSON re-encoding) so the
    # kgv2_* views reproduce the legacy rows exactly. Callers (_upsert_* and the
    # backfill) pass the already-canonical legacy column values.
    def _v2_project_node(
        self, conn: sqlite3.Connection, node_id: str, node_type: str, title: str,
        summary: Optional[str], metadata_json: Optional[str],
        *, created_at: Optional[str] = None, updated_at: Optional[str] = None,
    ) -> None:
        if KGStoreV2 is None:
            return
        ts = updated_at or _now()
        norm_type = NodeType.from_legacy(node_type).value if NodeType is not None else node_type
        try:
            conn.execute(
                """
                INSERT INTO nodes_v2(id, type, legacy_type, label, summary, attrs,
                                     owner_id, visibility, created_at, updated_at,
                                     importance_score)
                VALUES (?, ?, ?, ?, ?, ?, NULL, 'private', ?, ?, 0.0)
                ON CONFLICT(id) DO UPDATE SET
                  type=excluded.type, legacy_type=excluded.legacy_type,
                  label=excluded.label, summary=excluded.summary,
                  attrs=excluded.attrs, updated_at=excluded.updated_at
                """,
                (node_id, norm_type, node_type, title, summary,
                 metadata_json if metadata_json is not None else "{}",
                 created_at or ts, ts),
            )
        except Exception as ex:
            logging.debug("knowledge_graph: v2 node projection skipped (%s): %s", node_id, ex)

    def _v2_project_edge(
        self, conn: sqlite3.Connection, from_node: str, to_node: str, edge_type: str,
        weight: float, metadata_json: Optional[str],
        *, edge_id: Optional[str] = None, created_at: Optional[str] = None,
    ) -> None:
        if KGStoreV2 is None:
            return
        eid = edge_id or f"edge:{_sha256_text(f'{from_node}|{edge_type}|{to_node}')[:24]}"
        norm_type = EdgeType.from_legacy(edge_type).value if EdgeType is not None else edge_type
        meta_str = metadata_json if metadata_json is not None else "{}"
        confidence = float(_safe_loads(meta_str).get("confidence", 1.0))
        try:
            conn.execute(
                """
                INSERT INTO edges_v2(id, source, target, type, legacy_type, weight,
                                     confidence, evidence, metadata, created_by, created_at)
                VALUES (?, ?, ?, ?, ?, ?, ?, '[]', ?, 'legacy', ?)
                ON CONFLICT(source, target, legacy_type) DO UPDATE SET
                  type=excluded.type,
                  weight=max(edges_v2.weight, excluded.weight),
                  confidence=excluded.confidence,
                  metadata=excluded.metadata
                """,
                (eid, from_node, to_node, norm_type, edge_type, float(weight),
                 confidence, meta_str, created_at or _now()),
            )
        except Exception as ex:
            logging.debug("knowledge_graph: v2 edge projection skipped (%s->%s): %s", from_node, to_node, ex)

    def _v2_delete_nodes(self, conn: sqlite3.Connection, ids) -> None:
        """Mirror legacy node deletions into v2 (edges_v2 cascade on the FK)."""
        if KGStoreV2 is None:
            return
        ids = list(ids)
        if not ids:
            return
        ph = ",".join("?" * len(ids))
        try:
            conn.execute(f"DELETE FROM nodes_v2 WHERE id IN ({ph})", ids)
        except Exception as ex:
            logging.debug("knowledge_graph: v2 node delete mirror skipped: %s", ex)

    def _v2_delete_edges_from(self, conn: sqlite3.Connection, node_id: str) -> None:
        """Mirror a legacy ``DELETE FROM edges WHERE from_node=?`` into v2."""
        if KGStoreV2 is None:
            return
        try:
            conn.execute("DELETE FROM edges_v2 WHERE source=?", (node_id,))
        except Exception as ex:
            logging.debug("knowledge_graph: v2 edge delete mirror skipped: %s", ex)

    def _v2_sync_report(self) -> Dict[str, Any]:
        """Diagnose the dual-write invariant: legacy node/edge id sets must equal
        the v2 projection's. Returns counts + any drift (ids missing from / extra
        in v2). ``in_sync`` is True only when both id sets match exactly.

        All legacy writes go through _upsert_node/_upsert_edge (which dual-write)
        and every legacy delete is mirrored, so a non-empty drift signals a
        bypassed write path — this is the runtime guard for that invariant.
        """
        if KGStoreV2 is None:
            return {"available": False, "in_sync": True}
        with self._connect() as conn:
            legacy_nodes = {r[0] for r in conn.execute("SELECT id FROM nodes")}
            v2_nodes = {r[0] for r in conn.execute("SELECT id FROM nodes_v2")}
            legacy_edges = {r[0] for r in conn.execute("SELECT id FROM edges")}
            v2_edges = {r[0] for r in conn.execute("SELECT id FROM edges_v2")}
        return {
            "available": True,
            "in_sync": legacy_nodes == v2_nodes and legacy_edges == v2_edges,
            "nodes_legacy": len(legacy_nodes),
            "nodes_v2": len(v2_nodes),
            "edges_legacy": len(legacy_edges),
            "edges_v2": len(v2_edges),
            "nodes_missing_from_v2": sorted(legacy_nodes - v2_nodes),
            "nodes_extra_in_v2": sorted(v2_nodes - legacy_nodes),
            "edges_missing_from_v2": sorted(legacy_edges - v2_edges),
            "edges_extra_in_v2": sorted(v2_edges - legacy_edges),
        }

    def _upsert_node(
        self,
        conn: sqlite3.Connection,
        node_id: str,
        node_type: str,
        title: str,
        summary: str = "",
        metadata: Optional[Dict[str, Any]] = None,
        raw: Optional[Dict[str, Any]] = None,
    ) -> str:
        now = _now()
        # Canonical stored values, computed once and shared with the v2
        # projection so legacy and v2 hold byte-identical strings.
        title_s = title[:240]
        summary_s = summary[:1000]
        meta_json = _json(metadata)
        conn.execute(
            """
            INSERT INTO nodes(id, type, title, summary, metadata_json, raw_json, created_at, updated_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
            ON CONFLICT(id) DO UPDATE SET
              title=excluded.title,
              summary=excluded.summary,
              metadata_json=excluded.metadata_json,
              raw_json=excluded.raw_json,
              updated_at=excluded.updated_at
            """,
            (node_id, node_type, title_s, summary_s, meta_json, _json(raw), now, now),
        )
        # dual-write: project into the v2 graph on the same transaction
        self._v2_project_node(conn, node_id, node_type, title_s, summary_s, meta_json,
                              created_at=now, updated_at=now)
        if node_type != "Chunk":
            self._upsert_vector_item(
                conn,
                item_id=node_id,
                item_type="node",
                source_node=node_id,
                text=self._vector_text_for_node(title=title_s, summary=summary_s, metadata=metadata),
                metadata={"node_type": node_type, **(metadata or {})},
            )
        return node_id

    def _upsert_edge(
        self,
        conn: sqlite3.Connection,
        from_node: str,
        to_node: str,
        edge_type: str,
        weight: float = 1.0,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> str:
        edge_id = f"edge:{_sha256_text(f'{from_node}|{edge_type}|{to_node}')[:24]}"
        now = _now()
        meta_json = _json(metadata)   # canonical string shared with the projection
        conn.execute(
            """
            INSERT INTO edges(id, from_node, to_node, type, weight, metadata_json, created_at)
            VALUES (?, ?, ?, ?, ?, ?, ?)
            ON CONFLICT(from_node, to_node, type) DO UPDATE SET
              weight=max(edges.weight, excluded.weight),
              metadata_json=excluded.metadata_json
            """,
            (edge_id, from_node, to_node, edge_type, float(weight), meta_json, now),
        )
        # dual-write: project into the v2 graph on the same transaction
        self._v2_project_edge(conn, from_node, to_node, edge_type, float(weight), meta_json,
                              edge_id=edge_id, created_at=now)
        return edge_id

    def _vector_text_for_node(
        self,
        *,
        title: str,
        summary: str = "",
        metadata: Optional[Dict[str, Any]] = None,
    ) -> str:
        metadata = metadata or {}
        meta_parts = []
        for key in (
            "filename", "relative_path", "file_path", "conversation_id", "source",
            "category", "ext", "role",
        ):
            value = metadata.get(key)
            if value:
                meta_parts.append(str(value))
        return _clean_text("\n".join([str(title or ""), str(summary or ""), " ".join(meta_parts)]))

    def _upsert_vector_item(
        self,
        conn: sqlite3.Connection,
        *,
        item_id: str,
        item_type: str,
        source_node: str,
        text: str,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> bool:
        text = _clean_text(text)
        if len(text) < 2:
            conn.execute("DELETE FROM vector_embeddings WHERE item_id=?", (item_id,))
            return False
        text_hash = _sha256_text(text)
        existing = conn.execute(
            """
            SELECT text_hash, embedding_dim, embedding_model
            FROM vector_embeddings
            WHERE item_id=?
            """,
            (item_id,),
        ).fetchone()
        if (
            existing
            and existing["text_hash"] == text_hash
            and existing["embedding_dim"] == self._embedding_model.dim
            and existing["embedding_model"] == self._embedding_model.model_id
        ):
            return False
        embedding = self._embedding_model.encode(self._embedding_model.embed(text[:50_000]))
        conn.execute(
            """
            INSERT INTO vector_embeddings(
              item_id, item_type, source_node, text_hash, embedding,
              embedding_dim, embedding_model, metadata_json, indexed_at
            )
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
            ON CONFLICT(item_id) DO UPDATE SET
              item_type=excluded.item_type,
              source_node=excluded.source_node,
              text_hash=excluded.text_hash,
              embedding=excluded.embedding,
              embedding_dim=excluded.embedding_dim,
              embedding_model=excluded.embedding_model,
              metadata_json=excluded.metadata_json,
              indexed_at=excluded.indexed_at
            """,
            (
                item_id,
                item_type,
                source_node,
                text_hash,
                embedding,
                self._embedding_model.dim,
                self._embedding_model.model_id,
                _json(metadata),
                _now(),
            ),
        )
        return True

    def _upsert_chunk(
        self,
        conn: sqlite3.Connection,
        *,
        chunk_id: str,
        source_node: str,
        text: str,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> None:
        metadata = metadata or {}
        conn.execute(
            "INSERT OR REPLACE INTO chunks(id, source_node, text, metadata_json, created_at) "
            "VALUES (?, ?, ?, ?, ?)",
            (chunk_id, source_node, text, _json(metadata), _now()),
        )
        self._upsert_vector_item(
            conn,
            item_id=chunk_id,
            item_type="chunk",
            source_node=chunk_id,
            text=text,
            metadata={**metadata, "parent_source_node": source_node},
        )

    # ── Local folder sources → Graph RAG ──────────────────────────────────

    def discover_local_roots(self) -> Dict[str, Any]:
        """Return safe, cross-platform starting points for structure browsing."""
        os_type = _current_os_type()
        home = Path.home().expanduser()
        roots: List[Dict[str, Any]] = []
        seen: set = set()

        def add(label: str, path: Path, kind: str, *, recommended: bool = True, warning: Optional[str] = None) -> None:
            try:
                resolved = path.expanduser().resolve()
            except OSError:
                resolved = path.expanduser()
            key = str(resolved)
            if key in seen or not resolved.exists():
                return
            seen.add(key)
            roots.append({
                "id": f"{kind}:{_path_fingerprint(resolved)}",
                "label": label,
                "path": key,
                "kind": kind,
                "recommended": recommended,
                "warning": warning or _root_warning(resolved, os_type),
            })

        add("홈", home, "home", warning=_root_warning(home, os_type))
        for name, label in (
            ("Documents", "문서"),
            ("Desktop", "데스크탑"),
            ("Downloads", "다운로드"),
            ("Pictures", "사진"),
            ("Projects", "프로젝트"),
        ):
            add(label, home / name, name.lower())

        if os_type == "macos":
            volumes = Path("/Volumes")
            if volumes.exists():
                try:
                    for volume in sorted(volumes.iterdir(), key=lambda p: p.name.lower()):
                        add(volume.name, volume, "volume", recommended=False)
                except OSError:
                    pass
        elif os_type == "windows":
            for letter in "ABCDEFGHIJKLMNOPQRSTUVWXYZ":
                drive = Path(f"{letter}:\\")
                if drive.exists():
                    add(f"{letter}: 드라이브", drive, "drive", recommended=(letter != "C"))
            for env_name, label in (("OneDrive", "OneDrive"), ("OneDriveCommercial", "OneDrive")):
                raw = os.environ.get(env_name)
                if raw:
                    add(label, Path(raw), "cloud", recommended=False)
        elif os_type == "linux":
            for base in (Path("/mnt"), Path("/media")):
                add(str(base), base, "mounts", recommended=False)
                try:
                    if base.exists():
                        for mounted in sorted(base.iterdir(), key=lambda p: p.name.lower()):
                            add(mounted.name, mounted, "volume", recommended=False)
                except OSError:
                    pass

        return {
            "os_type": os_type,
            "computer": platform.node() or "local",
            "roots": roots,
            "privacy_notice": "처음에는 드라이브와 폴더 구조만 확인하며, 파일 내용은 사용자가 동의한 뒤에만 읽습니다.",
        }

    def preview_local_tree(self, path: Path, *, max_items: int = 200) -> Dict[str, Any]:
        """List one folder level using metadata only; file contents are not read."""
        root = Path(path).expanduser().resolve()
        if not root.exists():
            raise ValueError(f"경로가 존재하지 않습니다: {path}")
        if not root.is_dir():
            raise ValueError(f"폴더가 아닙니다: {path}")

        os_type = _current_os_type()
        max_items = max(1, min(int(max_items or 200), 1000))
        items: List[Dict[str, Any]] = []
        inaccessible = 0
        try:
            children = sorted(root.iterdir(), key=lambda p: (not p.is_dir(), p.name.lower()))
        except PermissionError as exc:
            return {
                "path": str(root),
                "items": [],
                "error": f"접근 권한 없음: {exc}",
                "privacy_notice": "현재 단계에서는 파일 내용을 읽지 않고, 폴더와 파일의 이름/크기/수정일만 확인합니다.",
            }

        for child in children[:max_items]:
            try:
                is_dir = child.is_dir()
                stat = child.stat()
                reason = _excluded_directory_reason(child, root=root, os_type=os_type) if is_dir else _sensitive_file_reason(child, root=root)
                items.append({
                    "name": child.name,
                    "path": str(child),
                    "type": "directory" if is_dir else "file",
                    "extension": "" if is_dir else child.suffix.lower(),
                    "size_bytes": None if is_dir else stat.st_size,
                    "modified_at": _safe_iso_from_stat_mtime(stat.st_mtime),
                    "hidden": _is_hidden_path(child, root),
                    "accessible": True,
                    "excluded_reason": reason,
                })
            except PermissionError:
                inaccessible += 1
                items.append({
                    "name": child.name,
                    "path": str(child),
                    "type": "unknown",
                    "accessible": False,
                    "excluded_reason": "permission_denied",
                })
            except OSError as exc:
                inaccessible += 1
                items.append({
                    "name": child.name,
                    "path": str(child),
                    "type": "unknown",
                    "accessible": False,
                    "excluded_reason": str(exc),
                })

        return {
            "path": str(root),
            "os_type": os_type,
            "items": items,
            "truncated": len(children) > max_items,
            "inaccessible": inaccessible,
            "warning": _root_warning(root, os_type),
            "privacy_notice": "현재 단계에서는 파일 내용을 읽지 않고, 폴더와 파일의 이름/크기/수정일만 확인합니다.",
        }

    def _iter_local_scan_entries(self, root: Path, *, max_files: int) -> Iterable[Dict[str, Any]]:
        os_type = _current_os_type()
        stack = [root]
        files_seen = 0
        while stack:
            current = stack.pop()
            try:
                children = sorted(current.iterdir(), key=lambda p: (not p.is_dir(), p.name.lower()))
            except PermissionError as exc:
                yield {"kind": "inaccessible_dir", "path": current, "reason": f"permission_denied: {exc}"}
                continue
            except OSError as exc:
                yield {"kind": "inaccessible_dir", "path": current, "reason": str(exc)}
                continue

            for child in children:
                if child.is_symlink():
                    yield {"kind": "excluded", "path": child, "reason": "symlink"}
                    continue
                try:
                    if child.is_dir():
                        reason = _excluded_directory_reason(child, root=root, os_type=os_type)
                        if reason:
                            yield {"kind": "excluded_dir", "path": child, "reason": reason}
                        else:
                            stack.append(child)
                        continue
                    if not child.is_file():
                        yield {"kind": "excluded", "path": child, "reason": "not_regular_file"}
                        continue
                    stat = child.stat()
                except PermissionError as exc:
                    yield {"kind": "inaccessible_file", "path": child, "reason": f"permission_denied: {exc}"}
                    continue
                except OSError as exc:
                    yield {"kind": "inaccessible_file", "path": child, "reason": str(exc)}
                    continue

                files_seen += 1
                if files_seen > max_files:
                    yield {"kind": "limit_reached", "path": child, "reason": "max_files"}
                    return
                yield {"kind": "file", "path": child, "stat": stat}

    def _local_file_decision(self, path: Path, root: Path, stat: os.stat_result) -> Dict[str, Any]:
        ext = path.suffix.lower()
        category = _file_category(ext)
        parser_type = _parser_type_for_category(category, ext)
        sensitive_reason = _sensitive_file_reason(path, root=root)
        if sensitive_reason:
            return {
                "status": "sensitive_blocked",
                "reason": sensitive_reason,
                "category": category,
                "parser_type": parser_type,
                "indexable": False,
            }
        if category == "unsupported":
            return {
                "status": "unsupported",
                "reason": "unsupported_extension",
                "category": category,
                "parser_type": parser_type,
                "indexable": False,
            }
        limit = _size_limit_for_category(category)
        if stat.st_size > limit:
            return {
                "status": "too_large",
                "reason": f"size>{limit}",
                "category": category,
                "parser_type": parser_type,
                "indexable": False,
            }
        return {
            "status": "pending",
            "reason": "",
            "category": category,
            "parser_type": parser_type,
            "indexable": True,
        }

    def audit_local_folder(self, path: Path, *, include_ocr: bool = False, max_files: int = 50_000) -> Dict[str, Any]:
        """Safety-check a folder using metadata only; file bodies are not read."""
        root = Path(path).expanduser().resolve()
        if not root.exists():
            raise ValueError(f"경로가 존재하지 않습니다: {path}")
        if not root.is_dir():
            raise ValueError(f"폴더가 아닙니다: {path}")

        os_type = _current_os_type()
        max_files = max(1, min(int(max_files or 50_000), 200_000))
        status_counts: Counter = Counter()
        category_counts: Counter = Counter()
        extension_counts: Counter = Counter()
        allowed_samples: List[Dict[str, Any]] = []
        excluded_samples: List[Dict[str, Any]] = []
        total_files = 0
        readable_files = 0
        inaccessible = 0
        excluded_dirs = 0
        limit_reached = False

        for entry in self._iter_local_scan_entries(root, max_files=max_files):
            kind = entry["kind"]
            path_obj = entry["path"]
            if kind == "limit_reached":
                limit_reached = True
                break
            if kind == "excluded_dir":
                excluded_dirs += 1
                if len(excluded_samples) < 25:
                    excluded_samples.append(_sample_file(path_obj, root, "excluded", entry.get("reason", "")))
                continue
            if kind in {"inaccessible_dir", "inaccessible_file"}:
                inaccessible += 1
                status_counts["failed"] += 1
                if len(excluded_samples) < 25:
                    excluded_samples.append(_sample_file(path_obj, root, "failed", entry.get("reason", "")))
                continue
            if kind == "excluded":
                status_counts["excluded"] += 1
                if len(excluded_samples) < 25:
                    excluded_samples.append(_sample_file(path_obj, root, "excluded", entry.get("reason", "")))
                continue
            if kind != "file":
                continue

            total_files += 1
            stat = entry["stat"]
            decision = self._local_file_decision(path_obj, root, stat)
            status = decision["status"]
            category = decision["category"]
            ext = path_obj.suffix.lower() or "(none)"
            category_counts[category] += 1
            extension_counts[ext] += 1
            if decision["indexable"]:
                readable_files += 1
                status_counts["readable"] += 1
                if len(allowed_samples) < 25:
                    allowed_samples.append(_sample_file(path_obj, root, "readable"))
            else:
                status_counts[status] += 1
                if len(excluded_samples) < 25:
                    excluded_samples.append(_sample_file(path_obj, root, status, decision["reason"]))

        doc_weight = category_counts["pdf"] * 1.4 + category_counts["document"] * 0.9 + category_counts["slide_deck"] * 1.0
        sheet_weight = category_counts["spreadsheet"] * 0.6
        ocr_weight = category_counts["image"] * (1.8 if include_ocr else 0.1)
        estimated_seconds = round(readable_files * 0.04 + doc_weight + sheet_weight + ocr_weight, 1)

        return {
            "path": str(root),
            "source_id": f"source:{_path_fingerprint(root)}",
            "os_type": os_type,
            "drive_id": _drive_id_for_path(root),
            "warning": _root_warning(root, os_type),
            "privacy_notice": "현재 단계에서는 파일 내용을 읽지 않고, 폴더와 파일의 이름/크기/수정일만 확인합니다.",
            "include_ocr_requested": bool(include_ocr),
            "summary": {
                "total_files": total_files,
                "readable_files": readable_files,
                "excluded_files": int(
                    status_counts["excluded"]
                    + status_counts["sensitive_blocked"]
                    + status_counts["too_large"]
                    + status_counts["unsupported"]
                ),
                "sensitive_files": int(status_counts["sensitive_blocked"]),
                "too_large_files": int(status_counts["too_large"]),
                "unsupported_files": int(status_counts["unsupported"]),
                "image_ocr_candidates": int(category_counts["image"]),
                "inaccessible_items": inaccessible,
                "excluded_dirs": excluded_dirs,
                "estimated_seconds": estimated_seconds,
                "storage_root": str(self.db_path.parent),
                "limit_reached": limit_reached,
            },
            "by_status": dict(status_counts),
            "by_category": dict(category_counts),
            "by_extension": dict(extension_counts.most_common(40)),
            "allowed_samples": allowed_samples,
            "excluded_samples": excluded_samples,
            "consent_required": {
                "knowledge_source": True,
                "image_ocr": bool(category_counts["image"]),
                "watch": True,
                "sensitive_files_default_excluded": True,
            },
        }

    def local_sources(self) -> Dict[str, Any]:
        with self._connect() as conn:
            sources = [
                {
                    "id": row["id"],
                    "root_path": row["root_path"],
                    "os_type": row["os_type"],
                    "drive_id": row["drive_id"],
                    "label": row["label"],
                    "status": row["status"],
                    "include_ocr": bool(row["include_ocr"]),
                    "watch_enabled": bool(row["watch_enabled"]),
                    "consent": _safe_loads(row["consent_json"]),
                    "created_at": row["created_at"],
                    "updated_at": row["updated_at"],
                    "last_scanned_at": row["last_scanned_at"],
                }
                for row in conn.execute(
                    """
                    SELECT id, root_path, os_type, drive_id, label, status, include_ocr,
                           watch_enabled, consent_json, created_at, updated_at, last_scanned_at
                    FROM knowledge_sources
                    ORDER BY updated_at DESC, id ASC
                    """
                )
            ]
            status_rows = conn.execute(
                "SELECT source_id, status, COUNT(*) AS count FROM local_file_index GROUP BY source_id, status"
            ).fetchall()
        counts: Dict[str, Dict[str, int]] = {}
        for row in status_rows:
            counts.setdefault(row["source_id"], {})[row["status"]] = row["count"]
        for source in sources:
            source["file_status"] = counts.get(source["id"], {})
        return {"sources": sources}

    def set_local_source_watch(self, source_id: str, enabled: bool) -> Dict[str, Any]:
        source_id = str(source_id or "").strip()
        if not source_id:
            raise ValueError("source_id required")
        with self._connect() as conn:
            row = conn.execute(
                "SELECT id FROM knowledge_sources WHERE id=?",
                (source_id,),
            ).fetchone()
            if not row:
                raise ValueError(f"knowledge source not found: {source_id}")
            conn.execute(
                "UPDATE knowledge_sources SET watch_enabled=?, updated_at=? WHERE id=?",
                (1 if enabled else 0, _now(), source_id),
            )
        return {"source_id": source_id, "watch_enabled": bool(enabled)}

    def remove_local_source(self, source_id: str) -> Dict[str, Any]:
        """Remove one approved local source and its derived graph projection.

        This is intentionally non-destructive for user files: only the LatticeAI
        index rows, graph nodes, edges, and chunks derived from the source are
        removed. The original folder and files are never touched.
        """
        source_id = str(source_id or "").strip()
        if not source_id:
            raise ValueError("source_id required")
        with self._connect() as conn:
            source = conn.execute(
                "SELECT id, root_path FROM knowledge_sources WHERE id=?",
                (source_id,),
            ).fetchone()
            if not source:
                raise ValueError(f"knowledge source not found: {source_id}")
            rows = conn.execute(
                "SELECT graph_node_id FROM local_file_index WHERE source_id=? AND graph_node_id IS NOT NULL",
                (source_id,),
            ).fetchall()
            graph_node_ids = [row["graph_node_id"] for row in rows if row["graph_node_id"]]
            for graph_node_id in graph_node_ids:
                self._delete_local_file_graph(conn, graph_node_id)
            conn.execute("DELETE FROM local_file_index WHERE source_id=?", (source_id,))
            conn.execute("DELETE FROM knowledge_sources WHERE id=?", (source_id,))
            self._cleanup_local_graph_orphans(conn, source_id)
        return {
            "source_id": source_id,
            "root_path": source["root_path"],
            "removed_graph_nodes": len(graph_node_ids),
        }

    def _extract_local_file_text(self, path: Path, category: str, *, include_ocr: bool) -> Tuple[str, Dict[str, Any]]:
        ext = path.suffix.lower()
        meta: Dict[str, Any] = {"parser": _parser_type_for_category(category, ext)}
        text = ""
        if category in {"text", "code"} or ext == ".csv":
            text = path.read_text(encoding="utf-8", errors="replace")
        elif ext == ".pdf":
            import pdfplumber
            with pdfplumber.open(str(path)) as pdf:
                meta["pages"] = len(pdf.pages)
                text = "\n\n".join((page.extract_text() or "") for page in pdf.pages)
        elif ext == ".docx":
            from docx import Document
            doc = Document(str(path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            table_lines = []
            for table in doc.tables:
                for row in table.rows:
                    cells = [_clean_text(cell.text) for cell in row.cells]
                    if any(cells):
                        table_lines.append("\t".join(cells))
            meta["paragraphs"] = len(paragraphs)
            meta["tables"] = len(doc.tables)
            meta["table_rows"] = len(table_lines)
            text = "\n\n".join([*paragraphs, *table_lines])
        elif ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(str(path), read_only=True, data_only=True)
            rows_all = []
            non_empty_rows = 0
            non_empty_cells = 0
            char_count = 0
            for ws in wb.worksheets:
                sheet_rows = []
                for row in ws.iter_rows(values_only=True):
                    cells = [str(cell).strip() if cell is not None else "" for cell in row]
                    if not any(cells):
                        continue
                    line = "\t".join(cells)
                    non_empty_rows += 1
                    non_empty_cells += sum(1 for cell in cells if cell)
                    sheet_rows.append(line)
                    char_count += len(line) + 1
                    if char_count > 200_000:
                        break
                if sheet_rows:
                    rows_all.append(f"[Sheet: {ws.title}]")
                    rows_all.extend(sheet_rows)
                if char_count > 200_000:
                    break
            meta["sheets"] = len(wb.worksheets)
            meta["rows"] = non_empty_rows
            meta["cells"] = non_empty_cells
            text = "\n".join(rows_all)
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            slides_text = []
            for index, slide in enumerate(prs.slides, 1):
                parts = []
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        slide_text = shape.text_frame.text.strip()
                        if slide_text:
                            parts.append(slide_text)
                if parts:
                    slides_text.append(f"[Slide {index}]\n" + "\n".join(parts))
            meta["slides"] = len(prs.slides)
            meta["text_slides"] = len(slides_text)
            text = "\n\n".join(slides_text)
        elif category == "image":
            from PIL import Image
            with Image.open(str(path)) as image:
                meta.update({
                    "width": image.width,
                    "height": image.height,
                    "format": image.format,
                    "mode": image.mode,
                    "ocr_enabled": bool(include_ocr),
                })
                if include_ocr:
                    try:
                        import pytesseract
                        text = pytesseract.image_to_string(image)
                        meta["ocr_chars"] = len(text)
                    except Exception as exc:  # pragma: no cover - depends on local OCR runtime
                        meta["ocr_error"] = str(exc)
                        text = ""
        return text[:200_000], meta

    def _ensure_local_hierarchy(
        self,
        conn: sqlite3.Connection,
        *,
        source_id: str,
        root: Path,
        file_path: Path,
        os_type: str,
        drive_id: str,
    ) -> str:
        computer_label = platform.node() or "내 컴퓨터"
        computer_id = f"computer:{_slug(computer_label)}"
        drive_node_id = f"drive:{_sha256_text(f'{os_type}:{drive_id}')[:24]}"
        root_folder_id = f"folder:{_sha256_text(f'{source_id}:root')[:24]}"
        self._upsert_node(conn, computer_id, "Computer", computer_label, metadata={"os_type": os_type})
        self._upsert_node(conn, drive_node_id, "Drive", drive_id, metadata={"os_type": os_type, "drive_id": drive_id})
        self._upsert_edge(conn, computer_id, drive_node_id, "포함함", metadata={"source": "local_scan"})
        self._upsert_node(
            conn,
            root_folder_id,
            "Folder",
            root.name or str(root),
            summary=str(root),
            metadata={"source_id": source_id, "path": str(root), "root": True},
        )
        self._upsert_edge(conn, drive_node_id, root_folder_id, "포함함", metadata={"source": "local_scan"})

        try:
            relative_parent = file_path.parent.relative_to(root)
        except ValueError:
            relative_parent = Path()
        parent_id = root_folder_id
        current_path = root
        for part in relative_parent.parts:
            current_path = current_path / part
            folder_id = f"folder:{_sha256_text(f'{source_id}:{current_path.as_posix()}')[:24]}"
            self._upsert_node(
                conn,
                folder_id,
                "Folder",
                part,
                summary=str(current_path),
                metadata={"source_id": source_id, "path": str(current_path), "root": False},
            )
            self._upsert_edge(conn, parent_id, folder_id, "포함함", metadata={"source": "local_scan"})
            parent_id = folder_id
        return parent_id

    def _upsert_local_file_index(
        self,
        conn: sqlite3.Connection,
        *,
        source_id: str,
        root: Path,
        file_path: Path,
        stat: Optional[os.stat_result],
        os_type: str,
        drive_id: str,
        status: str,
        parser_type: str,
        sha256: Optional[str] = None,
        graph_node_id: Optional[str] = None,
        error_message: Optional[str] = None,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> str:
        try:
            relative_path = file_path.relative_to(root).as_posix()
        except ValueError:
            relative_path = file_path.name
        index_id = f"local-index:{_sha256_text(f'{source_id}:{relative_path}')[:24]}"
        now = _now()
        size = stat.st_size if stat else None
        modified_at = _safe_iso_from_stat_mtime(stat.st_mtime) if stat else ""
        conn.execute(
            """
            INSERT INTO local_file_index(
              id, source_id, os_type, drive_id, root_path, file_path, relative_path,
              file_name, extension, size_bytes, modified_at, sha256, last_scanned_at,
              last_indexed_at, parser_type, status, error_message, graph_node_id,
              deleted, metadata_json
            )
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            ON CONFLICT(source_id, relative_path) DO UPDATE SET
              os_type=excluded.os_type,
              drive_id=excluded.drive_id,
              root_path=excluded.root_path,
              file_path=excluded.file_path,
              file_name=excluded.file_name,
              extension=excluded.extension,
              size_bytes=excluded.size_bytes,
              modified_at=excluded.modified_at,
              sha256=excluded.sha256,
              last_scanned_at=excluded.last_scanned_at,
              last_indexed_at=excluded.last_indexed_at,
              parser_type=excluded.parser_type,
              status=excluded.status,
              error_message=excluded.error_message,
              graph_node_id=excluded.graph_node_id,
              deleted=excluded.deleted,
              metadata_json=excluded.metadata_json
            """,
            (
                index_id, source_id, os_type, drive_id, str(root), str(file_path), relative_path,
                file_path.name, file_path.suffix.lower(), size, modified_at, sha256, now,
                now if status == "indexed" else None, parser_type, status, error_message,
                graph_node_id, 0 if status != "deleted" else 1, _json(metadata),
            ),
        )
        return index_id

    def _delete_local_file_graph(self, conn: sqlite3.Connection, file_node_id: Optional[str]) -> None:
        if not file_node_id:
            return

        file_row = conn.execute(
            "SELECT metadata_json FROM nodes WHERE id=?",
            (file_node_id,),
        ).fetchone()
        source_id = None
        if file_row:
            source_id = _safe_loads(file_row["metadata_json"]).get("source_id")

        linked_rows = conn.execute(
            """
            SELECT n.id, n.type, n.metadata_json
            FROM edges e
            JOIN nodes n ON n.id=e.to_node
            WHERE e.from_node=?
            """,
            (file_node_id,),
        ).fetchall()
        owned_ids: set = set()
        auto_candidate_ids: set = set()
        for row in linked_rows:
            metadata = _safe_loads(row["metadata_json"])
            if row["type"] in {"Chunk", "ImageText", "Section"} or metadata.get("source_node") == file_node_id:
                owned_ids.add(row["id"])
            elif metadata.get("auto_extracted") and metadata.get("source") == "local_folder":
                auto_candidate_ids.add(row["id"])

        conn.execute("DELETE FROM chunks WHERE source_node=?", (file_node_id,))
        conn.execute("DELETE FROM edges WHERE from_node=? OR to_node=?", (file_node_id, file_node_id))
        conn.execute("DELETE FROM nodes WHERE id=?", (file_node_id,))
        self._v2_delete_nodes(conn, [file_node_id])

        def delete_nodes(node_ids: set) -> None:
            if not node_ids:
                return
            placeholders = ",".join("?" * len(node_ids))
            params = list(node_ids)
            conn.execute(f"DELETE FROM chunks WHERE source_node IN ({placeholders})", params)
            conn.execute(f"DELETE FROM edges WHERE from_node IN ({placeholders}) OR to_node IN ({placeholders})", params * 2)
            conn.execute(f"DELETE FROM nodes WHERE id IN ({placeholders})", params)
            self._v2_delete_nodes(conn, params)

        delete_nodes(owned_ids)

        removable_auto_ids: set = set()
        for node_id in auto_candidate_ids:
            remaining_edges = conn.execute(
                "SELECT from_node, to_node FROM edges WHERE from_node=? OR to_node=?",
                (node_id, node_id),
            ).fetchall()
            if all(
                (row["from_node"] in auto_candidate_ids and row["to_node"] in auto_candidate_ids)
                for row in remaining_edges
            ):
                removable_auto_ids.add(node_id)
        delete_nodes(removable_auto_ids)
        if source_id:
            self._cleanup_local_graph_orphans(conn, str(source_id))

    def _cleanup_local_graph_orphans(self, conn: sqlite3.Connection, source_id: str) -> None:
        while True:
            folder_rows = conn.execute(
                "SELECT id, metadata_json FROM nodes WHERE type='Folder'"
            ).fetchall()
            leaf_ids = []
            for row in folder_rows:
                metadata = _safe_loads(row["metadata_json"])
                if metadata.get("source_id") != source_id:
                    continue
                has_children = conn.execute(
                    "SELECT 1 FROM edges WHERE from_node=? LIMIT 1",
                    (row["id"],),
                ).fetchone()
                if not has_children:
                    leaf_ids.append(row["id"])
            if not leaf_ids:
                break
            placeholders = ",".join("?" * len(leaf_ids))
            conn.execute(f"DELETE FROM edges WHERE from_node IN ({placeholders}) OR to_node IN ({placeholders})", leaf_ids * 2)
            conn.execute(f"DELETE FROM nodes WHERE id IN ({placeholders})", leaf_ids)
            self._v2_delete_nodes(conn, leaf_ids)

        for node_type in ("Drive", "Computer"):
            rows = conn.execute("SELECT id FROM nodes WHERE type=?", (node_type,)).fetchall()
            removable = []
            for row in rows:
                has_children = conn.execute(
                    "SELECT 1 FROM edges WHERE from_node=? LIMIT 1",
                    (row["id"],),
                ).fetchone()
                if not has_children:
                    removable.append(row["id"])
            if removable:
                placeholders = ",".join("?" * len(removable))
                conn.execute(f"DELETE FROM edges WHERE from_node IN ({placeholders}) OR to_node IN ({placeholders})", removable * 2)
                conn.execute(f"DELETE FROM nodes WHERE id IN ({placeholders})", removable)
                self._v2_delete_nodes(conn, removable)

    def _local_file_index_has_extracted_text(self, row: sqlite3.Row) -> bool:
        metadata = _safe_loads(row["metadata_json"])
        parser = metadata.get("parser") if isinstance(metadata, dict) else {}
        if not isinstance(parser, dict):
            return False
        try:
            return int(parser.get("extracted_chars") or 0) > 0
        except (TypeError, ValueError):
            return False

    def _upsert_local_file_node(
        self,
        conn: sqlite3.Connection,
        *,
        source_id: str,
        root: Path,
        file_path: Path,
        stat: os.stat_result,
        os_type: str,
        drive_id: str,
        sha256: str,
        category: str,
        parser_type: str,
        text: str,
        parser_meta: Dict[str, Any],
    ) -> str:
        text = _clean_text(text)
        if not text:
            raise ValueError("텍스트 추출 결과가 비어 있습니다.")
        try:
            relative_path = file_path.relative_to(root).as_posix()
        except ValueError:
            relative_path = file_path.name
        file_node_id = f"local-file:{_sha256_text(f'{source_id}:{relative_path}')[:24]}"
        parent_folder_id = self._ensure_local_hierarchy(
            conn,
            source_id=source_id,
            root=root,
            file_path=file_path,
            os_type=os_type,
            drive_id=drive_id,
        )
        child_rows = conn.execute(
            """
            SELECT e.to_node AS id
            FROM edges e
            JOIN nodes n ON n.id=e.to_node
            WHERE e.from_node=? AND n.type IN ('Chunk', 'ImageText', 'Section')
            """,
            (file_node_id,),
        ).fetchall()
        child_ids = [row["id"] for row in child_rows]
        conn.execute("DELETE FROM chunks WHERE source_node=?", (file_node_id,))
        if child_ids:
            placeholders = ",".join("?" * len(child_ids))
            conn.execute(f"DELETE FROM nodes WHERE id IN ({placeholders})", child_ids)
            self._v2_delete_nodes(conn, child_ids)
        conn.execute("DELETE FROM edges WHERE from_node=?", (file_node_id,))
        self._v2_delete_edges_from(conn, file_node_id)

        metadata = {
            "source": "local_folder",
            "source_id": source_id,
            "root_path": str(root),
            "file_path": str(file_path),
            "relative_path": relative_path,
            "filename": file_path.name,
            "ext": file_path.suffix.lower(),
            "category": category,
            "parser_type": parser_type,
            "bytes": stat.st_size,
            "modified_at": _safe_iso_from_stat_mtime(stat.st_mtime),
            "sha256": sha256,
            "parser": parser_meta,
        }
        self._upsert_node(
            conn,
            file_node_id,
            _node_type_for_category(category),
            file_path.name,
            summary=text[:700],
            metadata=metadata,
            raw=metadata,
        )
        self._upsert_edge(conn, parent_folder_id, file_node_id, "포함함", weight=1.0, metadata={"source": "local_scan"})

        target_for_concepts = text
        if category == "image" and text:
            image_text_id = f"imagetext:{_sha256_text(f'{file_node_id}:ocr')[:24]}"
            self._upsert_node(
                conn,
                image_text_id,
                "ImageText",
                f"{file_path.name} OCR",
                summary=_clean_text(text)[:700],
                metadata={"source_node": file_node_id, "source_id": source_id, "chars": len(text)},
            )
            self._upsert_edge(conn, file_node_id, image_text_id, "포함함", weight=0.8, metadata={"source": "ocr"})

        for index, chunk in enumerate(_chunks(text)):
            chunk_id = f"chunk:{_sha256_text(f'{file_node_id}:{index}:{chunk}')[:24]}"
            self._upsert_node(
                conn,
                chunk_id,
                "Chunk",
                f"{file_path.name} chunk {index + 1}",
                summary=chunk[:500],
                metadata={"index": index, "source_node": file_node_id, "source_id": source_id},
            )
            self._upsert_chunk(
                conn,
                chunk_id=chunk_id,
                source_node=file_node_id,
                text=chunk,
                metadata={"index": index, "source_node": file_node_id, "source_id": source_id},
            )
            self._upsert_edge(conn, file_node_id, chunk_id, "포함함", weight=0.7, metadata={"source": "local_scan"})

        concepts = _extract_concepts(target_for_concepts, limit=18)
        concept_ids: Dict[str, str] = {}
        for concept in concepts:
            node_t = _classify_node_type(concept, target_for_concepts)
            concept_id = f"{node_t.lower()}:{_slug(concept)}"
            concept_ids[concept.lower()] = concept_id
            self._upsert_node(
                conn,
                concept_id,
                node_t,
                concept,
                metadata={"auto_extracted": True, "source": "local_folder", "source_id": source_id},
            )
            self._upsert_edge(conn, file_node_id, concept_id, "언급함", weight=0.75, metadata={"source": "local_scan"})

        for triple in _extract_triples(target_for_concepts, concepts, limit=20):
            subj_id = concept_ids.get(triple["subject"].lower())
            obj_id = concept_ids.get(triple["object"].lower())
            if subj_id and obj_id and subj_id != obj_id:
                self._upsert_edge(
                    conn,
                    subj_id,
                    obj_id,
                    triple["relation"],
                    weight=0.9,
                    metadata={"context": triple.get("context", "")[:240], "source_id": source_id},
                )

        for item in _semantic_items(target_for_concepts):
            sem_type = item["type"]
            sem_title = item["title"]
            sem_id = f"{sem_type.lower()}:{_sha256_text(f'{file_node_id}:{sem_type}:{sem_title}')[:24]}"
            self._upsert_node(
                conn,
                sem_id,
                sem_type,
                sem_title,
                summary=item["summary"],
                metadata={"auto_extracted": True, "source_node": file_node_id, "filename": file_path.name},
                raw=item,
            )
            self._upsert_edge(conn, file_node_id, sem_id, "포함함", weight=0.9)

        return file_node_id

    def index_local_folder(
        self,
        path: Path,
        *,
        include_ocr: bool = False,
        watch_enabled: bool = False,
        user_email: Optional[str] = None,
        consent: Optional[Dict[str, Any]] = None,
        max_files: int = 5_000,
    ) -> Dict[str, Any]:
        """Read approved files from a local folder and connect them to Graph RAG."""
        root = Path(path).expanduser().resolve()
        if not root.exists():
            raise ValueError(f"경로가 존재하지 않습니다: {path}")
        if not root.is_dir():
            raise ValueError(f"폴더가 아닙니다: {path}")

        os_type = _current_os_type()
        drive_id = _drive_id_for_path(root)
        source_id = f"source:{_path_fingerprint(root)}"
        now = _now()
        max_files = max(1, min(int(max_files or 5_000), 50_000))
        consent_payload = {
            "approved_at": now,
            "approved_by": user_email,
            "knowledge_source": True,
            "include_ocr": bool(include_ocr),
            "watch_enabled": bool(watch_enabled),
            "sensitive_files_default_excluded": True,
            **(consent or {}),
        }
        counts: Counter = Counter()
        seen_relative_paths: set = set()
        indexed_nodes: List[str] = []
        errors: List[Dict[str, str]] = []
        limit_reached = False

        with self._connect() as conn:
            conn.execute(
                """
                INSERT INTO knowledge_sources(
                  id, root_path, os_type, drive_id, label, status, include_ocr,
                  watch_enabled, consent_json, created_at, updated_at, last_scanned_at
                )
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                ON CONFLICT(id) DO UPDATE SET
                  root_path=excluded.root_path,
                  os_type=excluded.os_type,
                  drive_id=excluded.drive_id,
                  label=excluded.label,
                  status=excluded.status,
                  include_ocr=excluded.include_ocr,
                  watch_enabled=excluded.watch_enabled,
                  consent_json=excluded.consent_json,
                  updated_at=excluded.updated_at,
                  last_scanned_at=excluded.last_scanned_at
                """,
                (
                    source_id, str(root), os_type, drive_id, root.name or str(root), "scanning",
                    1 if include_ocr else 0, 1 if watch_enabled else 0, _json(consent_payload),
                    now, now, now,
                ),
            )

            for entry in self._iter_local_scan_entries(root, max_files=max_files):
                kind = entry["kind"]
                file_path = entry["path"]
                if kind == "limit_reached":
                    counts["limit_reached"] += 1
                    limit_reached = True
                    break
                if kind in {"excluded_dir", "excluded"}:
                    counts["excluded"] += 1
                    continue
                if kind in {"inaccessible_dir", "inaccessible_file"}:
                    counts["failed"] += 1
                    errors.append({"path": str(file_path), "error": entry.get("reason", "inaccessible")})
                    continue
                if kind != "file":
                    continue

                stat = entry["stat"]
                try:
                    relative_path = file_path.relative_to(root).as_posix()
                except ValueError:
                    relative_path = file_path.name
                seen_relative_paths.add(relative_path)
                modified_at = _safe_iso_from_stat_mtime(stat.st_mtime)
                existing = conn.execute(
                    """
                    SELECT size_bytes, modified_at, sha256, graph_node_id, status, metadata_json
                    FROM local_file_index
                    WHERE source_id=? AND relative_path=?
                    """,
                    (source_id, relative_path),
                ).fetchone()
                decision = self._local_file_decision(file_path, root, stat)
                parser_type = decision["parser_type"]
                if not decision["indexable"]:
                    counts[decision["status"]] += 1
                    if existing and existing["graph_node_id"]:
                        self._delete_local_file_graph(conn, existing["graph_node_id"])
                    self._upsert_local_file_index(
                        conn,
                        source_id=source_id,
                        root=root,
                        file_path=file_path,
                        stat=stat,
                        os_type=os_type,
                        drive_id=drive_id,
                        status=decision["status"],
                        parser_type=parser_type,
                        metadata={"reason": decision["reason"], "category": decision["category"]},
                    )
                    continue

                if (
                    existing
                    and existing["status"] == "indexed"
                    and existing["graph_node_id"]
                    and self._local_file_index_has_extracted_text(existing)
                    and existing["size_bytes"] == stat.st_size
                    and existing["modified_at"] == modified_at
                ):
                    counts["skipped_unchanged"] += 1
                    self._upsert_local_file_index(
                        conn,
                        source_id=source_id,
                        root=root,
                        file_path=file_path,
                        stat=stat,
                        os_type=os_type,
                        drive_id=drive_id,
                        status="indexed",
                        parser_type=parser_type,
                        sha256=existing["sha256"],
                        graph_node_id=existing["graph_node_id"],
                        metadata={**_safe_loads(existing["metadata_json"]), "category": decision["category"], "unchanged": True},
                    )
                    continue

                try:
                    data = file_path.read_bytes()
                    digest = _sha256_bytes(data)
                except Exception as exc:
                    counts["failed"] += 1
                    errors.append({"path": str(file_path), "error": str(exc)})
                    if existing and existing["graph_node_id"]:
                        self._delete_local_file_graph(conn, existing["graph_node_id"])
                    self._upsert_local_file_index(
                        conn,
                        source_id=source_id,
                        root=root,
                        file_path=file_path,
                        stat=stat,
                        os_type=os_type,
                        drive_id=drive_id,
                        status="failed",
                        parser_type=parser_type,
                        error_message=str(exc),
                        metadata={"category": decision["category"]},
                    )
                    continue

                if (
                    existing
                    and existing["sha256"] == digest
                    and existing["graph_node_id"]
                    and self._local_file_index_has_extracted_text(existing)
                ):
                    counts["skipped_unchanged"] += 1
                    self._upsert_local_file_index(
                        conn,
                        source_id=source_id,
                        root=root,
                        file_path=file_path,
                        stat=stat,
                        os_type=os_type,
                        drive_id=drive_id,
                        status="indexed",
                        parser_type=parser_type,
                        sha256=digest,
                        graph_node_id=existing["graph_node_id"],
                        metadata={**_safe_loads(existing["metadata_json"]), "category": decision["category"], "sha256_unchanged": True},
                    )
                    continue

                try:
                    text, parser_meta = self._extract_local_file_text(
                        file_path,
                        decision["category"],
                        include_ocr=include_ocr,
                    )
                    text = _clean_text(text)
                    parser_meta = {**parser_meta, "extracted_chars": len(text)}
                    if not text:
                        counts["skipped_empty_text"] += 1
                        if existing and existing["graph_node_id"]:
                            self._delete_local_file_graph(conn, existing["graph_node_id"])
                        self._upsert_local_file_index(
                            conn,
                            source_id=source_id,
                            root=root,
                            file_path=file_path,
                            stat=stat,
                            os_type=os_type,
                            drive_id=drive_id,
                            status="skipped_empty_text",
                            parser_type=parser_type,
                            sha256=digest,
                            error_message="텍스트 추출 결과가 비어 있습니다.",
                            metadata={"category": decision["category"], "parser": parser_meta},
                        )
                        continue
                    graph_node_id = self._upsert_local_file_node(
                        conn,
                        source_id=source_id,
                        root=root,
                        file_path=file_path,
                        stat=stat,
                        os_type=os_type,
                        drive_id=drive_id,
                        sha256=digest,
                        category=decision["category"],
                        parser_type=parser_type,
                        text=text,
                        parser_meta=parser_meta,
                    )
                    self._upsert_local_file_index(
                        conn,
                        source_id=source_id,
                        root=root,
                        file_path=file_path,
                        stat=stat,
                        os_type=os_type,
                        drive_id=drive_id,
                        status="indexed",
                        parser_type=parser_type,
                        sha256=digest,
                        graph_node_id=graph_node_id,
                        metadata={"category": decision["category"], "parser": parser_meta},
                    )
                    counts["indexed"] += 1
                    indexed_nodes.append(graph_node_id)
                except Exception as exc:
                    counts["failed"] += 1
                    errors.append({"path": str(file_path), "error": str(exc)})
                    if existing and existing["graph_node_id"]:
                        self._delete_local_file_graph(conn, existing["graph_node_id"])
                    self._upsert_local_file_index(
                        conn,
                        source_id=source_id,
                        root=root,
                        file_path=file_path,
                        stat=stat,
                        os_type=os_type,
                        drive_id=drive_id,
                        status="failed",
                        parser_type=parser_type,
                        sha256=digest,
                        error_message=str(exc),
                        metadata={"category": decision["category"]},
                    )

            if not limit_reached:
                existing_rows = {
                    row["relative_path"]: row["graph_node_id"]
                    for row in conn.execute(
                        "SELECT relative_path, graph_node_id FROM local_file_index WHERE source_id=?",
                        (source_id,),
                    )
                }
                deleted_paths = set(existing_rows) - seen_relative_paths
                for relative_path in deleted_paths:
                    self._delete_local_file_graph(conn, existing_rows.get(relative_path))
                    conn.execute(
                        """
                        UPDATE local_file_index
                        SET status='deleted', deleted=1, last_scanned_at=?, error_message=NULL, graph_node_id=NULL
                        WHERE source_id=? AND relative_path=?
                        """,
                        (_now(), source_id, relative_path),
                    )
                counts["deleted"] = len(deleted_paths)
            conn.execute(
                """
                UPDATE knowledge_sources
                SET status='active', updated_at=?, last_scanned_at=?
                WHERE id=?
                """,
                (_now(), _now(), source_id),
            )

        return {
            "status": "ok",
            "source": {
                "id": source_id,
                "root_path": str(root),
                "os_type": os_type,
                "drive_id": drive_id,
                "include_ocr": bool(include_ocr),
                "watch_enabled": bool(watch_enabled),
            },
            "counts": dict(counts),
            "indexed_nodes": indexed_nodes[:100],
            "errors": errors[:50],
            "notice": "Lattice AI는 사용자가 선택한 폴더만 AI 지식으로 변환합니다.",
        }

    def ingest_message(
        self,
        role: str,
        content: str,
        *,
        user_email: Optional[str] = None,
        user_nickname: Optional[str] = None,
        source: Optional[str] = None,
        conversation_id: Optional[str] = None,
        raw: Optional[Dict[str, Any]] = None,
    ) -> Dict[str, Any]:
        content = str(content or "")
        digest = _sha256_text("|".join([role or "", content, conversation_id or "", user_email or ""]))[:24]
        node_type = "AIResponse" if role == "assistant" else "Message"
        node_id = f"{node_type.lower()}:{digest}"
        conv_id = f"conversation:{_slug(conversation_id or 'default')}"
        metadata = {
            "role": role,
            "source": source,
            "conversation_id": conversation_id,
            "user_email": user_email,
            "user_nickname": user_nickname,
            "chars": len(content),
        }
        concepts = _extract_concepts(content)
        triples  = _extract_triples(content, concepts)
        semantic = _semantic_items(content)

        with self._connect() as conn:
            # ── 1. Chat node  (점: 명사 — 대화 세션 단위) ─────────────────────
            #    One Chat node per conversation_id; title = first 80 chars of
            #    the first user message in this session (updated on each call).
            chat_title = _clean_text(content)[:80] or (conversation_id or "대화")
            self._upsert_node(
                conn, conv_id, "Chat",
                chat_title,
                summary=_clean_text(content)[:400],
                metadata={"source": source, "conversation_id": conversation_id},
            )

            # ── 2. Person node  (점: 명사 — 사람) ─────────────────────────────
            person_id = None
            if user_email or user_nickname:
                person_key = user_email or user_nickname or "unknown"
                person_id = f"person:{_slug(person_key)}"
                self._upsert_node(
                    conn, person_id, "Person",
                    user_nickname or user_email or "Unknown",
                    metadata={"email": user_email, "nickname": user_nickname},
                )
                # 선: 동사 — Person이 Chat을 "작성함"
                self._upsert_edge(conn, person_id, conv_id, "작성함",
                                  weight=1.0, metadata={"role": role})

            # ── 3. Raw message node  (RAG 검색용, 그래프에서 숨김) ─────────────
            self._upsert_node(
                conn, node_id, node_type,
                _clean_text(content)[:80] or role,
                summary=_clean_text(content)[:500],
                metadata=metadata,
                raw=raw or metadata,
            )
            # 선: Chat이 메시지를 "포함함"
            self._upsert_edge(conn, conv_id, node_id, "포함함",
                              weight=0.3, metadata={"role": role})

            # ── 4. RAG chunks  (검색용, 그래프에서 숨김) ──────────────────────
            for index, chunk in enumerate(_chunks(content)):
                chunk_id = f"chunk:{_sha256_text(f'{node_id}:{index}:{chunk}')[:24]}"
                self._upsert_node(
                    conn, chunk_id, "Chunk",
                    f"chunk {index + 1}",
                    summary=chunk[:500],
                    metadata={"index": index, "source_node": node_id},
                )
                self._upsert_chunk(
                    conn,
                    chunk_id=chunk_id,
                    source_node=node_id,
                    text=chunk,
                    metadata={"index": index, "source_node": node_id},
                )
                self._upsert_edge(conn, node_id, chunk_id, "포함함")

            # ── 5. Concept / Feature / Error / Code 노드  (점: 명사) ───────────
            concept_ids: Dict[str, str] = {}
            for concept in concepts:
                node_t = _classify_node_type(concept, content)
                cid = f"{node_t.lower()}:{_slug(concept)}"
                concept_ids[concept.lower()] = cid
                self._upsert_node(
                    conn, cid, node_t, concept,
                    metadata={"auto_extracted": True, "source": source},
                )
                # 선: Chat이 개념을 "언급함"
                self._upsert_edge(conn, conv_id, cid, "언급함",
                                  weight=0.7, metadata={"source": source})

            # ── 6. Concept–Concept 엣지  (선: 동사형) ─────────────────────────
            for triple in triples:
                subj_id = concept_ids.get(triple["subject"].lower())
                obj_id  = concept_ids.get(triple["object"].lower())
                if subj_id and obj_id and subj_id != obj_id:
                    self._upsert_edge(
                        conn, subj_id, obj_id,
                        triple["relation"],          # 동사형 레이블
                        weight=1.0,
                        metadata={"context": triple.get("context", "")[:240]},
                    )

            # ── 7. Task / Decision 노드  (점: 명사) ────────────────────────────
            for item in semantic:
                sem_type  = item["type"]
                sem_title = item["title"]
                sem_id = f"{sem_type.lower()}:{_sha256_text(f'{conv_id}:{sem_type}:{sem_title}')[:24]}"
                self._upsert_node(
                    conn, sem_id, sem_type, sem_title,
                    summary=item["summary"],
                    metadata={"auto_extracted": True, "source_node": node_id},
                    raw=item,
                )
                # 선: Chat이 Task/Decision을 "생성함"
                self._upsert_edge(conn, conv_id, sem_id, "생성함", weight=0.9)
                # Task/Decision이 관련 개념을 "언급함"
                for cid in list(concept_ids.values())[:3]:
                    self._upsert_edge(conn, sem_id, cid, "언급함", weight=0.6)

        return {"node_id": node_id, "type": node_type}

    def ingest_document(
        self,
        path: Path,
        *,
        original_filename: Optional[str] = None,
        mime_type: Optional[str] = None,
        uploader: Optional[str] = None,
        conversation_id: Optional[str] = None,
        extracted: Optional[Dict[str, Any]] = None,
    ) -> Dict[str, Any]:
        path = Path(path)
        data = path.read_bytes()
        digest = _sha256_bytes(data)
        ext = path.suffix.lower()
        filename = original_filename or path.name
        blob_path = self.blob_dir / digest[:2] / f"{digest}{ext}"
        blob_path.parent.mkdir(parents=True, exist_ok=True)
        if not blob_path.exists():
            shutil.copyfile(path, blob_path)

        doc_meta = self._document_structure(path, ext)
        text = str((extracted or {}).get("content") or (extracted or {}).get("preview") or "")
        file_id = f"file:{digest[:24]}"
        metadata = {
            "filename": filename,
            "ext": ext,
            "mime_type": mime_type,
            "bytes": len(data),
            "sha256": digest,
            "blob_path": str(blob_path),
            "uploader": uploader,
            "conversation_id": conversation_id,
            "extracted": {k: v for k, v in (extracted or {}).items() if k != "content"},
            "structure": doc_meta,
        }
        full_text = f"{filename}\n{text}"
        concepts = _extract_concepts(full_text, limit=15)
        triples  = _extract_triples(full_text, concepts)

        with self._connect() as conn:
            # ── Document 노드  (점: 명사 — 파일) ────────────────────────────────
            self._upsert_node(
                conn, file_id, "Document", filename,
                summary=(text or filename)[:500],
                metadata=metadata, raw=metadata,
            )
            self._ingest_structure_nodes(conn, file_id, filename, doc_meta)

            # ── Person 노드 + 동사형 엣지 ─────────────────────────────────────
            if uploader:
                person_id = f"person:{_slug(uploader)}"
                self._upsert_node(
                    conn, person_id, "Person", uploader,
                    metadata={"email": uploader},
                )
                # 선: 동사 — Person이 Document를 "업로드함"
                self._upsert_edge(conn, person_id, file_id, "업로드함", weight=1.0)

            # ── Chat 노드와 연결 ──────────────────────────────────────────────
            if conversation_id:
                conv_id = f"conversation:{_slug(conversation_id)}"
                self._upsert_node(conn, conv_id, "Chat", conversation_id)
                # 선: 동사 — Chat이 Document를 "언급함"
                self._upsert_edge(conn, conv_id, file_id, "언급함", weight=0.8)

            # ── RAG chunks (검색용, 그래프 비표시) ────────────────────────────
            for index, chunk in enumerate(_chunks(text)):
                chunk_id = f"chunk:{_sha256_text(f'{file_id}:{index}:{chunk}')[:24]}"
                self._upsert_node(
                    conn, chunk_id, "Chunk",
                    f"{filename} chunk {index + 1}",
                    summary=chunk[:500],
                    metadata={"index": index, "source_node": file_id},
                )
                self._upsert_chunk(
                    conn,
                    chunk_id=chunk_id,
                    source_node=file_id,
                    text=chunk,
                    metadata={"index": index, "source_node": file_id},
                )
                self._upsert_edge(conn, file_id, chunk_id, "포함함")

            # ── Concept / Feature / Error / Code 노드 + 동사형 엣지 ───────────
            concept_ids: Dict[str, str] = {}
            for concept in concepts:
                node_t = _classify_node_type(concept, full_text)
                cid = f"{node_t.lower()}:{_slug(concept)}"
                concept_ids[concept.lower()] = cid
                self._upsert_node(
                    conn, cid, node_t, concept,
                    metadata={"auto_extracted": True, "source_file": filename},
                )
                # 선: 동사 — Document가 Concept을 "포함함"
                self._upsert_edge(conn, file_id, cid, "포함함", weight=0.8)

            # ── Concept–Concept 엣지  (선: 동사형) ───────────────────────────
            for triple in triples:
                subj_id = concept_ids.get(triple["subject"].lower())
                obj_id  = concept_ids.get(triple["object"].lower())
                if subj_id and obj_id and subj_id != obj_id:
                    self._upsert_edge(
                        conn, subj_id, obj_id,
                        triple["relation"],
                        weight=1.0,
                        metadata={"context": triple.get("context", "")[:240]},
                    )

            # ── Task / Decision 노드 ──────────────────────────────────────────
            for item in _semantic_items(text):
                sem_type  = item["type"]
                sem_title = item["title"]
                sem_id = f"{sem_type.lower()}:{_sha256_text(f'{file_id}:{sem_type}:{sem_title}')[:24]}"
                self._upsert_node(
                    conn, sem_id, sem_type, sem_title,
                    summary=item["summary"],
                    metadata={"auto_extracted": True, "source_node": file_id, "filename": filename},
                    raw=item,
                )
                # 선: Document가 Task/Decision을 "포함함"
                self._upsert_edge(conn, file_id, sem_id, "포함함", weight=0.9)

        return {"node_id": file_id, "sha256": digest, "metadata": metadata}

    def ingest_event(
        self,
        event_type: str,
        title: str,
        *,
        user_email: Optional[str] = None,
        user_nickname: Optional[str] = None,
        source: Optional[str] = None,
        conversation_id: Optional[str] = None,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> Dict[str, Any]:
        event_type = str(event_type or "Event")
        title = str(title or event_type)
        payload = {
            "event_type": event_type,
            "title": title,
            "user_email": user_email,
            "user_nickname": user_nickname,
            "source": source,
            "conversation_id": conversation_id,
            "metadata": metadata or {},
            "timestamp": _now(),
        }
        event_id = f"event:{_sha256_text(_json(payload))[:24]}"
        conv_id = f"conversation:{_slug(conversation_id or 'default')}"
        with self._connect() as conn:
            self._upsert_node(conn, event_id, event_type, title, summary=title, metadata=payload, raw=payload)
            self._upsert_node(conn, conv_id, "Conversation", conversation_id or "Default conversation", metadata={"source": source})
            self._upsert_edge(conn, conv_id, event_id, "has_event", metadata={"source": source})
            if user_email or user_nickname:
                person_key = user_email or user_nickname or "unknown"
                person_id = f"person:{_slug(person_key)}"
                self._upsert_node(conn, person_id, "Person", user_nickname or user_email or "Unknown user", metadata={"email": user_email})
                self._upsert_edge(conn, person_id, event_id, "triggered", metadata={"event_type": event_type})
        return {"node_id": event_id, "type": event_type}

    def _ingest_structure_nodes(
        self,
        conn: sqlite3.Connection,
        file_id: str,
        filename: str,
        structure: Dict[str, Any],
    ) -> None:
        for slide in structure.get("slides") or []:
            index = slide.get("index")
            slide_id = f"slide:{_sha256_text(f'{file_id}:slide:{index}')[:24]}"
            title = f"{filename} slide {index}"
            summary = "\n".join(slide.get("texts") or [])[:800]
            self._upsert_node(conn, slide_id, "Slide", title, summary=summary, metadata=slide)
            self._upsert_edge(conn, file_id, slide_id, "has_slide")
            for text in slide.get("texts") or []:
                for topic in _topic_candidates(text, limit=4):
                    topic_id = f"topic:{_slug(topic)}"
                    self._upsert_node(conn, topic_id, "Topic", topic, metadata={"auto_extracted": True})
                    self._upsert_edge(conn, slide_id, topic_id, "discusses", weight=0.6)

        for page in structure.get("pages") or []:
            index = page.get("index")
            page_id = f"page:{_sha256_text(f'{file_id}:page:{index}')[:24]}"
            title = f"{filename} page {index}"
            self._upsert_node(conn, page_id, "Page", title, summary=page.get("preview") or "", metadata=page)
            self._upsert_edge(conn, file_id, page_id, "has_page")
            for topic in _topic_candidates(page.get("preview") or "", limit=4):
                topic_id = f"topic:{_slug(topic)}"
                self._upsert_node(conn, topic_id, "Topic", topic, metadata={"auto_extracted": True})
                self._upsert_edge(conn, page_id, topic_id, "discusses", weight=0.6)

        for sheet in (structure.get("sheets") or []):
            sheet_title = sheet.get("title")
            sheet_id = f"sheet:{_sha256_text(f'{file_id}:sheet:{sheet_title}')[:24]}"
            self._upsert_node(conn, sheet_id, "Sheet", f"{filename} / {sheet_title}", metadata=sheet)
            self._upsert_edge(conn, file_id, sheet_id, "has_sheet")

        for image in (structure.get("images") or []):
            image_key = image.get("sha256") or _sha256_text(json.dumps(image, ensure_ascii=False, sort_keys=True))
            image_id = f"image:{str(image_key)[:24]}"
            title_parts = [filename, "image"]
            if image.get("page"):
                title_parts.append(f"page {image.get('page')}")
            if image.get("name"):
                title_parts.append(str(image.get("name")).split("/")[-1])
            self._upsert_node(conn, image_id, "Image", " / ".join(title_parts), metadata=image)
            self._upsert_edge(conn, file_id, image_id, "contains_image")

    def _document_structure(self, path: Path, ext: str) -> Dict[str, Any]:
        try:
            if ext == ".pptx":
                return self._pptx_structure(path)
            if ext == ".pdf":
                return self._pdf_structure(path)
            if ext == ".docx":
                return self._docx_structure(path)
            if ext == ".xlsx":
                return self._xlsx_structure(path)
        except Exception as exc:
            return {"error": str(exc)}
        return {}

    def _pptx_structure(self, path: Path) -> Dict[str, Any]:
        result: Dict[str, Any] = {"slides": [], "images": []}
        try:
            from PIL import Image
            from pptx import Presentation
            prs = Presentation(str(path))
            for slide_index, slide in enumerate(prs.slides, start=1):
                slide_info = {"index": slide_index, "shapes": [], "texts": []}
                for shape_index, shape in enumerate(slide.shapes, start=1):
                    shape_info = {
                        "index": shape_index,
                        "name": getattr(shape, "name", ""),
                        "shape_type": str(getattr(shape, "shape_type", "")),
                        "bbox": {
                            "left": int(getattr(shape, "left", 0) or 0),
                            "top": int(getattr(shape, "top", 0) or 0),
                            "width": int(getattr(shape, "width", 0) or 0),
                            "height": int(getattr(shape, "height", 0) or 0),
                        },
                    }
                    if getattr(shape, "has_text_frame", False):
                        text = shape.text_frame.text.strip()
                        if text:
                            shape_info["text"] = text[:1000]
                            slide_info["texts"].append(text)
                    slide_info["shapes"].append(shape_info)
                result["slides"].append(slide_info)
            with zipfile.ZipFile(path) as zf:
                for name in zf.namelist():
                    if not name.startswith("ppt/media/"):
                        continue
                    data = zf.read(name)
                    image_info: Dict[str, Any] = {
                        "name": name,
                        "bytes": len(data),
                        "sha256": _sha256_bytes(data),
                    }
                    try:
                        from io import BytesIO
                        with Image.open(BytesIO(data)) as img:
                            image_info.update({"width": img.width, "height": img.height, "format": img.format})
                    except Exception:
                        pass
                    result["images"].append(image_info)
        except Exception as exc:
            result["error"] = str(exc)
        return result

    def _pdf_structure(self, path: Path) -> Dict[str, Any]:
        result: Dict[str, Any] = {"pages": [], "images": []}
        try:
            import pdfplumber
            with pdfplumber.open(str(path)) as pdf:
                metadata = dict(pdf.metadata or {})
                result["metadata"] = {str(k): str(v) for k, v in metadata.items()}
                for page_index, page in enumerate(pdf.pages, start=1):
                    text = page.extract_text() or ""
                    page_info = {
                        "index": page_index,
                        "width": float(page.width or 0),
                        "height": float(page.height or 0),
                        "chars": len(text),
                        "preview": _clean_text(text)[:500],
                        "image_count": len(page.images or []),
                    }
                    result["pages"].append(page_info)
                    for image_index, image in enumerate(page.images or [], start=1):
                        result["images"].append({
                            "page": page_index,
                            "index": image_index,
                            "name": image.get("name"),
                            "width": image.get("width"),
                            "height": image.get("height"),
                            "bbox": {
                                "x0": image.get("x0"),
                                "top": image.get("top"),
                                "x1": image.get("x1"),
                                "bottom": image.get("bottom"),
                            },
                        })
        except Exception as exc:
            result["error"] = str(exc)
        return result

    def _docx_structure(self, path: Path) -> Dict[str, Any]:
        from docx import Document
        doc = Document(str(path))
        headings = []
        paragraphs = 0
        for p in doc.paragraphs:
            text = p.text.strip()
            if not text:
                continue
            paragraphs += 1
            style = getattr(p.style, "name", "")
            if style.lower().startswith("heading"):
                headings.append({"style": style, "text": text[:240]})
        return {"paragraphs": paragraphs, "headings": headings[:80], "tables": len(doc.tables)}

    def _xlsx_structure(self, path: Path) -> Dict[str, Any]:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        sheets = []
        for ws in wb.worksheets:
            sheets.append({"title": ws.title, "max_row": ws.max_row, "max_column": ws.max_column})
        return {"sheets": sheets}

    # ── 그래프에 표시되는 노드 타입  (점 = 명사) ──────────────────────────────
    # Message / AIResponse / Chunk 는 RAG 검색용으로만 저장, 그래프에서 숨김.
    _GRAPH_VISIBLE_TYPES = (
        "Computer",   # 내 컴퓨터
        "Drive",      # 드라이브 / 볼륨
        "Folder",     # 폴더
        "File",       # 일반 파일
        "Chat",       # 대화 세션
        "Document",   # 파일 (PDF·PPT·Word·Excel·이미지)
        "CodeFile",   # 코드 파일
        "Spreadsheet",# 엑셀/CSV
        "SlideDeck",  # 프레젠테이션
        "Image",      # 이미지
        "ImageText",  # OCR 텍스트
        "Concept",    # 개념 / 아이디어 / 기술 용어
        "Person",     # 사람
        "Error",      # 오류 / 버그
        "Code",       # 코드 / 함수
        "Feature",    # 소프트웨어 기능
        "Task",       # 할 일
        "Decision",   # 결정 사항
    )

    def graph(self, limit: int = 300) -> Dict[str, Any]:
        limit = max(1, min(int(limit or 300), 2000))
        visible = ",".join(f"'{t}'" for t in self._GRAPH_VISIBLE_TYPES)
        nt, et = self._read_tables()
        with self._connect() as conn:
            nodes = [
                {
                    "id": row["id"],
                    "type": row["type"],
                    "title": row["title"],
                    "summary": row["summary"],
                    "metadata": _safe_loads(row["metadata_json"]),
                    "updated_at": row["updated_at"],
                }
                for row in conn.execute(
                    f"SELECT id, type, title, summary, metadata_json, updated_at FROM {nt} WHERE type IN ({visible}) ORDER BY updated_at DESC, id ASC LIMIT ?",
                    (limit,),
                )
            ]
            node_ids = {node["id"] for node in nodes}
            edges: List[Dict[str, Any]] = []
            if node_ids:
                edge_rows = conn.execute(
                    f"""
                    SELECT id, from_node, to_node, type, weight, metadata_json
                    FROM {et}
                    WHERE from_node IN (
                        SELECT id FROM {nt} WHERE type IN ({visible})
                        ORDER BY updated_at DESC, id ASC LIMIT ?
                    )
                    AND to_node IN (
                        SELECT id FROM {nt} WHERE type IN ({visible})
                        ORDER BY updated_at DESC, id ASC LIMIT ?
                    )
                    ORDER BY weight DESC, created_at DESC, id ASC
                    """,
                    (limit, limit),
                ).fetchall()
                edges = [
                    {
                        "id": row["id"],
                        "from": row["from_node"],
                        "to": row["to_node"],
                        "type": row["type"],
                        "weight": row["weight"],
                        "metadata": _safe_loads(row["metadata_json"]),
                    }
                    for row in edge_rows
                ]

        degree_map: Dict[str, int] = {}
        now = datetime.now()
        node_by_id = {node["id"]: node for node in nodes}
        topic_metrics: Dict[str, Dict[str, Any]] = {}

        for edge in edges:
            degree_map[edge["from"]] = degree_map.get(edge["from"], 0) + 1
            degree_map[edge["to"]] = degree_map.get(edge["to"], 0) + 1
            from_node = node_by_id.get(edge["from"])
            to_node = node_by_id.get(edge["to"])
            if not from_node or not to_node:
                continue
            for topic_node, other_node in ((from_node, to_node), (to_node, from_node)):
                if topic_node["type"] != "Topic":
                    continue
                metrics = topic_metrics.setdefault(topic_node["id"], {
                    "mention_count": 0.0,
                    "conversation_ids": set(),
                })
                if edge["type"] in {"mentions", "discusses"}:
                    metrics["mention_count"] += max(0.5, float(edge.get("weight") or 1.0))
                other_meta = other_node.get("metadata") or {}
                conversation_id = other_meta.get("conversation_id")
                if other_node["type"] == "Conversation":
                    conversation_id = other_node["id"]
                if conversation_id:
                    metrics["conversation_ids"].add(str(conversation_id))

        type_max_raw: Dict[str, float] = {}
        for node in nodes:
            degree = degree_map.get(node["id"], 0)
            recency = _recency_score(node.get("updated_at"), now=now)
            metrics = {
                "degree": degree,
                "recency_score": round(recency, 4),
            }
            if node["type"] == "Topic":
                topic_stat = topic_metrics.get(node["id"], {})
                mention_count = float(topic_stat.get("mention_count") or 0.0)
                conversation_count = len(topic_stat.get("conversation_ids") or ())
                raw_importance = (
                    math.log1p(mention_count) * 2.8
                    + math.log1p(conversation_count) * 2.2
                    + recency * 1.4
                    + math.sqrt(max(0, degree)) * 0.45
                )
                metrics.update({
                    "mention_count": round(mention_count, 2),
                    "conversation_count": conversation_count,
                })
            else:
                raw_importance = math.log1p(max(0, degree)) * 1.4 + recency * 0.9

            metrics["importance_raw"] = round(raw_importance, 4)
            node["importance"] = round(raw_importance, 4)
            node["_raw_importance"] = raw_importance
            node["metadata"] = {**(node.get("metadata") or {}), "graph_metrics": metrics}
            type_max_raw[node["type"]] = max(type_max_raw.get(node["type"], 0.0), raw_importance)

        for node in nodes:
            max_raw = max(type_max_raw.get(node["type"], 0.0), 0.0001)
            importance_norm = min(1.0, (node.get("_raw_importance") or 0.0) / max_raw)
            node["importance_norm"] = round(importance_norm, 4)
            node["metadata"]["graph_metrics"]["importance_norm"] = node["importance_norm"]
            node.pop("_raw_importance", None)
        return {"nodes": nodes, "edges": edges}

    def search(self, query: str, limit: int = 30) -> Dict[str, Any]:
        query = str(query or "").strip()
        q = f"%{query}%"
        limit = max(1, min(int(limit or 30), 100))
        nt, et = self._read_tables()
        with self._connect() as conn:
            rows = []
            if query:
                rows = conn.execute(
                    f"""
                    SELECT id, type, title, summary, metadata_json, updated_at
                    FROM {nt}
                    WHERE title LIKE ? OR summary LIKE ? OR metadata_json LIKE ?
                    ORDER BY updated_at DESC, id ASC
                    LIMIT ?
                    """,
                    (q, q, q, limit),
                ).fetchall()

            if len(rows) < limit:
                terms = _topic_candidates(query, limit=8)
                if terms:
                    clauses = []
                    params: List[str] = []
                    for term in terms:
                        clauses.append("(title LIKE ? OR summary LIKE ? OR metadata_json LIKE ?)")
                        params.extend([f"%{term}%", f"%{term}%", f"%{term}%"])
                    extra = conn.execute(
                        f"""
                        SELECT id, type, title, summary, metadata_json, updated_at
                        FROM {nt}
                        WHERE {' OR '.join(clauses)}
                        ORDER BY updated_at DESC, id ASC
                        LIMIT ?
                        """,
                        (*params, limit * 3),
                    ).fetchall()
                    by_id = {row["id"]: row for row in rows}
                    for row in extra:
                        by_id.setdefault(row["id"], row)
                    rows = list(by_id.values())

            terms_for_score = set(_topic_candidates(query, limit=12))
            def score(row: sqlite3.Row) -> tuple:
                haystack = f"{row['title']} {row['summary']} {row['metadata_json']}".lower()
                hits = sum(1 for term in terms_for_score if term.lower() in haystack)
                type_boost = 1 if row["type"] in {
                    "Decision", "Task", "File", "Document", "CodeFile",
                    "Spreadsheet", "SlideDeck", "Image", "ImageText", "Page", "Slide",
                } else 0
                return (hits, type_boost, row["updated_at"] or "")

            rows = sorted(rows, key=score, reverse=True)[:limit]
        return {
            "query": query,
            "matches": [
                {
                    "id": row["id"],
                    "type": row["type"],
                    "title": row["title"],
                    "summary": row["summary"],
                    "metadata": _safe_loads(row["metadata_json"]),
                    "updated_at": row["updated_at"],
                }
                for row in rows
            ],
        }

    def context_for_query(self, query: str, limit: int = 6) -> str:
        """Return compact graph-backed RAG context for chat generation."""
        query = str(query or "").strip()
        if not query:
            return ""
        matches = self.search(query, limit).get("matches", [])
        if not matches:
            topics = _topic_candidates(query, limit=4)
            if topics:
                nt, et = self._read_tables()
                with self._connect() as conn:
                    rows = []
                    for topic in topics:
                        rows.extend(conn.execute(
                            f"""
                            SELECT id, type, title, summary, metadata_json
                            FROM {nt}
                            WHERE title LIKE ? OR metadata_json LIKE ?
                            ORDER BY updated_at DESC, id ASC
                            LIMIT 3
                            """,
                            (f"%{topic}%", f"%{topic}%"),
                        ).fetchall())
                seen = set()
                matches = []
                for row in rows:
                    if row["id"] in seen:
                        continue
                    seen.add(row["id"])
                    matches.append({
                        "id": row["id"],
                        "type": row["type"],
                        "title": row["title"],
                        "summary": row["summary"],
                        "metadata": _safe_loads(row["metadata_json"]),
                    })
                    if len(matches) >= limit:
                        break
        lines = []
        for match in matches[:limit]:
            meta = match.get("metadata") or {}
            source = (
                meta.get("relative_path")
                or meta.get("filename")
                or meta.get("conversation_id")
                or meta.get("source")
                or match["id"]
            )
            summary = _clean_text(match.get("summary") or "")[:700]
            lines.append(f"- [{match['type']}] {match['title']} | source={source} | {summary}")
        return "\n".join(lines)

    def neighbors(self, node_id: str) -> Dict[str, Any]:
        """Return direct neighbors (1-hop) of a node."""
        nt, et = self._read_tables()
        with self._connect() as conn:
            edge_rows = conn.execute(
                f"SELECT from_node, to_node, type, weight FROM {et} WHERE from_node=? OR to_node=? ORDER BY id ASC",
                (node_id, node_id),
            ).fetchall()
            neighbor_ids: set = set()
            edges = []
            for row in edge_rows:
                neighbor_ids.add(row["from_node"])
                neighbor_ids.add(row["to_node"])
                edges.append({"from": row["from_node"], "to": row["to_node"], "type": row["type"], "weight": row["weight"]})
            neighbor_ids.discard(node_id)
            nodes = []
            if neighbor_ids:
                placeholders = ",".join("?" * len(neighbor_ids))
                nodes = [
                    {
                        "id": row["id"],
                        "type": row["type"],
                        "title": row["title"],
                        "summary": row["summary"],
                        "metadata": _safe_loads(row["metadata_json"]),
                    }
                    for row in conn.execute(
                        f"SELECT id, type, title, summary, metadata_json FROM {nt} WHERE id IN ({placeholders}) ORDER BY id ASC",
                        list(neighbor_ids),
                    )
                ]
        return {"node_id": node_id, "neighbors": nodes, "edges": edges}

    def get_node(self, node_id: str) -> Dict[str, Any]:
        node_id = str(node_id or "").strip()
        if not node_id:
            raise ValueError("node_id required")
        nt, et = self._read_tables()
        with self._connect() as conn:
            row = conn.execute(
                f"""
                SELECT id, type, title, summary, metadata_json, updated_at
                FROM {nt}
                WHERE id=?
                """,
                (node_id,),
            ).fetchone()
            if not row:
                raise ValueError(f"graph node not found: {node_id}")
            degree = conn.execute(
                f"SELECT COUNT(*) AS c FROM {et} WHERE from_node=? OR to_node=?",
                (node_id, node_id),
            ).fetchone()["c"]
        return {
            "id": row["id"],
            "type": row["type"],
            "title": row["title"],
            "summary": row["summary"],
            "metadata": _safe_loads(row["metadata_json"]),
            "updated_at": row["updated_at"],
            "degree": degree,
        }

    def relationship_search(
        self,
        *,
        query: str = "",
        node_id: str = "",
        relationship_type: str = "",
        limit: int = 30,
    ) -> Dict[str, Any]:
        query = str(query or "").strip()
        node_id = str(node_id or "").strip()
        relationship_type = str(relationship_type or "").strip()
        limit = max(1, min(int(limit or 30), 200))
        nt, et = self._read_tables()
        where = []
        params: List[Any] = []
        if node_id:
            where.append("(e.from_node=? OR e.to_node=?)")
            params.extend([node_id, node_id])
        if relationship_type:
            where.append("e.type LIKE ?")
            params.append(f"%{relationship_type}%")
        if query:
            where.append(
                "(e.type LIKE ? OR e.metadata_json LIKE ? OR src.title LIKE ? OR dst.title LIKE ? OR src.summary LIKE ? OR dst.summary LIKE ?)"
            )
            params.extend([f"%{query}%"] * 6)
        where_sql = "WHERE " + " AND ".join(where) if where else ""
        with self._connect() as conn:
            rows = conn.execute(
                f"""
                SELECT
                  e.id, e.from_node, e.to_node, e.type, e.weight, e.metadata_json, e.created_at,
                  src.type AS source_type, src.title AS source_title, src.summary AS source_summary,
                  src.metadata_json AS source_metadata,
                  dst.type AS target_type, dst.title AS target_title, dst.summary AS target_summary,
                  dst.metadata_json AS target_metadata
                FROM {et} e
                JOIN {nt} src ON src.id=e.from_node
                JOIN {nt} dst ON dst.id=e.to_node
                {where_sql}
                ORDER BY e.weight DESC, e.created_at DESC, e.id ASC
                LIMIT ?
                """,
                (*params, limit),
            ).fetchall()
        return {
            "query": query,
            "node_id": node_id,
            "relationship_type": relationship_type,
            "relationships": [
                {
                    "id": row["id"],
                    "type": row["type"],
                    "weight": row["weight"],
                    "metadata": _safe_loads(row["metadata_json"]),
                    "created_at": row["created_at"],
                    "source": {
                        "id": row["from_node"],
                        "type": row["source_type"],
                        "title": row["source_title"],
                        "summary": row["source_summary"],
                        "metadata": _safe_loads(row["source_metadata"]),
                    },
                    "target": {
                        "id": row["to_node"],
                        "type": row["target_type"],
                        "title": row["target_title"],
                        "summary": row["target_summary"],
                        "metadata": _safe_loads(row["target_metadata"]),
                    },
                }
                for row in rows
            ],
        }

    def traverse(self, node_id: str, *, depth: int = 1, limit: int = 100) -> Dict[str, Any]:
        node_id = str(node_id or "").strip()
        if not node_id:
            raise ValueError("node_id required")
        depth = max(0, min(int(depth or 1), 4))
        limit = max(1, min(int(limit or 100), 500))
        nt, et = self._read_tables()
        visited = {node_id}
        frontier = {node_id}
        edges_by_id: Dict[str, Dict[str, Any]] = {}
        with self._connect() as conn:
            for _ in range(depth):
                if not frontier or len(visited) >= limit:
                    break
                placeholders = ",".join("?" * len(frontier))
                rows = conn.execute(
                    f"""
                    SELECT id, from_node, to_node, type, weight, metadata_json
                    FROM {et}
                    WHERE from_node IN ({placeholders}) OR to_node IN ({placeholders})
                    ORDER BY weight DESC, id ASC
                    LIMIT ?
                    """,
                    (*frontier, *frontier, limit * 3),
                ).fetchall()
                next_frontier = set()
                for row in rows:
                    edges_by_id[row["id"]] = {
                        "id": row["id"],
                        "from": row["from_node"],
                        "to": row["to_node"],
                        "type": row["type"],
                        "weight": row["weight"],
                        "metadata": _safe_loads(row["metadata_json"]),
                    }
                    for candidate in (row["from_node"], row["to_node"]):
                        if candidate not in visited and len(visited) < limit:
                            visited.add(candidate)
                            next_frontier.add(candidate)
                frontier = next_frontier
            placeholders = ",".join("?" * len(visited))
            node_rows = conn.execute(
                f"""
                SELECT id, type, title, summary, metadata_json, updated_at
                FROM {nt}
                WHERE id IN ({placeholders})
                ORDER BY updated_at DESC, id ASC
                """,
                list(visited),
            ).fetchall()
        return {
            "root": node_id,
            "depth": depth,
            "nodes": [
                {
                    "id": row["id"],
                    "type": row["type"],
                    "title": row["title"],
                    "summary": row["summary"],
                    "metadata": _safe_loads(row["metadata_json"]),
                    "updated_at": row["updated_at"],
                }
                for row in node_rows
            ],
            "edges": list(edges_by_id.values()),
        }

    def _iter_vector_source_items(
        self,
        conn: sqlite3.Connection,
        *,
        include_nodes: bool = True,
        include_chunks: bool = True,
    ) -> List[Dict[str, Any]]:
        items: List[Dict[str, Any]] = []
        if include_nodes:
            for row in conn.execute(
                """
                SELECT id, type, title, summary, metadata_json
                FROM nodes
                WHERE type <> 'Chunk'
                ORDER BY updated_at DESC, id ASC
                """
            ).fetchall():
                metadata = _safe_loads(row["metadata_json"])
                text = self._vector_text_for_node(
                    title=row["title"],
                    summary=row["summary"] or "",
                    metadata=metadata,
                )
                if text:
                    items.append({
                        "item_id": row["id"],
                        "item_type": "node",
                        "source_node": row["id"],
                        "text": text,
                        "metadata": {"node_type": row["type"], **metadata},
                    })
        if include_chunks:
            for row in conn.execute(
                """
                SELECT c.id, c.source_node AS parent_source_node, c.text, c.metadata_json
                FROM chunks c
                JOIN nodes n ON n.id=c.id
                ORDER BY c.created_at DESC, c.id ASC
                """
            ).fetchall():
                metadata = _safe_loads(row["metadata_json"])
                text = _clean_text(row["text"] or "")
                if text:
                    items.append({
                        "item_id": row["id"],
                        "item_type": "chunk",
                        "source_node": row["id"],
                        "text": text,
                        "metadata": {**metadata, "parent_source_node": row["parent_source_node"]},
                    })
        return items

    def rebuild_vector_index(
        self,
        *,
        full: bool = False,
        include_nodes: bool = True,
        include_chunks: bool = True,
    ) -> Dict[str, Any]:
        """Rebuild the derived vector index without mutating graph content."""
        op_id = f"vector-op:{_sha256_text(f'{time.time()}:{os.getpid()}')[:24]}"
        requested_at = _now()
        started = time.perf_counter()
        try:
            with self._connect() as conn:
                conn.execute(
                    """
                    INSERT INTO vector_index_operations(
                      id, operation, status, requested_at, started_at, metadata_json
                    )
                    VALUES (?, ?, 'running', ?, ?, ?)
                    """,
                    (
                        op_id,
                        "rebuild_full" if full else "rebuild_incremental",
                        requested_at,
                        requested_at,
                        _json({"include_nodes": include_nodes, "include_chunks": include_chunks}),
                    ),
                )
                if full:
                    filters = []
                    if include_nodes:
                        filters.append("'node'")
                    if include_chunks:
                        filters.append("'chunk'")
                    if filters:
                        conn.execute(f"DELETE FROM vector_embeddings WHERE item_type IN ({','.join(filters)})")
                items = self._iter_vector_source_items(
                    conn,
                    include_nodes=include_nodes,
                    include_chunks=include_chunks,
                )
                indexed = skipped = 0
                for item in items:
                    changed = self._upsert_vector_item(conn, **item)
                    if changed:
                        indexed += 1
                    else:
                        skipped += 1
                duration_ms = round((time.perf_counter() - started) * 1000, 2)
                conn.execute(
                    """
                    UPDATE vector_index_operations
                    SET status='completed', completed_at=?, items_total=?,
                        items_indexed=?, items_skipped=?, metadata_json=?
                    WHERE id=?
                    """,
                    (
                        _now(),
                        len(items),
                        indexed,
                        skipped,
                        _json({
                            "include_nodes": include_nodes,
                            "include_chunks": include_chunks,
                            "duration_ms": duration_ms,
                            "embedding_model": self._embedding_model.model_id,
                            "embedding_dim": self._embedding_model.dim,
                        }),
                        op_id,
                    ),
                )
            return {
                "status": "completed",
                "operation_id": op_id,
                "full": bool(full),
                "items_total": len(items),
                "items_indexed": indexed,
                "items_skipped": skipped,
                "duration_ms": duration_ms,
                "embedding_model": self._embedding_model.model_id,
                "embedding_dim": self._embedding_model.dim,
            }
        except Exception as exc:
            duration_ms = round((time.perf_counter() - started) * 1000, 2)
            with self._connect() as conn:
                conn.execute(
                    """
                    INSERT INTO vector_index_operations(
                      id, operation, status, requested_at, started_at, completed_at,
                      error_message, metadata_json
                    )
                    VALUES (?, ?, 'failed', ?, ?, ?, ?, ?)
                    ON CONFLICT(id) DO UPDATE SET
                      status='failed',
                      completed_at=excluded.completed_at,
                      error_message=excluded.error_message,
                      metadata_json=excluded.metadata_json
                    """,
                    (
                        op_id,
                        "rebuild_full" if full else "rebuild_incremental",
                        requested_at,
                        requested_at,
                        _now(),
                        str(exc),
                        _json({"duration_ms": duration_ms}),
                    ),
                )
            raise

    def index_status(self) -> Dict[str, Any]:
        with self._connect() as conn:
            vector_counts = {
                row["item_type"]: row["count"]
                for row in conn.execute(
                    "SELECT item_type, COUNT(*) AS count FROM vector_embeddings GROUP BY item_type"
                )
            }
            source_items = self._iter_vector_source_items(conn)
            vector_rows = {
                row["item_id"]: row
                for row in conn.execute(
                    """
                    SELECT item_id, text_hash, embedding_dim, embedding_model, indexed_at
                    FROM vector_embeddings
                    """
                ).fetchall()
            }
            latest_rows = conn.execute(
                """
                SELECT id, operation, status, requested_at, started_at, completed_at,
                       items_total, items_indexed, items_skipped, error_message, metadata_json
                FROM vector_index_operations
                ORDER BY requested_at DESC, id DESC
                LIMIT 5
                """
            ).fetchall()
        missing = stale = ready = 0
        for item in source_items:
            vector_row = vector_rows.get(item["item_id"])
            expected_hash = _sha256_text(_clean_text(item["text"]))
            if not vector_row:
                missing += 1
            elif (
                vector_row["text_hash"] != expected_hash
                or vector_row["embedding_dim"] != self._embedding_model.dim
                or vector_row["embedding_model"] != self._embedding_model.model_id
            ):
                stale += 1
            else:
                ready += 1
        pending = missing + stale
        return {
            "status": "ready" if pending == 0 else "needs_reindex",
            "storage": {
                "db_path": str(self.db_path),
                "backend": "sqlite",
                "embedding_model": self._embedding_model.model_id,
                "embedding_dim": self._embedding_model.dim,
            },
            "source_items": len(source_items),
            "indexed_items": sum(vector_counts.values()),
            "ready_items": ready,
            "missing_items": missing,
            "stale_items": stale,
            "pending_items": pending,
            "by_item_type": vector_counts,
            "operations": [
                {
                    "id": row["id"],
                    "operation": row["operation"],
                    "status": row["status"],
                    "requested_at": row["requested_at"],
                    "started_at": row["started_at"],
                    "completed_at": row["completed_at"],
                    "items_total": row["items_total"],
                    "items_indexed": row["items_indexed"],
                    "items_skipped": row["items_skipped"],
                    "error_message": row["error_message"],
                    "metadata": _safe_loads(row["metadata_json"]),
                }
                for row in latest_rows
            ],
        }

    def vector_search(
        self,
        query: str,
        *,
        limit: int = 30,
        min_score: float = 0.0,
        max_candidates: int = 10_000,
    ) -> Dict[str, Any]:
        query = str(query or "").strip()
        limit = max(1, min(int(limit or 30), 100))
        min_score = float(min_score or 0.0)
        if not query:
            return {"query": query, "matches": []}
        query_vector = self._embedding_model.embed(query)
        max_candidates = max(limit, min(int(max_candidates or 10_000), 50_000))
        with self._connect() as conn:
            rows = conn.execute(
                """
                SELECT
                  ve.item_id, ve.item_type, ve.source_node, ve.embedding,
                  ve.embedding_dim, ve.embedding_model, ve.metadata_json AS vector_metadata,
                  n.type AS node_type, n.title AS node_title, n.summary AS node_summary,
                  n.metadata_json AS node_metadata, n.updated_at AS node_updated_at,
                  c.text AS chunk_text, c.source_node AS parent_node_id,
                  pn.type AS parent_type, pn.title AS parent_title,
                  pn.summary AS parent_summary, pn.metadata_json AS parent_metadata,
                  pn.updated_at AS parent_updated_at
                FROM vector_embeddings ve
                LEFT JOIN nodes n ON n.id=ve.source_node
                LEFT JOIN chunks c ON c.id=ve.item_id
                LEFT JOIN nodes pn ON pn.id=c.source_node
                WHERE ve.embedding_model=? AND ve.embedding_dim=?
                ORDER BY ve.indexed_at DESC
                LIMIT ?
                """,
                (self._embedding_model.model_id, self._embedding_model.dim, max_candidates),
            ).fetchall()
        scored = []
        for row in rows:
            vector = self._embedding_model.decode(row["embedding"], row["embedding_dim"])
            score = self._embedding_model.similarity(query_vector, vector)
            if score < min_score:
                continue
            is_chunk = row["item_type"] == "chunk"
            summary = row["chunk_text"] if is_chunk and row["chunk_text"] else row["node_summary"]
            parent_metadata = _safe_loads(row["parent_metadata"])
            node_metadata = _safe_loads(row["node_metadata"])
            scored.append({
                "id": row["item_id"],
                "node_id": row["parent_node_id"] if is_chunk and row["parent_node_id"] else row["source_node"],
                "item_type": row["item_type"],
                "type": "Chunk" if is_chunk else row["node_type"],
                "title": row["parent_title"] if is_chunk and row["parent_title"] else row["node_title"],
                "summary": _clean_text(summary or "")[:1000],
                "score": round(float(score), 6),
                "metadata": {
                    **(parent_metadata if is_chunk else node_metadata),
                    "vector": _safe_loads(row["vector_metadata"]),
                    "parent_node_id": row["parent_node_id"],
                    "parent_type": row["parent_type"],
                },
                "updated_at": row["parent_updated_at"] if is_chunk and row["parent_updated_at"] else row["node_updated_at"],
            })
        scored.sort(key=lambda item: (item["score"], item.get("updated_at") or ""), reverse=True)
        return {
            "query": query,
            "embedding_model": self._embedding_model.model_id,
            "embedding_dim": self._embedding_model.dim,
            "matches": scored[:limit],
        }

    def delete_conversation(self, conversation_id: str) -> Dict[str, Any]:
        conversation_id = str(conversation_id or "").strip()
        if not conversation_id:
            return {"status": "skipped", "removed_nodes": 0}
        conv_id = f"conversation:{_slug(conversation_id)}"
        with self._connect() as conn:
            direct_ids = [
                row["to_node"]
                for row in conn.execute(
                    "SELECT to_node FROM edges WHERE from_node=? AND type='contains'",
                    (conv_id,),
                )
            ]
            remove_ids = set(direct_ids)
            for source_id in list(direct_ids):
                for row in conn.execute(
                    """
                    SELECT to_node FROM edges
                    WHERE from_node=? AND type IN ('has_chunk', 'implies', 'contains_signal', 'has_page', 'has_slide', 'has_sheet', 'contains_image')
                    """,
                    (source_id,),
                ):
                    remove_ids.add(row["to_node"])
            remove_ids.add(conv_id)
            for node_id in remove_ids:
                conn.execute("DELETE FROM nodes WHERE id=?", (node_id,))
                if KGStoreV2 is not None:
                    conn.execute("DELETE FROM nodes_v2 WHERE id=?", (node_id,))  # edges_v2 cascade
            conn.execute(
                """
                DELETE FROM nodes
                WHERE type='Topic'
                  AND id NOT IN (SELECT to_node FROM edges)
                  AND id NOT IN (SELECT from_node FROM edges)
                """
            )
            if KGStoreV2 is not None:
                conn.execute(
                    """
                    DELETE FROM nodes_v2
                    WHERE legacy_type='Topic'
                      AND id NOT IN (SELECT target FROM edges_v2)
                      AND id NOT IN (SELECT source FROM edges_v2)
                    """
                )
        return {"status": "ok", "conversation_id": conversation_id, "removed_nodes": len(remove_ids)}

    def clear_all(self) -> Dict[str, Any]:
        with self._connect() as conn:
            counts = {
                "nodes": conn.execute("SELECT COUNT(*) AS c FROM nodes").fetchone()["c"],
                "edges": conn.execute("SELECT COUNT(*) AS c FROM edges").fetchone()["c"],
                "chunks": conn.execute("SELECT COUNT(*) AS c FROM chunks").fetchone()["c"],
                "knowledge_sources": conn.execute("SELECT COUNT(*) AS c FROM knowledge_sources").fetchone()["c"],
                "local_file_index": conn.execute("SELECT COUNT(*) AS c FROM local_file_index").fetchone()["c"],
            }
            conn.execute("DELETE FROM local_file_index")
            conn.execute("DELETE FROM knowledge_sources")
            conn.execute("DELETE FROM chunks")
            conn.execute("DELETE FROM edges")
            conn.execute("DELETE FROM nodes")
            if KGStoreV2 is not None:
                conn.execute("DELETE FROM edges_v2")
                conn.execute("DELETE FROM nodes_v2")
        if self.blob_dir.exists():
            shutil.rmtree(self.blob_dir, ignore_errors=True)
            self.blob_dir.mkdir(parents=True, exist_ok=True)
        return {"status": "ok", "removed": counts}

    def stats(self) -> Dict[str, Any]:
        nt, et = self._read_tables()
        with self._connect() as conn:
            node_counts = {
                row["type"]: row["count"]
                for row in conn.execute(f"SELECT type, COUNT(*) AS count FROM {nt} GROUP BY type")
            }
            edge_counts = {
                row["type"]: row["count"]
                for row in conn.execute(f"SELECT type, COUNT(*) AS count FROM {et} GROUP BY type")
            }
            local_sources = conn.execute("SELECT COUNT(*) AS c FROM knowledge_sources").fetchone()["c"]
            local_file_status = {
                row["status"]: row["count"]
                for row in conn.execute("SELECT status, COUNT(*) AS count FROM local_file_index GROUP BY status")
            }
        v2 = None
        if KGStoreV2 is not None:
            try:
                v2 = KGStoreV2(self.db_path).stats()
            except Exception as e:
                v2 = {"available": False, "error": str(e)}
        return {
            "db_path": str(self.db_path),
            "schema_version": GRAPH_SCHEMA_VERSION,
            "v2_schema_available": KGStoreV2 is not None,
            "nodes": node_counts,
            "edges": edge_counts,
            "local_sources": local_sources,
            "local_file_status": local_file_status,
            "v2": v2,
        }

    def search_for_document_generation(self, query: str, limit: int = 10) -> List[Dict[str, Any]]:
        """Hybrid retrieval optimized for document generation.

        Scoring: 0.5*text_relevance + 0.3*graph_relationship + 0.2*recency
        Returns nodes with rich context for document generation prompts.
        """
        query = str(query or "").strip()
        if not query:
            return []
        limit = max(1, min(int(limit or 10), 50))
        terms = _topic_candidates(query, limit=12)
        now = datetime.now()
        nt, et = self._read_tables()

        with self._connect() as conn:
            candidate_rows = []
            seen_ids = set()

            if query:
                q = f"%{query}%"
                rows = conn.execute(
                    f"""
                    SELECT id, type, title, summary, metadata_json, updated_at
                    FROM {nt}
                    WHERE (title LIKE ? OR summary LIKE ? OR metadata_json LIKE ?)
                      AND type IN ('Document', 'File', 'CodeFile', 'SlideDeck',
                                   'Spreadsheet', 'Image', 'ImageText', 'Chat',
                                   'Decision', 'Task', 'Concept', 'Feature',
                                   'Page', 'Slide')
                    ORDER BY updated_at DESC, id ASC
                    LIMIT ?
                    """,
                    (q, q, q, limit * 5),
                ).fetchall()
                for row in rows:
                    if row["id"] not in seen_ids:
                        seen_ids.add(row["id"])
                        candidate_rows.append(row)

            for term in terms:
                t = f"%{term}%"
                rows = conn.execute(
                    f"""
                    SELECT id, type, title, summary, metadata_json, updated_at
                    FROM {nt}
                    WHERE (title LIKE ? OR summary LIKE ? OR metadata_json LIKE ?)
                      AND type IN ('Document', 'File', 'CodeFile', 'SlideDeck',
                                   'Spreadsheet', 'Image', 'ImageText', 'Chat',
                                   'Decision', 'Task', 'Concept', 'Feature',
                                   'Page', 'Slide')
                    ORDER BY updated_at DESC, id ASC
                    LIMIT ?
                    """,
                    (t, t, t, limit * 3),
                ).fetchall()
                for row in rows:
                    if row["id"] not in seen_ids:
                        seen_ids.add(row["id"])
                        candidate_rows.append(row)

            scored_results = []
            for row in candidate_rows:
                haystack = f"{row['title']} {row['summary']} {row['metadata_json']}".lower()

                text_hits = sum(1 for term in terms if term.lower() in haystack)
                text_score = min(1.0, text_hits / max(len(terms), 1))

                edge_count = conn.execute(
                    f"SELECT COUNT(*) AS c FROM {et} WHERE from_node=? OR to_node=?",
                    (row["id"], row["id"]),
                ).fetchone()["c"]
                graph_score = min(1.0, math.log1p(edge_count) / 4.0)

                recency = _recency_score(row["updated_at"], now=now, half_life_days=14.0)

                doc_type_boost = 1.2 if row["type"] in (
                    "Document", "File", "SlideDeck", "Decision",
                ) else 1.0

                hybrid_score = (
                    0.5 * text_score
                    + 0.3 * graph_score
                    + 0.2 * recency
                ) * doc_type_boost

                meta = _safe_loads(row["metadata_json"])
                neighbor_concepts = []
                neighbor_rows = conn.execute(
                    f"""
                    SELECT n.title, n.type FROM {et} e
                    JOIN {nt} n ON n.id = CASE WHEN e.from_node = ? THEN e.to_node ELSE e.from_node END
                    WHERE (e.from_node = ? OR e.to_node = ?)
                      AND n.type IN ('Concept', 'Feature', 'Decision', 'Task')
                    LIMIT 8
                    """,
                    (row["id"], row["id"], row["id"]),
                ).fetchall()
                for nr in neighbor_rows:
                    neighbor_concepts.append({"title": nr["title"], "type": nr["type"]})

                scored_results.append({
                    "id": row["id"],
                    "type": row["type"],
                    "title": row["title"],
                    "summary": row["summary"],
                    "metadata": meta,
                    "updated_at": row["updated_at"],
                    "hybrid_score": round(hybrid_score, 4),
                    "scores": {
                        "text": round(text_score, 4),
                        "graph": round(graph_score, 4),
                        "recency": round(recency, 4),
                    },
                    "related_concepts": neighbor_concepts,
                })

            scored_results.sort(key=lambda x: x["hybrid_score"], reverse=True)
            return scored_results[:limit]

    def multi_hop_context(self, node_ids: List[str], max_hops: int = 2) -> Dict[str, Any]:
        """Multi-hop graph traversal from seed nodes for richer context."""
        visited_nodes = set()
        visited_edges = set()
        all_nodes = []
        all_edges = []
        frontier = set(node_ids)
        nt, et = self._read_tables()

        with self._connect() as conn:
            for hop in range(max_hops):
                if not frontier:
                    break
                next_frontier = set()
                for nid in frontier:
                    if nid in visited_nodes:
                        continue
                    visited_nodes.add(nid)
                    row = conn.execute(
                        f"SELECT id, type, title, summary, metadata_json, updated_at FROM {nt} WHERE id=?",
                        (nid,),
                    ).fetchone()
                    if row:
                        all_nodes.append({
                            "id": row["id"], "type": row["type"],
                            "title": row["title"], "summary": row["summary"],
                            "metadata": _safe_loads(row["metadata_json"]),
                            "hop": hop,
                        })
                    edge_rows = conn.execute(
                        f"""
                        SELECT id, from_node, to_node, type, weight
                        FROM {et} WHERE from_node=? OR to_node=?
                        ORDER BY id ASC
                        """,
                        (nid, nid),
                    ).fetchall()
                    for er in edge_rows:
                        if er["id"] not in visited_edges:
                            visited_edges.add(er["id"])
                            all_edges.append({
                                "from": er["from_node"], "to": er["to_node"],
                                "type": er["type"], "weight": er["weight"],
                            })
                            other = er["to_node"] if er["from_node"] == nid else er["from_node"]
                            if other not in visited_nodes:
                                next_frontier.add(other)
                frontier = next_frontier

        return {"nodes": all_nodes, "edges": all_edges}
