"""Document generation (docx/xlsx/pptx/pdf) and extraction (read_document)."""

from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, Optional, Tuple

from latticeai.tools import (
    _SUPPORTED_READ_EXTENSIONS,
    DOCUMENT_MAX_READ_BYTES,
    DOCUMENT_OUTPUT_DIR,
    PDF_OUTPUT_DIR,
    PRESENTATION_OUTPUT_DIR,
    SPREADSHEET_OUTPUT_DIR,
    ToolError,
)


def _safe_filename(name: str, suffix: str) -> str:
    base = Path(name or f"artifact{suffix}").name
    if not base.lower().endswith(suffix):
        base += suffix
    safe = "".join(ch if ch.isalnum() or ch in ("-", "_", ".", " ") else "_" for ch in base).strip()
    return safe or f"artifact{suffix}"


#: tool name → (output directory, enforced suffix). The creators below are the
#: only writers, so this table is the single source of truth for "where does
#: this document actually land?".
_DOCUMENT_TOOL_TARGETS: Dict[str, Tuple[str, str]] = {
    "create_docx": (DOCUMENT_OUTPUT_DIR, ".docx"),
    "create_xlsx": (SPREADSHEET_OUTPUT_DIR, ".xlsx"),
    "create_pptx": (PRESENTATION_OUTPUT_DIR, ".pptx"),
    "create_pdf": (PDF_OUTPUT_DIR, ".pdf"),
}


def document_output_target(tool_name: str, filename: str) -> Optional[str]:
    """Workspace-relative path ``tool_name`` will write ``filename`` to.

    Governance needs this: the creators sanitize the caller's ``filename``
    into their own output directory, so an "does the target already exist?"
    check against the raw argument looks at a path nothing ever writes — and
    the fail-closed overwrite guard would silently never fire.

    Returns ``None`` for tools that write wherever the caller points them, so
    callers fall back to the raw argument.
    """
    entry = _DOCUMENT_TOOL_TARGETS.get(tool_name)
    if entry is None:
        return None
    output_dir, suffix = entry
    return f"{output_dir}/{_safe_filename(filename, suffix)}"


def _body_to_str(body) -> str:
    if isinstance(body, list):
        return "\n\n".join(str(item) for item in body)
    return str(body or "")


def read_document(path: str) -> Dict[str, Any]:
    """Extract text from PDF, DOCX, XLSX, PPTX, TXT, MD, CSV files."""
    target = Path(path).expanduser().resolve()
    if not target.exists():
        raise ToolError(f"파일이 없습니다: {path}")
    if not target.is_file():
        raise ToolError(f"파일이 아닙니다: {path}")
    if target.stat().st_size > DOCUMENT_MAX_READ_BYTES:
        raise ToolError(f"파일이 너무 큽니다 ({target.stat().st_size:,} bytes).")

    ext = target.suffix.lower()
    if ext not in _SUPPORTED_READ_EXTENSIONS:
        raise ToolError(f"지원하지 않는 형식입니다: {ext}. 지원: {', '.join(_SUPPORTED_READ_EXTENSIONS)}")

    text = ""
    meta: Dict[str, Any] = {"path": str(target), "ext": ext}

    if ext == ".pdf":
        try:
            import pdfplumber
            with pdfplumber.open(str(target)) as pdf:
                meta["pages"] = len(pdf.pages)
                text = "\n\n".join(
                    (p.extract_text() or "") for p in pdf.pages
                ).strip()
        except Exception as exc:
            raise ToolError(f"PDF 읽기 실패: {exc}") from exc

    elif ext == ".docx":
        try:
            from docx import Document
            doc = Document(str(target))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            text = "\n\n".join(paragraphs)
            meta["paragraphs"] = len(paragraphs)
        except Exception as exc:
            raise ToolError(f"DOCX 읽기 실패: {exc}") from exc

    elif ext == ".xlsx":
        try:
            from openpyxl import load_workbook
            wb = load_workbook(str(target), data_only=True)
            rows_all = []
            for ws in wb.worksheets:
                rows_all.append(f"[Sheet: {ws.title}]")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    rows_all.append("\t".join(cells))
            text = "\n".join(rows_all)
            meta["sheets"] = len(wb.worksheets)
        except Exception as exc:
            raise ToolError(f"XLSX 읽기 실패: {exc}") from exc

    elif ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(target))
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                parts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
                slides_text.append(f"[Slide {i}]\n" + "\n".join(parts))
            text = "\n\n".join(slides_text)
            meta["slides"] = len(prs.slides)
        except Exception as exc:
            raise ToolError(f"PPTX 읽기 실패: {exc}") from exc

    elif ext in {".txt", ".md", ".csv"}:
        try:
            text = target.read_text(encoding="utf-8", errors="replace")
        except Exception as exc:
            raise ToolError(f"파일 읽기 실패: {exc}") from exc

    meta["chars"] = len(text)
    meta["preview"] = text[:500]
    meta["content"] = text[:50_000]   # 50K char cap for context
    return meta
