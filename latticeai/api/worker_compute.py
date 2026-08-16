"""The v11.6.0 pure-compute seams: data in, data out, nothing stored.

Plan §설계 결정 2 (revised) hands **every write** to Rust — platform state and
the knowledge graph both. What is left for Python is the part Rust cannot do
without shipping a model runtime and half of PyPI: turning text into vectors,
turning a document into text, turning a spec into document bytes, and turning
audio into a transcript.

A ninth seam, ``POST /worker/multimodal/describe``, was here until v11.8.0. It
wrapped :func:`lattice_brain.multimodal.extract_image_facts` for a native image
ingest that was never built, so nothing in the tree ever called it — not
``lattice-ingest``, not the gateway. The Brain Core functions behind it are
untouched; what went is the door nobody opened.

``worker_seams.py`` holds the *state* seams (``/worker/chat/record-turn``,
``/worker/graph/mutate``) that Wave 2 codes against and Wave 2.5 §W3 retires.
This module is the opposite half and is deliberately a separate file: nothing
here opens a database, writes a file, or reaches a store. Each handler is a
function of its request body, so a caller can retry it, cache it, or run two of
them at once without asking who else is writing.

The eight seams, and what each was extracted from:

``POST /worker/embed``
    ``EmbeddingProvider.embed_batch`` on the *resolved* provider — the same
    object ``phase_brain`` hands the graph store, so a vector produced here is
    the vector Python would have written. ``kind="passage"`` applies the
    50,000-character clamp ``write_master._upsert_vector_item`` applies before
    embedding; ``kind="query"`` does not, matching ``retrieval_vector.search``.

``POST /worker/parse``
    :func:`latticeai.tools.documents.read_document` over the posted bytes. The
    parser matrix (pdfplumber / python-docx / openpyxl / python-pptx / plain
    text) is unchanged; only the file it reads is a temporary one this seam
    writes and deletes.

``POST /worker/render/{docx,xlsx,pptx,pdf}``
    The *building* half of ``tools.documents.create_*`` — same libraries, same
    layout decisions — with ``save(path)`` replaced by ``save(BytesIO)``. Rust
    places the file, sanitises the name and resolves the target; this seam
    never learns where and answers with bytes only.

``POST /worker/asr``
    ``VoiceCaptureService._transcribe``'s contract without the ingest: the same
    injected transcriber port, the same container whitelist and size ceiling,
    and the same refusal to call an empty transcript "text".

``POST /worker/extract``
    :func:`~lattice_brain.graph._kg_common.extraction._extract_concepts` (LLM-first),
    :func:`~lattice_brain.graph._kg_common.extraction._extract_triples` and
    :func:`~lattice_brain.graph._kg_common.extraction._semantic_items` — the
    structures ``ingest_message`` / ``ingest_document`` / ``ingest_source``
    consume, already classified. Rust writes the Concept/Task/Decision subgraph;
    this seam never opens a store.

Gating is the seam gate the rest of the worker uses — ``LATTICEAI_AGENT_TOOL_SEAM``
read per request through :func:`latticeai.api.agent_worker_seam._seam_open`,
then ``require_user`` and the ``agent_seam`` rate bucket. Off ⇒ 404. Mounted
only by :func:`latticeai.runtime.build_phases.worker_profile.phase_worker_routes`,
so ``create_app`` and the committed OpenAPI contract are untouched.
"""

from __future__ import annotations

import asyncio
import base64
import binascii
import contextlib
import importlib.util
import io
import logging
import sys
import tempfile
from pathlib import Path
from typing import Any, Callable, Dict, Iterator, List, Optional, Tuple

from fastapi import APIRouter, Request
from pydantic import BaseModel, Field

from latticeai.api.agent_worker_seam import SEAM_RATE_BUCKET
from latticeai.api.agent_worker_seam import _seam_open as seam_open
from latticeai.core.messages import MESSAGES, http_error, resolve_language
from latticeai.core.quiet import quiet
from latticeai.services.voice_capture import (
    MAX_AUDIO_BYTES,
    SUPPORTED_AUDIO_EXTENSIONS,
)
from latticeai.tools import _CJK_FONT_CANDIDATES, ToolError
from latticeai.tools.documents import _body_to_str, read_document

logger = logging.getLogger(__name__)

#: What an embedding is *for*. The two differ in one thing only, and that thing
#: is real: the write path clamps to 50,000 characters before embedding
#: (``write_master._upsert_vector_item``) while the query path does not
#: (``retrieval_vector.search``). Naming it keeps Rust from having to know.
EMBED_KINDS: Tuple[str, ...] = ("query", "passage")

#: Which ingest door the caller is standing in. ``message`` is
#: ``ingest_message`` (concept limit 12); ``document`` is ``ingest_document``
#: / ``ingest_source`` (limit 15). The text itself is whatever that door
#: composed — the message body, or ``filename\\n{text}`` / ``title\\n{text}``.
EXTRACT_KINDS: Tuple[str, ...] = ("message", "document")
EXTRACT_LIMITS: Dict[str, int] = {"message": 12, "document": 15}

#: The clamp ``write_master`` applies. Copied as a constant rather than
#: imported: it is a literal at that call site, and importing a literal out of a
#: write-side module would tie this compute seam to a module W1 is replacing.
PASSAGE_MAX_CHARS = 50_000

#: The CJK-capable fonts ``create_pdf`` looks for, as a module attribute so a
#: test can point the probe at a font that exists (and at one that does not)
#: instead of asserting whatever the host machine happens to have installed.
CJK_FONT_CANDIDATES: List[str] = list(_CJK_FONT_CANDIDATES)

#: Messages this module owns. Registered into the one shared catalog with
#: ``setdefault`` — the same pattern ``worker_seams.py`` uses — so the entries
#: can be lifted into ``latticeai/core/messages.py`` verbatim later and this
#: registration becomes a no-op rather than an overwrite.
WORKER_COMPUTE_MESSAGES: Dict[str, Dict[str, str]] = {
    "worker_compute.embedder_unavailable": {
        "ko": "임베딩 공급자가 연결되어 있지 않습니다.",
        "en": "No embedding provider is connected to this worker.",
    },
    "worker_compute.kind_invalid": {
        "ko": "'{kind}' 은(는) 임베딩 종류가 아닙니다. {allowed} 중 하나여야 합니다.",
        "en": "'{kind}' is not an embedding kind. Use one of {allowed}.",
    },
    "worker_compute.content_invalid": {
        "ko": "본문을 base64로 읽지 못했습니다: {reason}",
        "en": "The payload is not valid base64: {reason}",
    },
    "worker_compute.parse_failed": {
        "ko": "문서를 읽지 못했습니다: {reason}",
        "en": "The document could not be parsed: {reason}",
    },
    "worker_compute.render_unavailable": {
        "ko": "이 워커는 '{kind}' 문서를 만들 수 없습니다: {reason}",
        "en": "This worker cannot render '{kind}' documents: {reason}",
    },
    "worker_compute.render_failed": {
        "ko": "'{kind}' 문서 생성이 실패했습니다: {reason}",
        "en": "Rendering the '{kind}' document failed: {reason}",
    },
    "worker_compute.audio_unsupported": {
        "ko": "'{suffix}' 은(는) 지원하는 오디오 형식이 아닙니다. {allowed} 만 가능합니다.",
        "en": "'{suffix}' is not a supported audio container. Supported: {allowed}.",
    },
    "worker_compute.audio_too_large": {
        "ko": "오디오가 {size} bytes 입니다. 한도는 {limit} bytes 입니다.",
        "en": "The audio is {size} bytes; the limit is {limit} bytes.",
    },
    "worker_compute.extract_kind_invalid": {
        "ko": "'{kind}' 은(는) 추출 종류가 아닙니다. {allowed} 중 하나여야 합니다.",
        "en": "'{kind}' is not an extraction kind. Use one of {allowed}.",
    },
}


def pointer_tools_available() -> bool:
    """Whether this interpreter can import ``pyautogui``.

    Cheap and side-effect free: ``find_spec`` does not load the module, so a
    headless worker without a display is not punished for answering the
    question. The platform computer-use status route reads this through
    ``GET /worker/sysinfo`` rather than guessing from its own process.
    """
    return importlib.util.find_spec("pyautogui") is not None


def sysinfo_payload_extras() -> Dict[str, Any]:
    """Additive fields for ``GET /worker/sysinfo``.

    Existing GPU keys stay owned by ``worker_seams.probe_gpu_memory``. This
    dict is merged in; it must never reuse those names.
    """
    return {
        "capabilities": {"pointer_tools": pointer_tools_available()},
        "python_version": "{}.{}.{}".format(*sys.version_info[:3]),
    }


def register_worker_compute_messages() -> None:
    """Publish this module's messages into the one shared catalog."""
    for key, entry in WORKER_COMPUTE_MESSAGES.items():
        MESSAGES.setdefault(key, entry)


register_worker_compute_messages()


# ── shared helpers ──────────────────────────────────────────────────────────


@contextlib.contextmanager
def _temp_payload(data: bytes, suffix: str) -> Iterator[str]:
    """The posted bytes as a real file, deleted whatever happens next.

    Every parser and model port in this file takes a *path*: pdfplumber, PIL and
    a local transcriber all open files, and rewriting them to take buffers would
    be a much larger change than this seam is. So the bytes become a file for
    the duration of one call and nothing survives it — which is what "owns zero
    state" has to mean in practice.
    """
    handle = tempfile.NamedTemporaryFile(  # noqa: SIM115 — closed below, path reused
        prefix="ltcai-worker-", suffix=suffix, delete=False
    )
    try:
        handle.write(data)
        handle.close()
        yield handle.name
    finally:
        try:
            Path(handle.name).unlink(missing_ok=True)
        except OSError:
            quiet("worker compute temp cleanup")


def _suffix_for(filename: Optional[str], mime: Optional[str], default: str) -> str:
    """The extension to give the temp file: filename first, then MIME, then a default.

    The filename wins because it is what the user actually named the thing; a
    MIME type is what some client claimed about it. Neither is trusted for
    *content* — the parsers and PIL sniff that themselves — so this only has to
    be good enough to route.
    """
    suffix = Path(str(filename or "")).suffix.lower()
    if not suffix:
        subtype = str(mime or "").split("/")[-1].split(";")[0].strip().lower()
        suffix = f".{subtype}" if subtype else ""
    return suffix or default


def _callable_identity(fn: Any) -> str:
    """A name for an injected port, so the receipt says *which* model answered."""
    module = str(getattr(fn, "__module__", "") or "")
    name = str(getattr(fn, "__qualname__", "") or type(fn).__name__)
    return f"{module}.{name}" if module else name


# ── document builders: create_* minus the filesystem ────────────────────────


def build_docx_bytes(title: str, body: Any) -> bytes:
    """``tools.documents.create_docx``'s document, returned instead of saved."""
    try:
        from docx import Document
    except Exception as exc:  # noqa: BLE001 — an absent library is a configuration
        raise ToolError(
            "python-docx is not installed. Run `pip install -r requirements.txt`."
        ) from exc

    document = Document()
    if title:
        document.add_heading(str(title), level=1)
    for block in _body_to_str(body).split("\n\n"):
        text = block.strip()
        if text:
            document.add_paragraph(text)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def build_xlsx_bytes(rows: Any, sheet_name: str = "Sheet1") -> bytes:
    """``tools.documents.create_xlsx``'s workbook, returned instead of saved."""
    try:
        from openpyxl import Workbook
    except Exception as exc:  # noqa: BLE001 — an absent library is a configuration
        raise ToolError(
            "openpyxl is not installed. Run `pip install -r requirements.txt`."
        ) from exc

    if not isinstance(rows, list) or not all(isinstance(row, list) for row in rows):
        raise ToolError("Rows must be a list of lists.")

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = (sheet_name or "Sheet1")[:31]
    for row in rows:
        sheet.append(row)
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def build_pptx_bytes(title: str, slides: Any) -> bytes:
    """``tools.documents.create_pptx``'s deck, returned instead of saved."""
    try:
        from pptx import Presentation
    except Exception as exc:  # noqa: BLE001 — an absent library is a configuration
        raise ToolError(
            "python-pptx is not installed. Run `pip install -r requirements.txt`."
        ) from exc

    presentation = Presentation()
    first_layout = presentation.slide_layouts[0]
    first = presentation.slides.add_slide(first_layout)
    first.shapes.title.text = title or "Presentation"
    first.placeholders[1].text = ""

    content_layout = presentation.slide_layouts[1]
    for slide_data in slides or []:
        slide = presentation.slides.add_slide(content_layout)
        slide.shapes.title.text = str(slide_data.get("title") or "Slide")
        body = slide.placeholders[1].text_frame
        body.clear()
        bullets = slide_data.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        for index, bullet in enumerate(bullets):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = str(bullet)
            paragraph.level = 0

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def build_pdf_bytes(title: str, body: Any) -> bytes:
    """``tools.documents.create_pdf``'s PDF, returned instead of saved.

    The CJK font hunt is kept verbatim: without it a Korean document renders as
    a page of empty boxes, and that is not a cosmetic difference.
    """
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle
        from reportlab.lib.units import mm
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
    except Exception as exc:  # noqa: BLE001 — an absent library is a configuration
        raise ToolError("reportlab is not installed. Run `pip install reportlab`.") from exc

    font_name = "Helvetica"
    for font_path in CJK_FONT_CANDIDATES:
        if Path(font_path).exists():
            try:
                pdfmetrics.registerFont(TTFont("KoreanFont", font_path))
                font_name = "KoreanFont"
            except Exception:  # noqa: BLE001 — a broken font file is not a failed render
                quiet("CJK font registration")
            break

    title_style = ParagraphStyle("Title", fontName=font_name, fontSize=18, spaceAfter=8, leading=24)
    body_style = ParagraphStyle("Body", fontName=font_name, fontSize=11, spaceAfter=6, leading=16)

    story: List[Any] = []
    if title:
        story.append(Paragraph(str(title), title_style))
        story.append(Spacer(1, 4 * mm))

    for block in _body_to_str(body).split("\n\n"):
        text = block.strip()
        if text:
            safe_text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe_text, body_style))
            story.append(Spacer(1, 2 * mm))

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        leftMargin=20 * mm,
        rightMargin=20 * mm,
        topMargin=20 * mm,
        bottomMargin=20 * mm,
    )
    doc.build(story)
    return buffer.getvalue()


# ── request bodies ──────────────────────────────────────────────────────────


class EmbedRequest(BaseModel):
    """Texts in, vectors out. ``kind`` selects the write-path clamp."""

    texts: List[str] = Field(default_factory=list)
    kind: str = "passage"


class ParseRequest(BaseModel):
    """One document, by name and by bytes. ``filename`` picks the parser."""

    filename: str
    content_b64: str


class RenderDocxRequest(BaseModel):
    """``ToolDocxRequest``, verbatim — W3 forwards its body unchanged."""

    title: str = ""
    body: str = ""
    filename: str = "document.docx"


class RenderXlsxRequest(BaseModel):
    """``ToolXlsxRequest``, verbatim."""

    rows: List[List] = Field(default_factory=list)
    filename: str = "spreadsheet.xlsx"
    sheet_name: str = "Sheet1"


class RenderPptxRequest(BaseModel):
    """``ToolPptxRequest``, verbatim."""

    title: str = ""
    slides: List[Dict] = Field(default_factory=list)
    filename: str = "presentation.pptx"


class RenderPdfRequest(BaseModel):
    """``ToolPdfRequest``, verbatim."""

    title: str = ""
    body: str = ""
    filename: str = "document.pdf"


class AsrRequest(BaseModel):
    """One recording. ``filename``/``mime`` only choose the temp suffix."""

    audio_b64: str
    mime: Optional[str] = None
    filename: Optional[str] = None


class ExtractRequest(BaseModel):
    """The text one ingest door would hand ``_extract_concepts``."""

    text: str = ""
    kind: str = "message"


def build_extract_reply(text: str, kind: str) -> Dict[str, Any]:
    """The structures ``ingest_message`` / ``ingest_document`` / ``ingest_source`` consume.

    Concepts are already classified (``_classify_node_type``), triples carry the
    evidence class the edge metadata keeps, and semantic items keep the raw
    dict the node stores. This is the golden generator's ``_concepts`` /
    ``_triples`` / ``_semantic`` over the same functions the Python write path
    calls — a test asserts equality against those call sites, not against a
    second implementation.
    """
    from lattice_brain.graph._kg_common import (
        _classify_node_type,
        _extract_concepts,
        _extract_triples,
        _semantic_items,
    )

    limit = EXTRACT_LIMITS[kind]
    raw_concepts = _extract_concepts(text, limit=limit)
    concepts = [
        {"text": concept, "node_type": _classify_node_type(concept, text)}
        for concept in raw_concepts
    ]
    triples: List[Dict[str, Any]] = []
    for triple in _extract_triples(text, raw_concepts):
        triples.append(
            {
                "subject": triple["subject"],
                "object": triple["object"],
                "relation": triple["relation"],
                "weight": float(triple.get("weight") or 1.0),
                "context": str(triple.get("context") or ""),
                "evidence": str(triple.get("evidence") or ""),
                "confidence": triple.get("confidence"),
            }
        )
    semantic = [
        {
            "item_type": item["type"],
            "title": item["title"],
            "summary": item["summary"],
            "raw": dict(item),
        }
        for item in _semantic_items(text)
    ]
    return {"concepts": concepts, "triples": triples, "semantic": semantic}


def create_worker_compute_router(
    *,
    embedder: Any,
    transcriber: Optional[Callable[[str], str]] = None,
    require_user: Callable[[Request], Any],
    enforce_rate_limit: Callable[[str, str], None],
) -> APIRouter:
    """The eight compute seams, wired to what this worker actually resolved.

    ``embedder`` is the :class:`~latticeai.core.embedding_providers.text.ResolvedEmbedder`
    ``phase_brain`` built (``None`` ⇒ 503, because a worker with no embedder is
    a configuration rather than a crash). ``transcriber`` is the injected port
    the voice path holds; absent, ``/worker/asr`` reports the absence instead of
    inventing a transcript.
    """
    router = APIRouter()

    def _require_seam(request: Request) -> None:
        """404 unless the host opened the seam for this worker."""
        if not seam_open():
            raise http_error(404, "agent_seam.disabled", resolve_language(request))

    def _admit(request: Request) -> str:
        """Authenticate and charge this call against the per-step budget."""
        current_user = require_user(request)
        enforce_rate_limit(current_user, SEAM_RATE_BUCKET)
        return str(current_user or "")

    def _decode(payload: str, language: str) -> bytes:
        """Base64 in, bytes out — a malformed body is 422, never a 500."""
        try:
            return base64.b64decode(str(payload or ""), validate=True)
        except (binascii.Error, ValueError) as exc:
            raise http_error(
                422, "worker_compute.content_invalid", language, reason=str(exc)
            ) from exc

    async def _render(
        kind: str,
        builder: Callable[[], bytes],
        language: str,
        **extra: Any,
    ) -> Dict[str, Any]:
        """Build one document off the event loop and hand back its bytes.

        The reply is the bytes and what they cost, and nothing about *where*
        they go: ``lattice-agent``'s ``documents::document_output_target`` runs
        its own ``safe_filename`` and ``Workspace::resolve`` over the request's
        ``filename`` **before** it calls here, and writes to the target it
        resolved. A second sanitisation on this side produced a name no caller
        ever read — two spellings of one rule, one of them invisible.
        """
        try:
            payload = await asyncio.to_thread(builder)
        except ToolError as exc:
            # The library is not installed here. That is a property of the
            # install, not of the request, so it is 503 and not 500.
            raise http_error(
                503, "worker_compute.render_unavailable", language, kind=kind, reason=str(exc)
            ) from exc
        except Exception as exc:  # noqa: BLE001 — reported with the kind, not swallowed
            logger.warning("worker render %s failed: %s", kind, exc)
            raise http_error(
                500, "worker_compute.render_failed", language, kind=kind, reason=str(exc)
            ) from exc
        return {
            "content_b64": base64.b64encode(payload).decode("ascii"),
            "bytes": len(payload),
            **extra,
        }

    @router.post("/worker/embed")
    async def worker_embed(req: EmbedRequest, request: Request):
        """Vectors for these texts, from the provider this worker resolved.

        ``dim`` and ``model_id`` are read *after* the call because a network
        provider locks its index identity to the width the model actually
        returned — reading them first would report the guess, and Rust files
        every vector under whatever this says.
        """
        _require_seam(request)
        _admit(request)
        language = resolve_language(request)
        if embedder is None:
            raise http_error(503, "worker_compute.embedder_unavailable", language)
        kind = req.kind.strip().lower()
        if kind not in EMBED_KINDS:
            raise http_error(
                422,
                "worker_compute.kind_invalid",
                language,
                kind=req.kind,
                allowed=", ".join(EMBED_KINDS),
            )
        provider = embedder.provider
        texts = list(req.texts)
        if kind == "passage":
            texts = [text[:PASSAGE_MAX_CHARS] for text in texts]
        vectors = await asyncio.to_thread(provider.embed_batch, texts)
        return {
            "vectors": vectors,
            "dim": provider.dim,
            "provider": embedder.active,
            "model_id": provider.model_id,
            "kind": kind,
        }

    @router.post("/worker/parse")
    async def worker_parse(req: ParseRequest, request: Request):
        """The parser matrix over posted bytes: text, counts, preview.

        The ``path`` key ``read_document`` reports is dropped: it named a
        temporary file that no longer exists by the time this answers, and
        echoing it would hand Rust a path it must not act on. ``filename`` — the
        one the caller named — takes its place.
        """
        _require_seam(request)
        _admit(request)
        language = resolve_language(request)
        data = _decode(req.content_b64, language)
        suffix = Path(req.filename or "").suffix.lower()
        with _temp_payload(data, suffix) as tmp_path:
            try:
                parsed = await asyncio.to_thread(read_document, tmp_path)
            except ToolError as exc:
                raise http_error(
                    400, "worker_compute.parse_failed", language, reason=str(exc)
                ) from exc
        parsed.pop("path", None)
        parsed["filename"] = req.filename
        return parsed

    @router.post("/worker/render/docx")
    async def worker_render_docx(req: RenderDocxRequest, request: Request):
        """A Word document's bytes. Rust decides where it lands."""
        _require_seam(request)
        _admit(request)
        return await _render(
            "docx",
            lambda: build_docx_bytes(req.title, req.body),
            resolve_language(request),
        )

    @router.post("/worker/render/xlsx")
    async def worker_render_xlsx(req: RenderXlsxRequest, request: Request):
        """A workbook's bytes, with the row count ``create_xlsx`` reported."""
        _require_seam(request)
        _admit(request)
        return await _render(
            "xlsx",
            lambda: build_xlsx_bytes(req.rows, req.sheet_name),
            resolve_language(request),
            rows=len(req.rows),
        )

    @router.post("/worker/render/pptx")
    async def worker_render_pptx(req: RenderPptxRequest, request: Request):
        """A deck's bytes. ``slides`` counts the title slide, as ``create_pptx`` did."""
        _require_seam(request)
        _admit(request)
        return await _render(
            "pptx",
            lambda: build_pptx_bytes(req.title, req.slides),
            resolve_language(request),
            slides=len(req.slides) + 1,
        )

    @router.post("/worker/render/pdf")
    async def worker_render_pdf(req: RenderPdfRequest, request: Request):
        """A PDF's bytes, CJK font and all."""
        _require_seam(request)
        _admit(request)
        return await _render(
            "pdf",
            lambda: build_pdf_bytes(req.title, req.body),
            resolve_language(request),
        )

    @router.post("/worker/asr")
    async def worker_asr(req: AsrRequest, request: Request):
        """Transcribe one recording — and never claim a transcript it did not get.

        ``status`` is the state machine ``VoiceCaptureService.capture`` records
        on the node: ``ok`` when there are words, ``unavailable`` when this
        machine has no transcriber or the transcriber returned nothing, and
        ``failed`` when it broke. All three are 200: "we could not hear it" is
        an answer about the audio, not an error about the request.
        """
        _require_seam(request)
        _admit(request)
        language = resolve_language(request)
        data = _decode(req.audio_b64, language)
        suffix = _suffix_for(req.filename, req.mime, ".m4a")
        if suffix not in SUPPORTED_AUDIO_EXTENSIONS:
            raise http_error(
                400,
                "worker_compute.audio_unsupported",
                language,
                suffix=suffix,
                allowed=", ".join(sorted(SUPPORTED_AUDIO_EXTENSIONS)),
            )
        if len(data) > MAX_AUDIO_BYTES:
            raise http_error(
                413,
                "worker_compute.audio_too_large",
                language,
                size=len(data),
                limit=MAX_AUDIO_BYTES,
            )
        if transcriber is None:
            return {
                "text": "",
                "segments": None,
                "provider": "",
                "status": "unavailable",
                "detail": "no local transcriber is configured",
            }
        provider = _callable_identity(transcriber)
        with _temp_payload(data, suffix) as tmp_path:
            try:
                raw = await asyncio.to_thread(transcriber, tmp_path)
            except Exception as exc:  # noqa: BLE001 — a broken transcriber is a state
                logger.warning("worker asr failed: %s", exc)
                return {
                    "text": "",
                    "segments": None,
                    "provider": provider,
                    "status": "failed",
                    "detail": str(exc),
                }
        text = str(raw or "").strip()
        if not text:
            # An empty transcript is not text: reporting it as one would put an
            # empty note in the Brain and call it a memory.
            return {
                "text": "",
                "segments": None,
                "provider": provider,
                "status": "unavailable",
                "detail": "the transcriber returned no text",
            }
        return {
            "text": text,
            # The port's contract is ``(path) -> str``. There are no timings to
            # report, and inventing an empty list would read as "no speech".
            "segments": None,
            "provider": provider,
            "status": "ok",
            "detail": "",
        }

    @router.post("/worker/extract")
    async def worker_extract(req: ExtractRequest, request: Request):
        """Concepts, triples and Task/Decision items for this text.

        The reply is field-for-field what ``GraphWriter``'s ingest doors take
        as ``concepts`` / ``triples`` / ``semantic``. Classification and the
        evidence class happen here because they are compute over the text, not
        a write.
        """
        _require_seam(request)
        _admit(request)
        language = resolve_language(request)
        kind = req.kind.strip().lower()
        if kind not in EXTRACT_KINDS:
            raise http_error(
                422,
                "worker_compute.extract_kind_invalid",
                language,
                kind=req.kind,
                allowed=", ".join(EXTRACT_KINDS),
            )
        return await asyncio.to_thread(build_extract_reply, req.text, kind)

    return router


__all__ = [
    "CJK_FONT_CANDIDATES",
    "EMBED_KINDS",
    "EXTRACT_KINDS",
    "EXTRACT_LIMITS",
    "PASSAGE_MAX_CHARS",
    "WORKER_COMPUTE_MESSAGES",
    "AsrRequest",
    "EmbedRequest",
    "ExtractRequest",
    "ParseRequest",
    "RenderDocxRequest",
    "RenderPdfRequest",
    "RenderPptxRequest",
    "RenderXlsxRequest",
    "build_docx_bytes",
    "build_extract_reply",
    "build_pdf_bytes",
    "build_pptx_bytes",
    "build_xlsx_bytes",
    "create_worker_compute_router",
    "pointer_tools_available",
    "register_worker_compute_messages",
    "sysinfo_payload_extras",
]
