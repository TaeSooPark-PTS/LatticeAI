"""Local file → text: the parsers behind local indexing.

PDF/Word/Excel/PowerPoint text plus the image signal path (dimensions, OCR,
and — only when a VLM exists — a caption). Moved verbatim out of
``discovery_index.py`` (v11.3.0 decomposition).
"""

from __future__ import annotations

from typing import TYPE_CHECKING

# ruff: noqa: F403,F405
from .._kg_common import *  # noqa: F403,F401

# The cross-mixin surface (`_connect`, `_upsert_node`, …) is declared in
# `_kg_contract.KnowledgeGraphCore`. It is a typing-only base: at runtime this
# is `object`, so the MRO of `KnowledgeGraphStore` is unchanged.
if TYPE_CHECKING:
    from .._kg_contract import KnowledgeGraphCore as _Core
else:
    _Core = object


class _LocalExtractMixin(_Core):
    """File-text and image-signal extraction. Composed into the public mixin."""

    def _extract_local_file_text(
        self, path: Path, category: str, *, include_ocr: bool
    ) -> Tuple[str, Dict[str, Any]]:
        ext = path.suffix.lower()
        meta: Dict[str, Any] = {"parser": _parser_type_for_category(category, ext)}
        text = ""
        if category in {"text", "code"} or ext == ".csv":
            text = path.read_text(encoding="utf-8", errors="replace")
        elif ext == ".pdf":
            import pdfplumber

            with pdfplumber.open(str(path)) as pdf:
                meta["pages"] = len(pdf.pages)
                text = "\n\n".join((page.extract_text() or "") for page in pdf.pages)
        elif ext == ".docx":
            from docx import Document

            doc = Document(str(path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            table_lines = []
            for table in doc.tables:
                for row in table.rows:
                    cells = [_clean_text(cell.text) for cell in row.cells]
                    if any(cells):
                        table_lines.append("\t".join(cells))
            meta["paragraphs"] = len(paragraphs)
            meta["tables"] = len(doc.tables)
            meta["table_rows"] = len(table_lines)
            text = "\n\n".join([*paragraphs, *table_lines])
        elif ext == ".xlsx":
            from openpyxl import load_workbook

            wb = load_workbook(str(path), read_only=True, data_only=True)
            rows_all = []
            non_empty_rows = 0
            non_empty_cells = 0
            char_count = 0
            for ws in wb.worksheets:
                sheet_rows = []
                for row in ws.iter_rows(values_only=True):
                    cells = [
                        str(cell).strip() if cell is not None else "" for cell in row
                    ]
                    if not any(cells):
                        continue
                    line = "\t".join(cells)
                    non_empty_rows += 1
                    non_empty_cells += sum(1 for cell in cells if cell)
                    sheet_rows.append(line)
                    char_count += len(line) + 1
                    if char_count > 200_000:
                        break
                if sheet_rows:
                    rows_all.append(f"[Sheet: {ws.title}]")
                    rows_all.extend(sheet_rows)
                if char_count > 200_000:
                    break
            meta["sheets"] = len(wb.worksheets)
            meta["rows"] = non_empty_rows
            meta["cells"] = non_empty_cells
            text = "\n".join(rows_all)
        elif ext == ".pptx":
            from pptx import Presentation

            prs = Presentation(str(path))
            slides_text = []
            for index, slide in enumerate(prs.slides, 1):
                parts = []
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        slide_text = shape.text_frame.text.strip()
                        if slide_text:
                            parts.append(slide_text)
                if parts:
                    slides_text.append(f"[Slide {index}]\n" + "\n".join(parts))
            meta["slides"] = len(prs.slides)
            meta["text_slides"] = len(slides_text)
            text = "\n\n".join(slides_text)
        elif category == "image":
            text = self._extract_image_signals(path, meta, include_ocr=include_ocr)
        return text[:200_000], meta

    def _extract_image_signals(
        self, path: Path, meta: Dict[str, Any], *, include_ocr: bool
    ) -> str:
        """Dimensions, OCR text, and — only if a VLM exists — a caption.

        Until v11.1.0 this path always attached a ``vision_caption`` built out
        of the filename and the pixel dimensions (``Image pic.png (PNG 12x8)``)
        and used it as the retrieval text. Nothing downstream could tell that
        string apart from something a vision model had actually said about the
        picture, so every screenshot in the graph carried a fake description.

        Now the caption comes from the injected port and from nowhere else. A
        picture with no OCR text and no model still gets indexed — under its
        filename, which is a fact — and ``caption_status`` says why there is no
        caption.
        """
        from ...multimodal import MultimodalPorts, extract_image_facts

        ports = getattr(self, "multimodal_ports", None) or MultimodalPorts()
        facts = extract_image_facts(str(path), ports=ports, ocr=include_ocr)
        meta.update(facts.as_metadata())
        meta["ocr_enabled"] = bool(include_ocr)
        if facts.ocr_text:
            meta["ocr_chars"] = len(facts.ocr_text)
        if facts.ocr_status == "failed":
            meta["ocr_error"] = facts.ocr_detail
        if not facts.readable:
            return ""
        return facts.index_text() or path.name
