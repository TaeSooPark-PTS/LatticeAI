from __future__ import annotations

from typing import TYPE_CHECKING

# ruff: noqa: F403,F405
from ._kg_common import *  # noqa: F403,F401

# The cross-mixin surface (`_connect`, `_upsert_node`, …) is declared in
# `_kg_contract.KnowledgeGraphCore`. It is a typing-only base: at runtime this
# is `object`, so the MRO of `KnowledgeGraphStore` is unchanged.
if TYPE_CHECKING:
    from ._kg_contract import KnowledgeGraphCore as _Core
else:
    _Core = object



def _local_scoped_slug(prefix: str, value: str, workspace_id: Optional[str]) -> str:
    """Preserve legacy IDs while isolating newly workspace-scoped nodes."""
    slug = _slug(value)
    if not workspace_id:
        return f"{prefix}:{slug}"
    scope = _sha256_text(str(workspace_id))[:12]
    return f"{prefix}:{scope}:{slug}"


class KnowledgeGraphLocalIndexMixin(_Core):
    """Local file → graph indexing (text extraction, node/index upserts,
    graph-node deletion, orphan cleanup, and the index_local_folder driver),
    split out of discovery. Composed into KnowledgeGraphStore alongside
    KnowledgeGraphDiscoveryMixin; both share the instance so these methods
    still reach sibling discovery/write helpers through the class MRO.
    """

    def _extract_local_file_text(
        self, path: Path, category: str, *, include_ocr: bool
    ) -> Tuple[str, Dict[str, Any]]:
        ext = path.suffix.lower()
        meta: Dict[str, Any] = {"parser": _parser_type_for_category(category, ext)}
        text = ""
        if category in {"text", "code"} or ext == ".csv":
            text = path.read_text(encoding="utf-8", errors="replace")
        elif ext == ".pdf":
            import pdfplumber

            with pdfplumber.open(str(path)) as pdf:
                meta["pages"] = len(pdf.pages)
                text = "\n\n".join((page.extract_text() or "") for page in pdf.pages)
        elif ext == ".docx":
            from docx import Document

            doc = Document(str(path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            table_lines = []
            for table in doc.tables:
                for row in table.rows:
                    cells = [_clean_text(cell.text) for cell in row.cells]
                    if any(cells):
                        table_lines.append("\t".join(cells))
            meta["paragraphs"] = len(paragraphs)
            meta["tables"] = len(doc.tables)
            meta["table_rows"] = len(table_lines)
            text = "\n\n".join([*paragraphs, *table_lines])
        elif ext == ".xlsx":
            from openpyxl import load_workbook

            wb = load_workbook(str(path), read_only=True, data_only=True)
            rows_all = []
            non_empty_rows = 0
            non_empty_cells = 0
            char_count = 0
            for ws in wb.worksheets:
                sheet_rows = []
                for row in ws.iter_rows(values_only=True):
                    cells = [
                        str(cell).strip() if cell is not None else "" for cell in row
                    ]
                    if not any(cells):
                        continue
                    line = "\t".join(cells)
                    non_empty_rows += 1
                    non_empty_cells += sum(1 for cell in cells if cell)
                    sheet_rows.append(line)
                    char_count += len(line) + 1
                    if char_count > 200_000:
                        break
                if sheet_rows:
                    rows_all.append(f"[Sheet: {ws.title}]")
                    rows_all.extend(sheet_rows)
                if char_count > 200_000:
                    break
            meta["sheets"] = len(wb.worksheets)
            meta["rows"] = non_empty_rows
            meta["cells"] = non_empty_cells
            text = "\n".join(rows_all)
        elif ext == ".pptx":
            from pptx import Presentation

            prs = Presentation(str(path))
            slides_text = []
            for index, slide in enumerate(prs.slides, 1):
                parts = []
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        slide_text = shape.text_frame.text.strip()
                        if slide_text:
                            parts.append(slide_text)
                if parts:
                    slides_text.append(f"[Slide {index}]\n" + "\n".join(parts))
            meta["slides"] = len(prs.slides)
            meta["text_slides"] = len(slides_text)
            text = "\n\n".join(slides_text)
        elif category == "image":
            text = self._extract_image_signals(path, meta, include_ocr=include_ocr)
        return text[:200_000], meta

    def _extract_image_signals(
        self, path: Path, meta: Dict[str, Any], *, include_ocr: bool
    ) -> str:
        """Dimensions, OCR text, and — only if a VLM exists — a caption.

        Until v11.1.0 this path always attached a ``vision_caption`` built out
        of the filename and the pixel dimensions (``Image pic.png (PNG 12x8)``)
        and used it as the retrieval text. Nothing downstream could tell that
        string apart from something a vision model had actually said about the
        picture, so every screenshot in the graph carried a fake description.

        Now the caption comes from the injected port and from nowhere else. A
        picture with no OCR text and no model still gets indexed — under its
        filename, which is a fact — and ``caption_status`` says why there is no
        caption.
        """
        from ..multimodal import MultimodalPorts, extract_image_facts

        ports = getattr(self, "multimodal_ports", None) or MultimodalPorts()
        facts = extract_image_facts(str(path), ports=ports, ocr=include_ocr)
        meta.update(facts.as_metadata())
        meta["ocr_enabled"] = bool(include_ocr)
        if facts.ocr_text:
            meta["ocr_chars"] = len(facts.ocr_text)
        if facts.ocr_status == "failed":
            meta["ocr_error"] = facts.ocr_detail
        if not facts.readable:
            return ""
        return facts.index_text() or path.name

    def _ensure_local_hierarchy(
        self,
        conn: sqlite3.Connection,
        *,
        source_id: str,
        root: Path,
        file_path: Path,
        os_type: str,
        drive_id: str,
        user_email: Optional[str] = None,
        workspace_id: Optional[str] = None,
    ) -> str:
        computer_label = platform.node() or "내 컴퓨터"
        computer_id = _local_scoped_slug("computer", computer_label, workspace_id)
        drive_identity = f"{workspace_id}|{os_type}:{drive_id}" if workspace_id else f"{os_type}:{drive_id}"
        drive_node_id = f"drive:{_sha256_text(drive_identity)[:24]}"
        root_folder_id = f"folder:{_sha256_text(f'{source_id}:root')[:24]}"
        self._upsert_node(
            conn,
            computer_id,
            "Computer",
            computer_label,
            metadata={"os_type": os_type, "workspace_id": workspace_id},
            owner=user_email,
            workspace_id=workspace_id,
        )
        self._upsert_node(
            conn,
            drive_node_id,
            "Drive",
            drive_id,
            metadata={"os_type": os_type, "drive_id": drive_id, "workspace_id": workspace_id},
            owner=user_email,
            workspace_id=workspace_id,
        )
        stale_parents = conn.execute(
            """
                SELECT e.from_node
                FROM edges e
                JOIN nodes n ON n.id=e.from_node
                WHERE e.to_node=? AND n.type='Drive' AND e.from_node<>?
                """,
            (root_folder_id, drive_node_id),
        ).fetchall()
        for row in stale_parents:
            conn.execute(
                "DELETE FROM edges WHERE from_node=? AND to_node=?",
                (row["from_node"], root_folder_id),
            )
            conn.execute(
                "DELETE FROM edges_v2 WHERE source=? AND target=?",
                (row["from_node"], root_folder_id),
            )
        self._upsert_edge(
            conn,
            computer_id,
            drive_node_id,
            "포함함",
            metadata={"source": "local_scan", "workspace_id": workspace_id},
        )
        self._upsert_node(
            conn,
            root_folder_id,
            "Folder",
            root.name or str(root),
            summary=str(root),
            metadata={"source_id": source_id, "path": str(root), "root": True, "workspace_id": workspace_id},
            owner=user_email,
            workspace_id=workspace_id,
        )
        self._upsert_edge(
            conn,
            drive_node_id,
            root_folder_id,
            "포함함",
            metadata={"source": "local_scan", "workspace_id": workspace_id},
        )

        try:
            relative_parent = file_path.parent.relative_to(root)
        except ValueError:
            relative_parent = Path()
        parent_id = root_folder_id
        current_path = root
        for part in relative_parent.parts:
            current_path = current_path / part
            folder_id = (
                f"folder:{_sha256_text(f'{source_id}:{current_path.as_posix()}')[:24]}"
            )
            self._upsert_node(
                conn,
                folder_id,
                "Folder",
                part,
                summary=str(current_path),
                metadata={
                    "source_id": source_id,
                    "path": str(current_path),
                    "root": False,
                    "workspace_id": workspace_id,
                },
                owner=user_email,
                workspace_id=workspace_id,
            )
            self._upsert_edge(
                conn,
                parent_id,
                folder_id,
                "포함함",
                metadata={"source": "local_scan", "workspace_id": workspace_id},
            )
            parent_id = folder_id
        return parent_id

    def _upsert_local_file_index(
        self,
        conn: sqlite3.Connection,
        *,
        source_id: str,
        root: Path,
        file_path: Path,
        stat: Optional[os.stat_result],
        os_type: str,
        drive_id: str,
        status: str,
        parser_type: str,
        sha256: Optional[str] = None,
        graph_node_id: Optional[str] = None,
        error_message: Optional[str] = None,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> str:
        try:
            relative_path = file_path.relative_to(root).as_posix()
        except ValueError:
            relative_path = file_path.name
        index_id = f"local-index:{_sha256_text(f'{source_id}:{relative_path}')[:24]}"
        now = _now()
        size = stat.st_size if stat else None
        modified_at = _safe_iso_from_stat_mtime(stat.st_mtime) if stat else ""
        conn.execute(
            """
                INSERT INTO local_file_index(
                  id, source_id, os_type, drive_id, root_path, file_path, relative_path,
                  file_name, extension, size_bytes, modified_at, sha256, last_scanned_at,
                  last_indexed_at, parser_type, status, error_message, graph_node_id,
                  deleted, metadata_json
                )
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                ON CONFLICT(source_id, relative_path) DO UPDATE SET
                  os_type=excluded.os_type,
                  drive_id=excluded.drive_id,
                  root_path=excluded.root_path,
                  file_path=excluded.file_path,
                  file_name=excluded.file_name,
                  extension=excluded.extension,
                  size_bytes=excluded.size_bytes,
                  modified_at=excluded.modified_at,
                  sha256=excluded.sha256,
                  last_scanned_at=excluded.last_scanned_at,
                  last_indexed_at=excluded.last_indexed_at,
                  parser_type=excluded.parser_type,
                  status=excluded.status,
                  error_message=excluded.error_message,
                  graph_node_id=excluded.graph_node_id,
                  deleted=excluded.deleted,
                  metadata_json=excluded.metadata_json
                """,
            (
                index_id,
                source_id,
                os_type,
                drive_id,
                str(root),
                str(file_path),
                relative_path,
                file_path.name,
                file_path.suffix.lower(),
                size,
                modified_at,
                sha256,
                now,
                now if status == "indexed" else None,
                parser_type,
                status,
                error_message,
                graph_node_id,
                0 if status != "deleted" else 1,
                _json(metadata),
            ),
        )
        return index_id

    def _upsert_local_file_node(
        self,
        conn: sqlite3.Connection,
        *,
        source_id: str,
        root: Path,
        file_path: Path,
        stat: os.stat_result,
        os_type: str,
        drive_id: str,
        sha256: str,
        category: str,
        parser_type: str,
        text: str,
        parser_meta: Dict[str, Any],
        user_email: Optional[str] = None,
        workspace_id: Optional[str] = None,
    ) -> str:
        text = _clean_text(text)
        if not text:
            raise ValueError("텍스트 추출 결과가 비어 있습니다.")
        try:
            relative_path = file_path.relative_to(root).as_posix()
        except ValueError:
            relative_path = file_path.name
        file_node_id = f"local-file:{_sha256_text(f'{source_id}:{relative_path}')[:24]}"
        parent_folder_id = self._ensure_local_hierarchy(
            conn,
            source_id=source_id,
            root=root,
            file_path=file_path,
            os_type=os_type,
            drive_id=drive_id,
            user_email=user_email,
            workspace_id=workspace_id,
        )
        linked_rows = conn.execute(
            """
                SELECT e.to_node AS id, n.type, n.metadata_json
                FROM edges e
                JOIN nodes n ON n.id=e.to_node
                WHERE e.from_node=?
                """,
            (file_node_id,),
        ).fetchall()
        child_ids = []
        auto_candidate_ids = set()
        for row in linked_rows:
            linked_metadata = _safe_loads(row["metadata_json"])
            if row["type"] in {"Chunk", "ImageText", "Section"} or linked_metadata.get("source_node") == file_node_id:
                child_ids.append(row["id"])
            elif linked_metadata.get("auto_extracted") and linked_metadata.get("source") == "local_folder":
                auto_candidate_ids.add(row["id"])
        conn.execute("DELETE FROM chunks WHERE source_node=?", (file_node_id,))
        if child_ids:
            placeholders = ",".join("?" * len(child_ids))
            conn.execute(f"DELETE FROM nodes WHERE id IN ({placeholders})", child_ids)
            self._v2_delete_nodes(conn, child_ids)
        conn.execute("DELETE FROM edges WHERE from_node=?", (file_node_id,))
        self._v2_delete_edges_from(conn, file_node_id)
        removable_auto_ids = set()
        for node_id in auto_candidate_ids:
            remaining_edges = conn.execute(
                "SELECT from_node, to_node FROM edges WHERE from_node=? OR to_node=?",
                (node_id, node_id),
            ).fetchall()
            if all(
                row["from_node"] in auto_candidate_ids
                and row["to_node"] in auto_candidate_ids
                for row in remaining_edges
            ):
                removable_auto_ids.add(node_id)
        if removable_auto_ids:
            placeholders = ",".join("?" * len(removable_auto_ids))
            params = list(removable_auto_ids)
            conn.execute(
                f"DELETE FROM edges WHERE from_node IN ({placeholders}) OR to_node IN ({placeholders})",
                params * 2,
            )
            conn.execute(f"DELETE FROM nodes WHERE id IN ({placeholders})", params)
            self._v2_delete_nodes(conn, params)

        metadata = {
            "source": "local_folder",
            "source_id": source_id,
            "root_path": str(root),
            "file_path": str(file_path),
            "relative_path": relative_path,
            "filename": file_path.name,
            "ext": file_path.suffix.lower(),
            "category": category,
            "parser_type": parser_type,
            "bytes": stat.st_size,
            "modified_at": _safe_iso_from_stat_mtime(stat.st_mtime),
            "sha256": sha256,
            "parser": parser_meta,
            "workspace_id": workspace_id,
        }
        self._upsert_node(
            conn,
            file_node_id,
            _node_type_for_category(category),
            file_path.name,
            summary=text[:700],
            metadata=metadata,
            raw=metadata,
            owner=user_email,
            workspace_id=workspace_id,
        )
        self._upsert_edge(
            conn,
            parent_folder_id,
            file_node_id,
            "포함함",
            weight=1.0,
            metadata={"source": "local_scan", "workspace_id": workspace_id},
        )
        self._cleanup_local_graph_orphans(conn, source_id)

        target_for_concepts = text
        if category == "image" and text:
            image_text_id = f"imagetext:{_sha256_text(f'{file_node_id}:ocr')[:24]}"
            self._upsert_node(
                conn,
                image_text_id,
                "ImageText",
                f"{file_path.name} OCR",
                summary=_clean_text(text)[:700],
                metadata={
                    "source_node": file_node_id,
                    "source_id": source_id,
                    "chars": len(text),
                    "workspace_id": workspace_id,
                },
                owner=user_email,
                workspace_id=workspace_id,
            )
            self._upsert_edge(
                conn,
                file_node_id,
                image_text_id,
                "포함함",
                weight=0.8,
                metadata={"source": "ocr", "workspace_id": workspace_id},
            )

        # Typed chunking by file extension (markdown/code/plain); plain files
        # keep legacy _chunks boundaries and therefore identical chunk ids.
        chunk_strategy = chunk_strategy_for(file_path.name)
        for index, piece in enumerate(typed_chunks(text, strategy=chunk_strategy)):
            chunk = piece["text"]
            chunk_fields = typed_chunk_meta_fields(piece)
            chunk_id = f"chunk:{_sha256_text(f'{file_node_id}:{index}:{chunk}')[:24]}"
            self._upsert_node(
                conn,
                chunk_id,
                "Chunk",
                f"{file_path.name} chunk {index + 1}",
                summary=chunk[:500],
                metadata={
                    "index": index,
                    "source_node": file_node_id,
                    "source_id": source_id,
                    "workspace_id": workspace_id,
                    **chunk_fields,
                },
                owner=user_email,
                workspace_id=workspace_id,
            )
            self._upsert_chunk(
                conn,
                chunk_id=chunk_id,
                source_node=file_node_id,
                text=chunk,
                metadata={
                    "index": index,
                    "source_node": file_node_id,
                    "source_id": source_id,
                    "workspace_id": workspace_id,
                    **chunk_fields,
                },
            )
            self._upsert_edge(
                conn,
                file_node_id,
                chunk_id,
                "포함함",
                weight=0.7,
                metadata={"source": "local_scan", "workspace_id": workspace_id},
            )

        concepts = _extract_concepts(target_for_concepts, limit=18)
        concept_ids: Dict[str, str] = {}
        for concept in concepts:
            node_t = _classify_node_type(concept, target_for_concepts)
            concept_id = _local_scoped_slug(node_t.lower(), concept, workspace_id)
            concept_ids[concept.lower()] = concept_id
            self._upsert_node(
                conn,
                concept_id,
                node_t,
                concept,
                metadata={
                    "auto_extracted": True,
                    "source": "local_folder",
                    "source_id": source_id,
                    "workspace_id": workspace_id,
                },
                owner=user_email,
                workspace_id=workspace_id,
            )
            self._upsert_edge(
                conn,
                file_node_id,
                concept_id,
                "언급함",
                weight=0.75,
                metadata={"source": "local_scan", "workspace_id": workspace_id},
            )

        for triple in _extract_triples(target_for_concepts, concepts, limit=20):
            subj_id = concept_ids.get(triple["subject"].lower())
            obj_id = concept_ids.get(triple["object"].lower())
            if subj_id and obj_id and subj_id != obj_id:
                self._upsert_edge(
                    conn,
                    subj_id,
                    obj_id,
                    triple["relation"],
                    weight=0.9,
                    metadata={
                        "context": triple.get("context", "")[:240],
                        "source_id": source_id,
                        "workspace_id": workspace_id,
                    },
                )

        for item in _semantic_items(target_for_concepts):
            sem_type = item["type"]
            sem_title = item["title"]
            sem_id = f"{sem_type.lower()}:{_sha256_text(f'{file_node_id}:{sem_type}:{sem_title}')[:24]}"
            self._upsert_node(
                conn,
                sem_id,
                sem_type,
                sem_title,
                summary=item["summary"],
                metadata={
                    "auto_extracted": True,
                    "source_node": file_node_id,
                    "filename": file_path.name,
                    "workspace_id": workspace_id,
                },
                raw=item,
                owner=user_email,
                workspace_id=workspace_id,
            )
            self._upsert_edge(
                conn,
                file_node_id,
                sem_id,
                "포함함",
                weight=0.9,
                metadata={"source": "local_scan", "workspace_id": workspace_id},
            )

        return file_node_id

    def _delete_local_file_graph(
        self, conn: sqlite3.Connection, file_node_id: Optional[str]
    ) -> None:
        if not file_node_id:
            return

        file_row = conn.execute(
            "SELECT metadata_json FROM nodes WHERE id=?",
            (file_node_id,),
        ).fetchone()
        source_id = None
        if file_row:
            source_id = _safe_loads(file_row["metadata_json"]).get("source_id")

        linked_rows = conn.execute(
            """
                SELECT n.id, n.type, n.metadata_json
                FROM edges e
                JOIN nodes n ON n.id=e.to_node
                WHERE e.from_node=?
                """,
            (file_node_id,),
        ).fetchall()
        owned_ids: set = set()
        auto_candidate_ids: set = set()
        for row in linked_rows:
            metadata = _safe_loads(row["metadata_json"])
            if (
                row["type"] in {"Chunk", "ImageText", "Section"}
                or metadata.get("source_node") == file_node_id
            ):
                owned_ids.add(row["id"])
            elif (
                metadata.get("auto_extracted")
                and metadata.get("source") == "local_folder"
            ):
                auto_candidate_ids.add(row["id"])

        conn.execute("DELETE FROM chunks WHERE source_node=?", (file_node_id,))
        conn.execute(
            "DELETE FROM edges WHERE from_node=? OR to_node=?",
            (file_node_id, file_node_id),
        )
        conn.execute("DELETE FROM nodes WHERE id=?", (file_node_id,))
        self._v2_delete_nodes(conn, [file_node_id])

        def delete_nodes(node_ids: set) -> None:
            if not node_ids:
                return
            placeholders = ",".join("?" * len(node_ids))
            params = list(node_ids)
            conn.execute(
                f"DELETE FROM chunks WHERE source_node IN ({placeholders})", params
            )
            conn.execute(
                f"DELETE FROM edges WHERE from_node IN ({placeholders}) OR to_node IN ({placeholders})",
                params * 2,
            )
            conn.execute(f"DELETE FROM nodes WHERE id IN ({placeholders})", params)
            self._v2_delete_nodes(conn, params)

        delete_nodes(owned_ids)

        removable_auto_ids: set = set()
        for node_id in auto_candidate_ids:
            remaining_edges = conn.execute(
                "SELECT from_node, to_node FROM edges WHERE from_node=? OR to_node=?",
                (node_id, node_id),
            ).fetchall()
            if all(
                (
                    row["from_node"] in auto_candidate_ids
                    and row["to_node"] in auto_candidate_ids
                )
                for row in remaining_edges
            ):
                removable_auto_ids.add(node_id)
        delete_nodes(removable_auto_ids)
        if source_id:
            self._cleanup_local_graph_orphans(conn, str(source_id))

    def _cleanup_local_graph_orphans(
        self, conn: sqlite3.Connection, source_id: str
    ) -> None:
        while True:
            folder_rows = conn.execute(
                "SELECT id, metadata_json FROM nodes WHERE type='Folder'"
            ).fetchall()
            leaf_ids = []
            for row in folder_rows:
                metadata = _safe_loads(row["metadata_json"])
                if metadata.get("source_id") != source_id:
                    continue
                has_children = conn.execute(
                    "SELECT 1 FROM edges WHERE from_node=? LIMIT 1",
                    (row["id"],),
                ).fetchone()
                if not has_children:
                    leaf_ids.append(row["id"])
            if not leaf_ids:
                break
            placeholders = ",".join("?" * len(leaf_ids))
            conn.execute(
                f"DELETE FROM edges WHERE from_node IN ({placeholders}) OR to_node IN ({placeholders})",
                leaf_ids * 2,
            )
            conn.execute(f"DELETE FROM nodes WHERE id IN ({placeholders})", leaf_ids)
            self._v2_delete_nodes(conn, leaf_ids)

        for node_type in ("Drive", "Computer"):
            rows = conn.execute(
                "SELECT id FROM nodes WHERE type=?", (node_type,)
            ).fetchall()
            removable = []
            for row in rows:
                has_children = conn.execute(
                    "SELECT 1 FROM edges WHERE from_node=? LIMIT 1",
                    (row["id"],),
                ).fetchone()
                if not has_children:
                    removable.append(row["id"])
            if removable:
                placeholders = ",".join("?" * len(removable))
                conn.execute(
                    f"DELETE FROM edges WHERE from_node IN ({placeholders}) OR to_node IN ({placeholders})",
                    removable * 2,
                )
                conn.execute(
                    f"DELETE FROM nodes WHERE id IN ({placeholders})", removable
                )
                self._v2_delete_nodes(conn, removable)

    def _local_file_index_has_extracted_text(self, row: sqlite3.Row) -> bool:
        metadata = _safe_loads(row["metadata_json"])
        parser = metadata.get("parser") if isinstance(metadata, dict) else {}
        if not isinstance(parser, dict):
            return False
        try:
            return int(parser.get("extracted_chars") or 0) > 0
        except (TypeError, ValueError):
            return False

    @staticmethod
    def _node_matches_workspace(
        conn: sqlite3.Connection,
        node_id: Optional[str],
        workspace_id: Optional[str],
    ) -> bool:
        """Return true only when the projected node has the expected scope."""
        if not node_id:
            return False
        row = conn.execute(
            "SELECT workspace_id FROM nodes_v2 WHERE id=?",
            (node_id,),
        ).fetchone()
        return bool(row is not None and row["workspace_id"] == workspace_id)

    def index_local_folder(
        self,
        path: Path,
        *,
        include_ocr: bool = False,
        watch_enabled: bool = False,
        user_email: Optional[str] = None,
        workspace_id: Optional[str] = None,
        consent: Optional[Dict[str, Any]] = None,
        max_files: int = 5_000,
        source_id_override: Optional[str] = None,
    ) -> Dict[str, Any]:
        """Read approved files from a local folder and connect them to Graph RAG."""
        root = Path(path).expanduser().resolve()
        if not root.exists():
            raise ValueError(f"경로가 존재하지 않습니다: {path}")
        if not root.is_dir():
            raise ValueError(f"폴더가 아닙니다: {path}")

        os_type = _current_os_type()
        drive_id = _drive_id_for_path(root)
        path_fingerprint = _path_fingerprint(root)
        source_id = str(source_id_override or "").strip()
        if not source_id:
            source_id = (
                f"source:{_sha256_text(f'{workspace_id}|{path_fingerprint}')[:24]}"
                if workspace_id
                else f"source:{path_fingerprint}"
            )
        now = _now()
        max_files = max(1, min(int(max_files or 5_000), 50_000))
        consent_payload = {
            "approved_at": now,
            "knowledge_source": True,
            "include_ocr": bool(include_ocr),
            "watch_enabled": bool(watch_enabled),
            "sensitive_files_default_excluded": True,
            **(consent or {}),
            "approved_by": user_email or (consent or {}).get("approved_by"),
            "workspace_id": workspace_id or (consent or {}).get("workspace_id"),
        }
        counts: Counter = Counter()
        seen_relative_paths: set = set()
        indexed_nodes: List[str] = []
        errors: List[Dict[str, str]] = []
        limit_reached = False

        with self._connect() as conn:
            existing_source = conn.execute(
                "SELECT id, consent_json FROM knowledge_sources WHERE root_path=?",
                (str(root),),
            ).fetchone()
            if existing_source is not None:
                existing_consent = _safe_loads(existing_source["consent_json"])
                existing_scope = existing_consent.get("workspace_id") or "personal"
                requested_scope = workspace_id or consent_payload.get("workspace_id") or "personal"
                if existing_scope != requested_scope:
                    raise ValueError(
                        "This folder is already connected to another workspace. "
                        "Disconnect it there before assigning it to a different Brain."
                    )
                if existing_source["id"] != source_id:
                    # Reuse the legacy source identity so a personal source is
                    # reprojected in place instead of duplicated during upgrade.
                    source_id = existing_source["id"]
            conn.execute(
                """
                    INSERT INTO knowledge_sources(
                      id, root_path, os_type, drive_id, label, status, include_ocr,
                      watch_enabled, consent_json, created_at, updated_at, last_scanned_at
                    )
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                    ON CONFLICT(id) DO UPDATE SET
                      root_path=excluded.root_path,
                      os_type=excluded.os_type,
                      drive_id=excluded.drive_id,
                      label=excluded.label,
                      status=excluded.status,
                      include_ocr=excluded.include_ocr,
                      watch_enabled=excluded.watch_enabled,
                      consent_json=excluded.consent_json,
                      updated_at=excluded.updated_at,
                      last_scanned_at=excluded.last_scanned_at
                    """,
                (
                    source_id,
                    str(root),
                    os_type,
                    drive_id,
                    root.name or str(root),
                    "scanning",
                    1 if include_ocr else 0,
                    1 if watch_enabled else 0,
                    _json(consent_payload),
                    now,
                    now,
                    now,
                ),
            )

            for entry in self._iter_local_scan_entries(root, max_files=max_files):
                kind = entry["kind"]
                file_path = entry["path"]
                if kind == "limit_reached":
                    counts["limit_reached"] += 1
                    limit_reached = True
                    break
                if kind in {"excluded_dir", "excluded"}:
                    counts["excluded"] += 1
                    continue
                if kind in {"inaccessible_dir", "inaccessible_file"}:
                    counts["failed"] += 1
                    errors.append(
                        {
                            "path": str(file_path),
                            "error": entry.get("reason", "inaccessible"),
                        }
                    )
                    continue
                if kind != "file":
                    continue

                stat = entry["stat"]
                try:
                    relative_path = file_path.relative_to(root).as_posix()
                except ValueError:
                    relative_path = file_path.name
                seen_relative_paths.add(relative_path)
                modified_at = _safe_iso_from_stat_mtime(stat.st_mtime)
                existing = conn.execute(
                    """
                        SELECT size_bytes, modified_at, sha256, graph_node_id, status, metadata_json
                        FROM local_file_index
                        WHERE source_id=? AND relative_path=?
                        """,
                    (source_id, relative_path),
                ).fetchone()
                decision = self._local_file_decision(file_path, root, stat)
                parser_type = decision["parser_type"]
                if not decision["indexable"]:
                    counts[decision["status"]] += 1
                    if existing and existing["graph_node_id"]:
                        self._delete_local_file_graph(conn, existing["graph_node_id"])
                    self._upsert_local_file_index(
                        conn,
                        source_id=source_id,
                        root=root,
                        file_path=file_path,
                        stat=stat,
                        os_type=os_type,
                        drive_id=drive_id,
                        status=decision["status"],
                        parser_type=parser_type,
                        metadata={
                            "reason": decision["reason"],
                            "category": decision["category"],
                        },
                    )
                    continue

                if (
                    existing
                    and existing["status"] == "indexed"
                    and existing["graph_node_id"]
                    and self._local_file_index_has_extracted_text(existing)
                    and self._node_matches_workspace(
                        conn, existing["graph_node_id"], workspace_id
                    )
                    and existing["size_bytes"] == stat.st_size
                    and existing["modified_at"] == modified_at
                ):
                    counts["skipped_unchanged"] += 1
                    self._upsert_local_file_index(
                        conn,
                        source_id=source_id,
                        root=root,
                        file_path=file_path,
                        stat=stat,
                        os_type=os_type,
                        drive_id=drive_id,
                        status="indexed",
                        parser_type=parser_type,
                        sha256=existing["sha256"],
                        graph_node_id=existing["graph_node_id"],
                        metadata={
                            **_safe_loads(existing["metadata_json"]),
                            "category": decision["category"],
                            "unchanged": True,
                        },
                    )
                    continue

                try:
                    data = file_path.read_bytes()
                    digest = _sha256_bytes(data)
                except Exception as exc:
                    counts["failed"] += 1
                    errors.append({"path": str(file_path), "error": str(exc)})
                    if existing and existing["graph_node_id"]:
                        self._delete_local_file_graph(conn, existing["graph_node_id"])
                    self._upsert_local_file_index(
                        conn,
                        source_id=source_id,
                        root=root,
                        file_path=file_path,
                        stat=stat,
                        os_type=os_type,
                        drive_id=drive_id,
                        status="failed",
                        parser_type=parser_type,
                        error_message=str(exc),
                        metadata={"category": decision["category"]},
                    )
                    continue

                if (
                    existing
                    and existing["sha256"] == digest
                    and existing["graph_node_id"]
                    and self._local_file_index_has_extracted_text(existing)
                    and self._node_matches_workspace(
                        conn, existing["graph_node_id"], workspace_id
                    )
                ):
                    counts["skipped_unchanged"] += 1
                    self._upsert_local_file_index(
                        conn,
                        source_id=source_id,
                        root=root,
                        file_path=file_path,
                        stat=stat,
                        os_type=os_type,
                        drive_id=drive_id,
                        status="indexed",
                        parser_type=parser_type,
                        sha256=digest,
                        graph_node_id=existing["graph_node_id"],
                        metadata={
                            **_safe_loads(existing["metadata_json"]),
                            "category": decision["category"],
                            "sha256_unchanged": True,
                        },
                    )
                    continue

                try:
                    text, parser_meta = self._extract_local_file_text(
                        file_path,
                        decision["category"],
                        include_ocr=include_ocr,
                    )
                    text = _clean_text(text)
                    parser_meta = {**parser_meta, "extracted_chars": len(text)}
                    if not text:
                        counts["skipped_empty_text"] += 1
                        if existing and existing["graph_node_id"]:
                            self._delete_local_file_graph(
                                conn, existing["graph_node_id"]
                            )
                        self._upsert_local_file_index(
                            conn,
                            source_id=source_id,
                            root=root,
                            file_path=file_path,
                            stat=stat,
                            os_type=os_type,
                            drive_id=drive_id,
                            status="skipped_empty_text",
                            parser_type=parser_type,
                            sha256=digest,
                            error_message="텍스트 추출 결과가 비어 있습니다.",
                            metadata={
                                "category": decision["category"],
                                "parser": parser_meta,
                            },
                        )
                        continue
                    graph_node_id = self._upsert_local_file_node(
                        conn,
                        source_id=source_id,
                        root=root,
                        file_path=file_path,
                        stat=stat,
                        os_type=os_type,
                        drive_id=drive_id,
                        sha256=digest,
                        category=decision["category"],
                        parser_type=parser_type,
                        text=text,
                        parser_meta=parser_meta,
                        user_email=user_email,
                        workspace_id=workspace_id,
                    )
                    self._upsert_local_file_index(
                        conn,
                        source_id=source_id,
                        root=root,
                        file_path=file_path,
                        stat=stat,
                        os_type=os_type,
                        drive_id=drive_id,
                        status="indexed",
                        parser_type=parser_type,
                        sha256=digest,
                        graph_node_id=graph_node_id,
                        metadata={
                            "category": decision["category"],
                            "parser": parser_meta,
                        },
                    )
                    counts["indexed"] += 1
                    indexed_nodes.append(graph_node_id)
                except Exception as exc:
                    counts["failed"] += 1
                    errors.append({"path": str(file_path), "error": str(exc)})
                    if existing and existing["graph_node_id"]:
                        self._delete_local_file_graph(conn, existing["graph_node_id"])
                    self._upsert_local_file_index(
                        conn,
                        source_id=source_id,
                        root=root,
                        file_path=file_path,
                        stat=stat,
                        os_type=os_type,
                        drive_id=drive_id,
                        status="failed",
                        parser_type=parser_type,
                        sha256=digest,
                        error_message=str(exc),
                        metadata={"category": decision["category"]},
                    )

            if not limit_reached:
                existing_rows = {
                    row["relative_path"]: row["graph_node_id"]
                    for row in conn.execute(
                        "SELECT relative_path, graph_node_id FROM local_file_index WHERE source_id=?",
                        (source_id,),
                    )
                }
                deleted_paths = set(existing_rows) - seen_relative_paths
                for relative_path in deleted_paths:
                    self._delete_local_file_graph(
                        conn, existing_rows.get(relative_path)
                    )
                    conn.execute(
                        """
                            UPDATE local_file_index
                            SET status='deleted', deleted=1, last_scanned_at=?, error_message=NULL, graph_node_id=NULL
                            WHERE source_id=? AND relative_path=?
                            """,
                        (_now(), source_id, relative_path),
                    )
                counts["deleted"] = len(deleted_paths)
            conn.execute(
                """
                    UPDATE knowledge_sources
                    SET status='active', updated_at=?, last_scanned_at=?
                    WHERE id=?
                    """,
                (_now(), _now(), source_id),
            )

        return {
            "status": "ok",
            "source": {
                "id": source_id,
                "root_path": str(root),
                "os_type": os_type,
                "drive_id": drive_id,
                "include_ocr": bool(include_ocr),
                "watch_enabled": bool(watch_enabled),
            },
            "counts": dict(counts),
            "indexed_nodes": indexed_nodes[:100],
            "errors": errors[:50],
            "notice": "Lattice AI는 사용자가 선택한 폴더만 AI 지식으로 변환합니다.",
        }
