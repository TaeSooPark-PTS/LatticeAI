from __future__ import annotations

from typing import TYPE_CHECKING

from ..quiet import quiet

# ruff: noqa: F403,F405
from ._kg_common import *  # noqa: F403,F401

# The cross-mixin surface (`_connect`, `_upsert_node`, …) is declared in
# `_kg_contract.KnowledgeGraphCore`. It is a typing-only base: at runtime this
# is `object`, so the MRO of `KnowledgeGraphStore` is unchanged.
if TYPE_CHECKING:
    from ._kg_contract import KnowledgeGraphCore as _Core
else:
    _Core = object



class KnowledgeGraphDocumentsMixin(_Core):
    def find_documents_by_uri_prefix(self, prefix: str) -> List[Dict[str, Any]]:
        """Content nodes whose ``metadata.source_uri`` starts with ``prefix``.

        Powers removable seeded corpora (e.g. the ``demo://`` First Value Loop
        documents): callers stamp a URI scheme at ingest time and can later
        enumerate exactly those nodes without guessing ids.
        """
        prefix = str(prefix or "").strip()
        if not prefix:
            return []
        nt, _ = self._read_tables()
        documents: List[Dict[str, Any]] = []
        with self._connect() as conn:
            rows = conn.execute(
                f"""
                SELECT id, type, title, metadata_json, created_at, updated_at
                FROM {nt}
                WHERE json_extract(metadata_json, '$.source_uri') LIKE ? || '%'
                  AND type NOT IN ('Source', 'Chunk')
                ORDER BY created_at ASC, id ASC
                """,
                (prefix,),
            ).fetchall()
            for row in rows:
                meta = _safe_loads(row["metadata_json"]) or {}
                documents.append(
                    {
                        "id": row["id"],
                        "type": row["type"],
                        "title": row["title"],
                        "source_uri": meta.get("source_uri"),
                        "workspace_id": meta.get("workspace_id"),
                        "created_at": row["created_at"],
                        "updated_at": row["updated_at"],
                    }
                )
        return documents

    def delete_document_tree(self, node_id: str) -> Dict[str, Any]:
        """Delete an ingested content node plus everything it owns.

        Removes the node, its retrieval chunks (``chunks`` table + ``Chunk``
        nodes), auto-extracted Task/Decision nodes whose ``source_node`` is
        this document, every touching edge, the vector rows, and — when it
        becomes orphaned — the linked ``Source`` node. Shared Concept nodes
        are intentionally left alone (they may be cited by other content).
        Mirrored into the v2 projection via ``_v2_delete_nodes``.
        """
        node_id = str(node_id or "").strip()
        if not node_id:
            return {"status": "skipped", "removed_nodes": 0}
        with self._connect() as conn:
            if not conn.execute(
                "SELECT 1 FROM nodes WHERE id=?", (node_id,)
            ).fetchone():
                return {"status": "not_found", "node_id": node_id, "removed_nodes": 0}

            remove_ids = {node_id}
            # Owned children: chunks + auto-extracted semantic nodes that
            # explicitly point back at this document via metadata.source_node.
            for row in conn.execute(
                """
                SELECT id FROM nodes
                WHERE json_extract(metadata_json, '$.source_node') = ?
                """,
                (node_id,),
            ):
                remove_ids.add(row["id"])
            # The Source node this content was indexed from (if it ends up
            # referenced by nothing else after the delete, drop it too).
            source_ids = [
                row["to_node"]
                for row in conn.execute(
                    "SELECT to_node FROM edges WHERE from_node=? AND type IN ('indexed_from', 'INDEXED_FROM')",
                    (node_id,),
                )
            ]

            ids = list(remove_ids)
            ph = ",".join("?" * len(ids))
            conn.execute(f"DELETE FROM chunks WHERE source_node IN ({ph})", ids)
            conn.execute(
                f"DELETE FROM edges WHERE from_node IN ({ph}) OR to_node IN ({ph})",
                ids * 2,
            )
            conn.execute(
                f"DELETE FROM vector_embeddings WHERE item_id IN ({ph}) OR source_node IN ({ph})",
                ids * 2,
            )
            conn.execute(f"DELETE FROM nodes WHERE id IN ({ph})", ids)
            self._v2_delete_nodes(conn, ids)

            removed_sources = 0
            for source_id in source_ids:
                still_linked = conn.execute(
                    "SELECT 1 FROM edges WHERE from_node=? OR to_node=? LIMIT 1",
                    (source_id, source_id),
                ).fetchone()
                if still_linked:
                    continue
                conn.execute("DELETE FROM vector_embeddings WHERE item_id=?", (source_id,))
                conn.execute("DELETE FROM nodes WHERE id=?", (source_id,))
                self._v2_delete_nodes(conn, [source_id])
                removed_sources += 1

        return {
            "status": "ok",
            "node_id": node_id,
            "removed_nodes": len(remove_ids) + removed_sources,
        }

    def _ingest_structure_nodes(
        self,
        conn: sqlite3.Connection,
        file_id: str,
        filename: str,
        structure: Dict[str, Any],
        *,
        owner: Optional[str] = None,
        workspace_id: Optional[str] = None,
    ) -> None:
        for slide in structure.get("slides") or []:
            index = slide.get("index")
            slide_id = f"slide:{_sha256_text(f'{file_id}:slide:{index}')[:24]}"
            title = f"{filename} slide {index}"
            summary = "\n".join(slide.get("texts") or [])[:800]
            self._upsert_node(
                conn,
                slide_id,
                "Slide",
                title,
                summary=summary,
                metadata={**slide, "workspace_id": workspace_id},
                owner=owner,
                workspace_id=workspace_id,
            )
            self._upsert_edge(conn, file_id, slide_id, "has_slide")
            for text in slide.get("texts") or []:
                for topic in _topic_candidates(text, limit=4):
                    topic_key = f"{workspace_id}|{topic}" if workspace_id else topic
                    topic_id = f"topic:{_sha256_text(topic_key)[:24]}" if workspace_id else f"topic:{_slug(topic)}"
                    self._upsert_node(
                        conn,
                        topic_id,
                        "Topic",
                        topic,
                        metadata={"auto_extracted": True, "workspace_id": workspace_id},
                        owner=owner,
                        workspace_id=workspace_id,
                    )
                    self._upsert_edge(conn, slide_id, topic_id, "discusses", weight=0.6)

        for page in structure.get("pages") or []:
            index = page.get("index")
            page_id = f"page:{_sha256_text(f'{file_id}:page:{index}')[:24]}"
            title = f"{filename} page {index}"
            self._upsert_node(
                conn,
                page_id,
                "Page",
                title,
                summary=page.get("preview") or "",
                metadata={**page, "workspace_id": workspace_id},
                owner=owner,
                workspace_id=workspace_id,
            )
            self._upsert_edge(conn, file_id, page_id, "has_page")
            for topic in _topic_candidates(page.get("preview") or "", limit=4):
                topic_key = f"{workspace_id}|{topic}" if workspace_id else topic
                topic_id = f"topic:{_sha256_text(topic_key)[:24]}" if workspace_id else f"topic:{_slug(topic)}"
                self._upsert_node(
                    conn,
                    topic_id,
                    "Topic",
                    topic,
                    metadata={"auto_extracted": True, "workspace_id": workspace_id},
                    owner=owner,
                    workspace_id=workspace_id,
                )
                self._upsert_edge(conn, page_id, topic_id, "discusses", weight=0.6)

        for sheet in structure.get("sheets") or []:
            sheet_title = sheet.get("title")
            sheet_id = f"sheet:{_sha256_text(f'{file_id}:sheet:{sheet_title}')[:24]}"
            self._upsert_node(
                conn,
                sheet_id,
                "Sheet",
                f"{filename} / {sheet_title}",
                metadata={**sheet, "workspace_id": workspace_id},
                owner=owner,
                workspace_id=workspace_id,
            )
            self._upsert_edge(conn, file_id, sheet_id, "has_sheet")

        for image in structure.get("images") or []:
            image_key = image.get("sha256") or _sha256_text(
                json.dumps(image, ensure_ascii=False, sort_keys=True)
            )
            scoped_image_key = f"{workspace_id}|{image_key}" if workspace_id else str(image_key)
            image_id = f"image:{_sha256_text(scoped_image_key)[:24]}" if workspace_id else f"image:{str(image_key)[:24]}"
            title_parts = [filename, "image"]
            if image.get("page"):
                title_parts.append(f"page {image.get('page')}")
            if image.get("name"):
                title_parts.append(str(image.get("name")).split("/")[-1])
            self._upsert_node(
                conn,
                image_id,
                "Image",
                " / ".join(title_parts),
                metadata={**image, "workspace_id": workspace_id},
                owner=owner,
                workspace_id=workspace_id,
            )
            self._upsert_edge(conn, file_id, image_id, "contains_image")

    def _document_structure(self, path: Path, ext: str) -> Dict[str, Any]:
        try:
            if ext == ".pptx":
                return self._pptx_structure(path)
            if ext == ".pdf":
                return self._pdf_structure(path)
            if ext == ".docx":
                return self._docx_structure(path)
            if ext == ".xlsx":
                return self._xlsx_structure(path)
        except Exception as exc:
            return {"error": str(exc)}
        return {}

    def _pptx_structure(self, path: Path) -> Dict[str, Any]:
        result: Dict[str, Any] = {"slides": [], "images": []}
        try:
            from PIL import Image
            from pptx import Presentation

            prs = Presentation(str(path))
            for slide_index, slide in enumerate(prs.slides, start=1):
                slide_info: Dict[str, Any] = {
                    "index": slide_index,
                    "shapes": [],
                    "texts": [],
                }
                for shape_index, shape in enumerate(slide.shapes, start=1):
                    shape_info = {
                        "index": shape_index,
                        "name": getattr(shape, "name", ""),
                        "shape_type": str(getattr(shape, "shape_type", "")),
                        "bbox": {
                            "left": int(getattr(shape, "left", 0) or 0),
                            "top": int(getattr(shape, "top", 0) or 0),
                            "width": int(getattr(shape, "width", 0) or 0),
                            "height": int(getattr(shape, "height", 0) or 0),
                        },
                    }
                    if getattr(shape, "has_text_frame", False):
                        text = shape.text_frame.text.strip()
                        if text:
                            shape_info["text"] = text[:1000]
                            slide_info["texts"].append(text)
                    slide_info["shapes"].append(shape_info)
                result["slides"].append(slide_info)
            with zipfile.ZipFile(path) as zf:
                for name in zf.namelist():
                    if not name.startswith("ppt/media/"):
                        continue
                    data = zf.read(name)
                    image_info: Dict[str, Any] = {
                        "name": name,
                        "bytes": len(data),
                        "sha256": _sha256_bytes(data),
                    }
                    try:
                        from io import BytesIO

                        with Image.open(BytesIO(data)) as img:
                            image_info.update(
                                {
                                    "width": img.width,
                                    "height": img.height,
                                    "format": img.format,
                                }
                            )
                    except Exception:
                        quiet()
                    result["images"].append(image_info)
        except Exception as exc:
            result["error"] = str(exc)
        return result

    def _pdf_structure(self, path: Path) -> Dict[str, Any]:
        result: Dict[str, Any] = {"pages": [], "images": []}
        try:
            import pdfplumber

            with pdfplumber.open(str(path)) as pdf:
                metadata = dict(pdf.metadata or {})
                result["metadata"] = {str(k): str(v) for k, v in metadata.items()}
                for page_index, page in enumerate(pdf.pages, start=1):
                    text = page.extract_text() or ""
                    page_info = {
                        "index": page_index,
                        "width": float(page.width or 0),
                        "height": float(page.height or 0),
                        "chars": len(text),
                        "preview": _clean_text(text)[:500],
                        "image_count": len(page.images or []),
                    }
                    result["pages"].append(page_info)
                    for image_index, image in enumerate(page.images or [], start=1):
                        result["images"].append(
                            {
                                "page": page_index,
                                "index": image_index,
                                "name": image.get("name"),
                                "width": image.get("width"),
                                "height": image.get("height"),
                                "bbox": {
                                    "x0": image.get("x0"),
                                    "top": image.get("top"),
                                    "x1": image.get("x1"),
                                    "bottom": image.get("bottom"),
                                },
                            }
                        )
        except Exception as exc:
            result["error"] = str(exc)
        return result

    def _docx_structure(self, path: Path) -> Dict[str, Any]:
        from docx import Document

        doc = Document(str(path))
        headings = []
        paragraphs = 0
        for p in doc.paragraphs:
            text = p.text.strip()
            if not text:
                continue
            paragraphs += 1
            style = getattr(p.style, "name", "")
            if style.lower().startswith("heading"):
                headings.append({"style": style, "text": text[:240]})
        return {
            "paragraphs": paragraphs,
            "headings": headings[:80],
            "tables": len(doc.tables),
        }

    def _xlsx_structure(self, path: Path) -> Dict[str, Any]:
        from openpyxl import load_workbook

        wb = load_workbook(str(path), read_only=True, data_only=True)
        sheets = []
        for ws in wb.worksheets:
            sheets.append(
                {"title": ws.title, "max_row": ws.max_row, "max_column": ws.max_column}
            )
        return {"sheets": sheets}
