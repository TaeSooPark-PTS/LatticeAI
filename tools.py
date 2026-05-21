"""
Safe local tools for Lattice AI agent mode.

All filesystem operations are confined to LATTICEAI_AGENT_ROOT, defaulting to
./agent_workspace. Command execution runs without a shell and from inside that
workspace.
"""

import base64
import os
import platform
import re
import shlex
import socket
import subprocess
import tempfile
import json
from html.parser import HTMLParser
from pathlib import Path
from typing import Any, Dict, List, Optional

_PLATFORM = platform.system()  # "Darwin" | "Windows" | "Linux"

from p_reinforce import BRAIN_DIR, STRUCTURE

# ── Computer Use ──────────────────────────────────────────────────────────────
_CU_AVAILABLE = False
_pyautogui = None

def _init_computer_use():
    global _CU_AVAILABLE, _pyautogui
    try:
        import pyautogui as _pag
        _pag.FAILSAFE = True   # 마우스를 좌상단 코너로 이동하면 중단
        _pag.PAUSE = 0.25
        _pyautogui = _pag
        _CU_AVAILABLE = True
    except Exception:
        pass

_init_computer_use()


def computer_screenshot() -> Dict[str, Any]:
    """현재 화면 전체를 캡처하여 base64 PNG로 반환합니다."""
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_file:
        tmp = tmp_file.name
    try:
        if _PLATFORM == "Darwin":
            r = subprocess.run(
                ["screencapture", "-x", "-t", "png", tmp],
                capture_output=True, timeout=10, check=False,
            )
            if r.returncode != 0:
                raise ToolError(f"screencapture 실패: {r.stderr.decode()}")
        elif _CU_AVAILABLE:
            # Windows / Linux: use pyautogui screenshot
            screenshot = _pyautogui.screenshot()
            screenshot.save(tmp)
        else:
            raise ToolError("스크린샷 불가: macOS 전용 screencapture 또는 pyautogui 필요")
        with open(tmp, "rb") as f:
            b64 = base64.b64encode(f.read()).decode()
        size = os.path.getsize(tmp)
        w, h = (_pyautogui.size() if _CU_AVAILABLE else (0, 0))
        return {
            "screenshot_b64": b64,
            "format": "png",
            "bytes": size,
            "screen_width": int(w),
            "screen_height": int(h),
        }
    finally:
        try:
            if os.path.exists(tmp):
                os.unlink(tmp)
        except OSError:
            pass


def computer_open_app(app: str = "Google Chrome") -> Dict[str, Any]:
    """앱을 실행하거나 앞으로 가져옵니다 (macOS/Windows/Linux)."""
    app = str(app or "Google Chrome").strip()
    if not app:
        raise ToolError("앱 이름이 필요합니다.")
    if _PLATFORM == "Darwin":
        cmd = ["open", "-a", app]
    elif _PLATFORM == "Windows":
        cmd = ["cmd", "/c", "start", "", app]
    else:
        cmd = ["xdg-open", app]
    r = subprocess.run(cmd, capture_output=True, timeout=10, check=False)
    if r.returncode != 0:
        err = r.stderr.decode("utf-8", errors="replace").strip()
        raise ToolError(f"앱 열기 실패: {err or app}")
    return {"action": "open_app", "app": app}


def computer_open_url(url: str, app: str = "Google Chrome") -> Dict[str, Any]:
    """URL을 브라우저로 엽니다 (macOS/Windows/Linux)."""
    url = str(url or "").strip()
    app = str(app or "").strip()
    if not url:
        raise ToolError("URL이 필요합니다.")
    if "://" not in url and not url.startswith(("localhost", "127.0.0.1")):
        url = "https://" + url
    if _PLATFORM == "Darwin":
        cmd = ["open", "-a", app, url] if app else ["open", url]
    elif _PLATFORM == "Windows":
        cmd = ["cmd", "/c", "start", "", url]
    else:
        cmd = ["xdg-open", url]
    r = subprocess.run(cmd, capture_output=True, timeout=10, check=False)
    if r.returncode != 0:
        err = r.stderr.decode("utf-8", errors="replace").strip()
        raise ToolError(f"URL 열기 실패: {err or url}")
    return {"action": "open_url", "app": app or "default", "url": url}


def computer_click(x: int, y: int, button: str = "left", double: bool = False) -> Dict[str, Any]:
    """화면 좌표 (x, y)를 클릭합니다."""
    if not _CU_AVAILABLE:
        raise ToolError("pyautogui를 사용할 수 없습니다.")
    x, y = int(x), int(y)
    if double:
        _pyautogui.doubleClick(x, y)
    elif button == "right":
        _pyautogui.rightClick(x, y)
    elif button == "middle":
        _pyautogui.middleClick(x, y)
    else:
        _pyautogui.click(x, y)
    return {"action": "click", "x": x, "y": y, "button": button, "double": double}


def computer_type(text: str, interval: float = 0.04) -> Dict[str, Any]:
    """현재 포커스된 위치에 텍스트를 입력합니다."""
    if not _CU_AVAILABLE:
        raise ToolError("pyautogui를 사용할 수 없습니다.")
    _pyautogui.write(str(text), interval=float(interval))
    return {"action": "type", "text": (text[:60] + "...") if len(text) > 60 else text, "chars": len(text)}


def computer_key(key: str) -> Dict[str, Any]:
    """키보드 키를 누릅니다. 예: 'return', 'escape', 'command+c', 'tab'"""
    if not _CU_AVAILABLE:
        raise ToolError("pyautogui를 사용할 수 없습니다.")
    key = str(key)
    if "+" in key:
        _pyautogui.hotkey(*key.split("+"))
    else:
        _pyautogui.press(key)
    return {"action": "key", "key": key}


def computer_scroll(x: int, y: int, direction: str = "down", clicks: int = 3) -> Dict[str, Any]:
    """화면 좌표에서 스크롤합니다."""
    if not _CU_AVAILABLE:
        raise ToolError("pyautogui를 사용할 수 없습니다.")
    x, y, clicks = int(x), int(y), int(clicks)
    _pyautogui.moveTo(x, y)
    amount = -clicks if direction == "down" else clicks
    _pyautogui.scroll(amount)
    return {"action": "scroll", "x": x, "y": y, "direction": direction, "clicks": clicks}


def computer_move(x: int, y: int) -> Dict[str, Any]:
    """마우스를 좌표로 이동합니다 (클릭 없음)."""
    if not _CU_AVAILABLE:
        raise ToolError("pyautogui를 사용할 수 없습니다.")
    _pyautogui.moveTo(int(x), int(y), duration=0.2)
    return {"action": "move", "x": int(x), "y": int(y)}


def computer_drag(x1: int, y1: int, x2: int, y2: int) -> Dict[str, Any]:
    """(x1,y1)에서 (x2,y2)로 드래그합니다."""
    if not _CU_AVAILABLE:
        raise ToolError("pyautogui를 사용할 수 없습니다.")
    _pyautogui.moveTo(int(x1), int(y1))
    _pyautogui.dragTo(int(x2), int(y2), duration=0.35, button="left")
    return {"action": "drag", "from": [int(x1), int(y1)], "to": [int(x2), int(y2)]}


def computer_status() -> Dict[str, Any]:
    """Computer Use 기능 사용 가능 여부를 확인합니다."""
    if not _CU_AVAILABLE:
        return {"available": False, "reason": "pyautogui not installed"}
    w, h = _pyautogui.size()
    return {
        "available": True,
        "screen_size": {"width": int(w), "height": int(h)},
        "failsafe": _pyautogui.FAILSAFE,
        "note": "macOS Accessibility 권한이 필요합니다 (시스템 설정 > 개인 정보 보호 > 손쉬운 사용)",
    }


AGENT_ROOT = Path(os.getenv("LATTICEAI_AGENT_ROOT") or "agent_workspace").resolve()
MAX_FILE_BYTES = 512_000
MAX_COMMAND_SECONDS = 30
MAX_BUILD_SECONDS = 180
MAX_DEPLOY_SECONDS = 300
MAX_COMMAND_OUTPUT = 12_000

BLOCKED_COMMANDS = {
    "rm",
    "rmdir",
    "sudo",
    "su",
    "chmod",
    "chown",
    "curl",
    "wget",
    "ssh",
    "scp",
    "rsync",
    "dd",
    "mkfs",
    "diskutil",
    "launchctl",
}

ALLOWED_COMMANDS = {
    "pwd",
    "ls",
    "find",
    "cat",
    "sed",
    "head",
    "tail",
    "wc",
    "rg",
    "python",
    "python3",
    "node",
    "npm",
    "npx",
    "git",
}

BUILD_SCRIPT_NAMES = {"build", "compile", "typecheck", "test"}
DEPLOY_SCRIPT_NAMES = {
    "deploy",
    "preview",
    "release",
    "package",
    "dist",
    "make",
    "build:installer",
    "build:pkg",
    "build:exe",
    "package:mac",
    "package:win",
}

ALLOWED_GIT_SUBCOMMANDS = {"status", "diff", "log", "show"}

TEXT_EXTENSIONS = {
    ".css",
    ".csv",
    ".html",
    ".js",
    ".json",
    ".jsx",
    ".md",
    ".py",
    ".ts",
    ".tsx",
    ".txt",
    ".xml",
    ".yaml",
    ".yml",
}

DOCUMENT_OUTPUT_DIR = "generated_documents"
PRESENTATION_OUTPUT_DIR = "generated_presentations"
SPREADSHEET_OUTPUT_DIR = "generated_spreadsheets"


class ToolError(ValueError):
    pass


def ensure_agent_root() -> Path:
    AGENT_ROOT.mkdir(parents=True, exist_ok=True)
    return AGENT_ROOT


def _resolve_path(path: str = "") -> Path:
    ensure_agent_root()
    if not path:
        return AGENT_ROOT
    candidate = (AGENT_ROOT / path).resolve()
    if candidate != AGENT_ROOT and AGENT_ROOT not in candidate.parents:
        raise ToolError("Path escapes the agent workspace.")
    return candidate


def _relative(path: Path) -> str:
    return str(path.relative_to(AGENT_ROOT))


def list_dir(path: str = ".") -> Dict[str, Any]:
    target = _resolve_path(path)
    if not target.exists():
        raise ToolError("Directory does not exist.")
    if not target.is_dir():
        raise ToolError("Path is not a directory.")

    items: List[Dict[str, Any]] = []
    for child in sorted(target.iterdir(), key=lambda p: (not p.is_dir(), p.name.lower())):
        items.append(
            {
                "name": child.name,
                "path": _relative(child),
                "type": "directory" if child.is_dir() else "file",
                "size": child.stat().st_size if child.is_file() else None,
            }
        )
    return {"root": str(AGENT_ROOT), "path": _relative(target) if target != AGENT_ROOT else ".", "items": items}


def workspace_tree(path: str = ".", max_depth: int = 3) -> Dict[str, Any]:
    target = _resolve_path(path)
    if not target.exists() or not target.is_dir():
        raise ToolError("Path is not a directory.")

    max_depth = max(1, min(int(max_depth), 8))
    entries: List[Dict[str, Any]] = []

    def walk(current: Path, depth: int) -> None:
        if depth > max_depth:
            return
        for child in sorted(current.iterdir(), key=lambda p: (not p.is_dir(), p.name.lower())):
            rel = _relative(child)
            entries.append(
                {
                    "path": rel,
                    "type": "directory" if child.is_dir() else "file",
                    "size": child.stat().st_size if child.is_file() else None,
                    "depth": depth,
                }
            )
            if child.is_dir():
                walk(child, depth + 1)

    walk(target, 1)
    return {"root": str(AGENT_ROOT), "path": _relative(target) if target != AGENT_ROOT else ".", "entries": entries}


def read_file(path: str, offset: int = 0, limit: int = 0, line_numbers: bool = True) -> Dict[str, Any]:
    """Read a file from the agent workspace.

    Returns content as plain text. When line_numbers is True (default), also
    returns a numbered view (`numbered`) plus `total_lines` so the agent can
    cite file:line locations precisely.

    offset is 0-indexed (the first line is offset=0). limit=0 reads to the end.
    """
    target = _resolve_path(path)
    if not target.exists():
        raise ToolError("File does not exist.")
    if not target.is_file():
        raise ToolError("Path is not a file.")
    size = target.stat().st_size
    if size > MAX_FILE_BYTES:
        raise ToolError(f"File is too large to read ({size} bytes).")
    text = target.read_text(encoding="utf-8")
    all_lines = text.splitlines()
    total_lines = len(all_lines)

    offset = max(0, int(offset or 0))
    limit = max(0, int(limit or 0))
    end = total_lines if limit == 0 else min(total_lines, offset + limit)
    sliced = all_lines[offset:end]
    sliced_text = "\n".join(sliced)
    if offset == 0 and limit == 0 and text.endswith("\n"):
        sliced_text += "\n"

    result: Dict[str, Any] = {
        "path": _relative(target),
        "content": sliced_text,
        "total_lines": total_lines,
        "start_line": offset + 1,
        "end_line": end,
    }
    if line_numbers:
        width = max(4, len(str(end or total_lines)))
        result["numbered"] = "\n".join(
            f"{(offset + i + 1):>{width}}\t{line}" for i, line in enumerate(sliced)
        )
    return result


def write_file(path: str, content: str) -> Dict[str, Any]:
    target = _resolve_path(path)
    if len(content.encode("utf-8")) > MAX_FILE_BYTES:
        raise ToolError("Content is too large to write.")
    target.parent.mkdir(parents=True, exist_ok=True)
    target.write_text(content, encoding="utf-8")
    return {"path": _relative(target), "bytes": target.stat().st_size}


def edit_file(path: str, old_string: str, new_string: str, replace_all: bool = False) -> Dict[str, Any]:
    """Precise diff-style edit: replace `old_string` with `new_string` in `path`.

    Fails when `old_string` is missing or appears more than once (unless
    replace_all=True). This forces the caller to read the file first and pass
    enough surrounding context to uniquely identify the edit site — the same
    discipline Claude Code uses for safe edits.
    """
    if old_string == new_string:
        raise ToolError("old_string and new_string are identical; nothing to change.")
    target = _resolve_path(path)
    if not target.exists() or not target.is_file():
        raise ToolError("File does not exist.")
    if target.stat().st_size > MAX_FILE_BYTES:
        raise ToolError("File is too large to edit.")

    original = target.read_text(encoding="utf-8")
    occurrences = original.count(old_string)
    if occurrences == 0:
        raise ToolError("old_string not found in file. Read the file first and copy the exact bytes (including whitespace).")
    if occurrences > 1 and not replace_all:
        raise ToolError(f"old_string is ambiguous: appears {occurrences} times. Add more context to make it unique, or pass replace_all=true.")

    updated = original.replace(old_string, new_string) if replace_all else original.replace(old_string, new_string, 1)
    if len(updated.encode("utf-8")) > MAX_FILE_BYTES:
        raise ToolError("Resulting file would exceed the workspace size limit.")
    target.write_text(updated, encoding="utf-8")

    edited_line = original[: original.find(old_string)].count("\n") + 1
    return {
        "path": _relative(target),
        "replacements": occurrences if replace_all else 1,
        "bytes": target.stat().st_size,
        "first_edit_line": edited_line,
    }


_GREP_BINARY_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".pdf", ".zip", ".tar",
                    ".gz", ".bz2", ".xz", ".7z", ".mp3", ".mp4", ".mov", ".wav",
                    ".woff", ".woff2", ".ttf", ".eot", ".ico", ".db", ".sqlite",
                    ".pyc", ".pyo", ".o", ".so", ".dylib", ".dll", ".exe", ".bin"}
_GREP_BINARY_DIRS = {"node_modules", ".git", ".venv", "venv", "__pycache__",
                    ".pytest_cache", "dist", "build", ".next", ".cache"}


def grep(
    pattern: str,
    path: str = ".",
    glob: Optional[str] = None,
    max_results: int = 50,
    case_insensitive: bool = False,
    context_lines: int = 0,
) -> Dict[str, Any]:
    """Regex search across the agent workspace.

    Unlike `search_files` (single line, 9 extensions, substring only), this
    walks all text files, supports regex, returns line numbers, and can
    optionally include surrounding context lines. Skips obvious binary
    files/directories.
    """
    if not pattern:
        raise ToolError("Pattern is required.")
    try:
        flags = re.IGNORECASE if case_insensitive else 0
        regex = re.compile(pattern, flags)
    except re.error as exc:
        raise ToolError(f"Invalid regex: {exc}") from exc

    target = _resolve_path(path)
    if not target.exists() or not target.is_dir():
        raise ToolError("Path is not a directory.")

    max_results = max(1, min(int(max_results), 500))
    context_lines = max(0, min(int(context_lines), 8))
    matches: List[Dict[str, Any]] = []
    files_scanned = 0
    files_with_matches = 0

    iterator = target.rglob(glob) if glob else target.rglob("*")
    for file_path in iterator:
        if len(matches) >= max_results:
            break
        if not file_path.is_file():
            continue
        if file_path.suffix.lower() in _GREP_BINARY_EXTS:
            continue
        if any(part in _GREP_BINARY_DIRS for part in file_path.parts):
            continue
        if file_path.stat().st_size > MAX_FILE_BYTES:
            continue
        try:
            lines = file_path.read_text(encoding="utf-8").splitlines()
        except (UnicodeDecodeError, OSError):
            continue

        files_scanned += 1
        file_had_match = False
        for index, line in enumerate(lines, start=1):
            if len(matches) >= max_results:
                break
            if not regex.search(line):
                continue
            file_had_match = True
            entry: Dict[str, Any] = {
                "path": _relative(file_path),
                "line": index,
                "match": line[:400],
            }
            if context_lines:
                lo = max(0, index - 1 - context_lines)
                hi = min(len(lines), index + context_lines)
                entry["context"] = [
                    {"line": lo + i + 1, "text": lines[lo + i][:200]}
                    for i in range(hi - lo)
                ]
            matches.append(entry)
        if file_had_match:
            files_with_matches += 1

    return {
        "pattern": pattern,
        "matches": matches,
        "files_scanned": files_scanned,
        "files_with_matches": files_with_matches,
        "truncated": len(matches) >= max_results,
    }


_TODO_REL_PATH = ".lattice/todos.json"
_TODO_ALLOWED_STATUS = {"pending", "in_progress", "completed"}


def _todo_file() -> Path:
    ensure_agent_root()
    target = AGENT_ROOT / _TODO_REL_PATH
    target.parent.mkdir(parents=True, exist_ok=True)
    return target


def todo_read() -> Dict[str, Any]:
    """Read the agent's persistent TODO list (per-workspace)."""
    target = _todo_file()
    if not target.exists():
        return {"todos": [], "path": _TODO_REL_PATH}
    try:
        todos = json.loads(target.read_text(encoding="utf-8"))
    except (json.JSONDecodeError, OSError):
        todos = []
    if not isinstance(todos, list):
        todos = []
    return {"todos": todos, "path": _TODO_REL_PATH}


def todo_write(todos: List[Dict[str, Any]]) -> Dict[str, Any]:
    """Replace the agent's TODO list. Each todo: {id, content, status}.

    Status must be one of: pending, in_progress, completed.
    At most one todo should be in_progress at any time — the agent enforces
    this convention; the tool only warns if violated.
    """
    if not isinstance(todos, list):
        raise ToolError("todos must be a list.")
    if len(todos) > 50:
        raise ToolError("Too many todos (max 50). Split into smaller batches.")

    cleaned: List[Dict[str, Any]] = []
    in_progress_count = 0
    for idx, raw in enumerate(todos, start=1):
        if not isinstance(raw, dict):
            raise ToolError(f"Todo #{idx} is not an object.")
        content = str(raw.get("content") or "").strip()
        if not content:
            raise ToolError(f"Todo #{idx} is missing 'content'.")
        status = str(raw.get("status") or "pending").strip().lower()
        if status not in _TODO_ALLOWED_STATUS:
            raise ToolError(f"Todo #{idx} has invalid status '{status}'. Use one of {sorted(_TODO_ALLOWED_STATUS)}.")
        if status == "in_progress":
            in_progress_count += 1
        cleaned.append({
            "id": str(raw.get("id") or idx),
            "content": content[:240],
            "status": status,
        })

    target = _todo_file()
    target.write_text(json.dumps(cleaned, ensure_ascii=False, indent=2), encoding="utf-8")
    return {
        "todos": cleaned,
        "path": _TODO_REL_PATH,
        "warning": "More than one todo is in_progress; keep only one active at a time." if in_progress_count > 1 else None,
    }


def search_files(query: str, path: str = ".", max_results: int = 20) -> Dict[str, Any]:
    if not query:
        raise ToolError("Query is required.")
    target = _resolve_path(path)
    if not target.exists() or not target.is_dir():
        raise ToolError("Path is not a directory.")

    max_results = max(1, min(int(max_results), 100))
    matches: List[Dict[str, Any]] = []
    query_lower = query.lower()

    for file_path in target.rglob("*"):
        if len(matches) >= max_results:
            break
        if not file_path.is_file() or file_path.stat().st_size > MAX_FILE_BYTES:
            continue
        if file_path.suffix.lower() not in TEXT_EXTENSIONS:
            continue
        try:
            lines = file_path.read_text(encoding="utf-8").splitlines()
        except UnicodeDecodeError:
            continue
        for index, line in enumerate(lines, start=1):
            if query_lower in line.lower():
                matches.append({"path": _relative(file_path), "line": index, "preview": line[:240]})
                break

    return {"query": query, "matches": matches}


class _HTMLInspector(HTMLParser):
    def __init__(self):
        super().__init__()
        self.title = ""
        self._in_title = False
        self.links: List[str] = []
        self.scripts: List[str] = []
        self.stylesheets: List[str] = []
        self.images: List[str] = []
        self.forms = 0
        self.headings: List[Dict[str, str]] = []

    def handle_starttag(self, tag: str, attrs: List[tuple]) -> None:
        attr = dict(attrs)
        if tag == "title":
            self._in_title = True
        elif tag == "a" and attr.get("href"):
            self.links.append(attr["href"])
        elif tag == "script" and attr.get("src"):
            self.scripts.append(attr["src"])
        elif tag == "link" and attr.get("rel") and "stylesheet" in " ".join(attr.get("rel", [])):
            if attr.get("href"):
                self.stylesheets.append(attr["href"])
        elif tag == "img" and attr.get("src"):
            self.images.append(attr["src"])
        elif tag == "form":
            self.forms += 1
        elif tag in {"h1", "h2", "h3"}:
            self.headings.append({"level": tag, "text": ""})

    def handle_endtag(self, tag: str) -> None:
        if tag == "title":
            self._in_title = False

    def handle_data(self, data: str) -> None:
        text = data.strip()
        if not text:
            return
        if self._in_title:
            self.title += text
        elif self.headings and not self.headings[-1]["text"]:
            self.headings[-1]["text"] = text[:120]


def inspect_html(path: str) -> Dict[str, Any]:
    target = _resolve_path(path)
    if not target.exists() or not target.is_file():
        raise ToolError("HTML file does not exist.")
    if target.suffix.lower() not in {".html", ".htm"}:
        raise ToolError("Path is not an HTML file.")
    if target.stat().st_size > MAX_FILE_BYTES:
        raise ToolError("HTML file is too large to inspect.")

    parser = _HTMLInspector()
    parser.feed(target.read_text(encoding="utf-8"))
    return {
        "path": _relative(target),
        "title": parser.title,
        "links": parser.links[:50],
        "scripts": parser.scripts[:50],
        "stylesheets": parser.stylesheets[:50],
        "images": parser.images[:50],
        "forms": parser.forms,
        "headings": [h for h in parser.headings if h["text"]][:30],
    }


def preview_url(path: str = "index.html") -> Dict[str, Any]:
    target = _resolve_path(path)
    if not target.exists() or not target.is_file():
        raise ToolError("Preview file does not exist.")
    rel = _relative(target)
    return {
        "path": rel,
        "local_url": f"http://127.0.0.1:4825/agent-files/{rel}",
        "note": "Use the server host or /web Telegram link host instead of 127.0.0.1 from a phone.",
    }


def create_web_project(path: str, framework: str = "react", template: str = "vite") -> Dict[str, Any]:
    framework = str(framework or "").strip().lower()
    template = str(template or "").strip().lower()
    if framework != "react" or template != "vite":
        raise ToolError("Only React + Vite template is currently supported.")
    if not path:
        raise ToolError("Project path is required.")

    root = _resolve_path(path)
    root.mkdir(parents=True, exist_ok=True)

    files = {
        "package.json": json.dumps(
            {
                "name": Path(path).name.replace(" ", "-").lower() or "vite-react-app",
                "private": True,
                "version": "0.0.0",
                "type": "module",
                "scripts": {
                    "dev": "vite",
                    "build": "vite build",
                    "preview": "vite preview",
                },
                "dependencies": {
                    "react": "^18.3.1",
                    "react-dom": "^18.3.1",
                },
                "devDependencies": {
                    "@vitejs/plugin-react": "^4.3.1",
                    "vite": "^5.4.0",
                },
            },
            ensure_ascii=False,
            indent=2,
        ) + "\n",
        "index.html": """<!doctype html>
<html lang="en">
  <head>
    <meta charset="UTF-8" />
    <meta name="viewport" content="width=device-width, initial-scale=1.0" />
    <title>Vite React App</title>
  </head>
  <body>
    <div id="root"></div>
    <script type="module" src="/src/main.jsx"></script>
  </body>
</html>
""",
        "vite.config.js": """import { defineConfig } from 'vite'
import react from '@vitejs/plugin-react'

export default defineConfig({
  plugins: [react()],
})
""",
        "README.md": """# Vite React App

## Run

```bash
npm install
npm run dev
```

## Build

```bash
npm run build
npm run preview
```

## Lattice AI Notes

- Inspect `package.json` and existing config files before adding new libraries.
- If you add Tailwind CSS, framer-motion, TypeScript, or other tooling, add the required config files too.
- Do not report the app as complete until `npm run build` succeeds.
""",
        "src/main.jsx": """import React from 'react'
import ReactDOM from 'react-dom/client'
import App from './App.jsx'
import './index.css'

ReactDOM.createRoot(document.getElementById('root')).render(
  <React.StrictMode>
    <App />
  </React.StrictMode>,
)
""",
        "src/App.jsx": """import { useState } from 'react'

export default function App() {
  const [count, setCount] = useState(0)
  return (
    <main style={{ maxWidth: 760, margin: '48px auto', padding: '0 20px', fontFamily: 'system-ui, sans-serif' }}>
      <h1>Vite + React</h1>
      <p>Starter generated by Lattice AI agent.</p>
      <p style={{ color: '#555', lineHeight: 1.6 }}>
        Inspect the current setup before adding new UI libraries, then verify
        changes with <code>npm run build</code>.
      </p>
      <button onClick={() => setCount((c) => c + 1)}>count is {count}</button>
    </main>
  )
}
""",
        "src/index.css": """* { box-sizing: border-box; }
body { margin: 0; background: #f6f7fb; color: #111; }
button { padding: 10px 14px; border-radius: 10px; border: 1px solid #d6d6d6; background: #fff; cursor: pointer; }
""",
    }

    created: List[str] = []
    total_bytes = 0
    for rel_path, content in files.items():
        target = (root / rel_path).resolve()
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_text(content, encoding="utf-8")
        created.append(_relative(target))
        total_bytes += target.stat().st_size

    return {
        "path": _relative(root),
        "framework": framework,
        "template": template,
        "created_files": created,
        "file_count": len(created),
        "bytes": total_bytes,
    }


def _safe_filename(name: str, suffix: str) -> str:
    base = Path(name or f"artifact{suffix}").name
    if not base.lower().endswith(suffix):
        base += suffix
    safe = "".join(ch if ch.isalnum() or ch in ("-", "_", ".", " ") else "_" for ch in base).strip()
    return safe or f"artifact{suffix}"


def _body_to_str(body) -> str:
    if isinstance(body, list):
        return "\n\n".join(str(item) for item in body)
    return str(body or "")


def create_docx(title: str, body, filename: str = "document.docx") -> Dict[str, Any]:
    try:
        from docx import Document
    except Exception as exc:
        raise ToolError("python-docx is not installed. Run `pip install -r requirements.txt`.") from exc

    output_dir = _resolve_path(DOCUMENT_OUTPUT_DIR)
    output_dir.mkdir(parents=True, exist_ok=True)
    target = output_dir / _safe_filename(filename, ".docx")

    document = Document()
    if title:
        document.add_heading(str(title), level=1)
    for block in _body_to_str(body).split("\n\n"):
        text = block.strip()
        if text:
            document.add_paragraph(text)
    document.save(target)
    return {"path": _relative(target), "bytes": target.stat().st_size}


def create_xlsx(rows: List[List[Any]], filename: str = "spreadsheet.xlsx", sheet_name: str = "Sheet1") -> Dict[str, Any]:
    try:
        from openpyxl import Workbook
    except Exception as exc:
        raise ToolError("openpyxl is not installed. Run `pip install -r requirements.txt`.") from exc

    if not isinstance(rows, list) or not all(isinstance(row, list) for row in rows):
        raise ToolError("Rows must be a list of lists.")

    output_dir = _resolve_path(SPREADSHEET_OUTPUT_DIR)
    output_dir.mkdir(parents=True, exist_ok=True)
    target = output_dir / _safe_filename(filename, ".xlsx")

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = (sheet_name or "Sheet1")[:31]
    for row in rows:
        sheet.append(row)
    workbook.save(target)
    return {"path": _relative(target), "rows": len(rows), "bytes": target.stat().st_size}


def create_pptx(title: str, slides: List[Dict[str, Any]], filename: str = "presentation.pptx") -> Dict[str, Any]:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise ToolError("python-pptx is not installed. Run `pip install -r requirements.txt`.") from exc

    output_dir = _resolve_path(PRESENTATION_OUTPUT_DIR)
    output_dir.mkdir(parents=True, exist_ok=True)
    target = output_dir / _safe_filename(filename, ".pptx")

    presentation = Presentation()
    first_layout = presentation.slide_layouts[0]
    first = presentation.slides.add_slide(first_layout)
    first.shapes.title.text = title or "Presentation"
    first.placeholders[1].text = ""

    content_layout = presentation.slide_layouts[1]
    for slide_data in slides or []:
        slide = presentation.slides.add_slide(content_layout)
        slide.shapes.title.text = str(slide_data.get("title") or "Slide")
        body = slide.placeholders[1].text_frame
        body.clear()
        bullets = slide_data.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        for index, bullet in enumerate(bullets):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = str(bullet)
            paragraph.level = 0

    presentation.save(target)
    return {"path": _relative(target), "slides": len(presentation.slides), "bytes": target.stat().st_size}


PDF_OUTPUT_DIR = "generated_pdfs"
LOCAL_MAX_FILE_BYTES = 2_000_000  # 2 MB cap for local reads


# CJK-capable fonts (Korean + Chinese + Japanese)
_CJK_FONT_CANDIDATES = [
    "/System/Library/Fonts/AppleSDGothicNeo.ttc",   # Korean (macOS)
    "/System/Library/Fonts/STHeiti Light.ttc",       # Chinese (macOS)
    "/System/Library/Fonts/PingFang.ttc",            # Chinese (macOS)
    "/Library/Fonts/NanumGothic.ttf",               # Korean
    "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",
]

_SUPPORTED_READ_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx", ".txt", ".md", ".csv"}
DOCUMENT_MAX_READ_BYTES = 10_000_000  # 10 MB


def create_pdf(title: str, body, filename: str = "document.pdf") -> Dict[str, Any]:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import mm
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
    except Exception as exc:
        raise ToolError("reportlab is not installed. Run `pip install reportlab`.") from exc

    output_dir = _resolve_path(PDF_OUTPUT_DIR)
    output_dir.mkdir(parents=True, exist_ok=True)
    target = output_dir / _safe_filename(filename, ".pdf")

    # CJK 폰트 등록
    font_name = "Helvetica"
    for font_path in _CJK_FONT_CANDIDATES:
        if Path(font_path).exists():
            try:
                pdfmetrics.registerFont(TTFont("KoreanFont", font_path))
                font_name = "KoreanFont"
            except Exception:
                pass
            break

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("Title", fontName=font_name, fontSize=18, spaceAfter=8, leading=24)
    body_style  = ParagraphStyle("Body",  fontName=font_name, fontSize=11, spaceAfter=6, leading=16)

    story = []
    if title:
        story.append(Paragraph(str(title), title_style))
        story.append(Spacer(1, 4 * mm))

    for block in _body_to_str(body).split("\n\n"):
        text = block.strip()
        if text:
            safe_text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe_text, body_style))
            story.append(Spacer(1, 2 * mm))

    doc = SimpleDocTemplate(str(target), pagesize=A4,
                             leftMargin=20*mm, rightMargin=20*mm,
                             topMargin=20*mm, bottomMargin=20*mm)
    doc.build(story)
    return {"path": _relative(target), "bytes": target.stat().st_size}


def local_list(path: str) -> Dict[str, Any]:
    """List any directory on the local filesystem (requires user approval via UI)."""
    target = Path(path).expanduser().resolve()
    if not target.exists():
        raise ToolError(f"경로가 존재하지 않습니다: {path}")
    if not target.is_dir():
        raise ToolError(f"폴더가 아닙니다: {path}")
    items = []
    try:
        for child in sorted(target.iterdir(), key=lambda p: (not p.is_dir(), p.name.lower())):
            stat = child.stat()
            items.append({
                "name": child.name,
                "path": str(child),
                "type": "directory" if child.is_dir() else "file",
                "size": stat.st_size if child.is_file() else None,
            })
    except PermissionError as exc:
        raise ToolError(f"접근 권한 없음: {exc}") from exc
    return {"path": str(target), "items": items}


def local_read(path: str) -> Dict[str, Any]:
    """Read any file on the local filesystem (requires user approval via UI)."""
    target = Path(path).expanduser().resolve()
    if not target.exists():
        raise ToolError(f"파일이 존재하지 않습니다: {path}")
    if not target.is_file():
        raise ToolError(f"파일이 아닙니다: {path}")
    size = target.stat().st_size
    if size > LOCAL_MAX_FILE_BYTES:
        raise ToolError(f"파일이 너무 큽니다 ({size:,} bytes). 최대 {LOCAL_MAX_FILE_BYTES:,} bytes.")
    try:
        content = target.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:
        raise ToolError(f"파일 읽기 실패: {exc}") from exc
    return {"path": str(target), "size": size, "content": content}


def local_write(path: str, content: str) -> Dict[str, Any]:
    """Write content to any path on the local filesystem (requires user approval via UI)."""
    target = Path(path).expanduser().resolve()
    if len(content.encode("utf-8")) > LOCAL_MAX_FILE_BYTES:
        raise ToolError("내용이 너무 큽니다.")
    try:
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_text(content, encoding="utf-8")
    except PermissionError as exc:
        raise ToolError(f"쓰기 권한 없음: {exc}") from exc
    return {"path": str(target), "bytes": target.stat().st_size}


def read_document(path: str) -> Dict[str, Any]:
    """Extract text from PDF, DOCX, XLSX, PPTX, TXT, MD, CSV files."""
    target = Path(path).expanduser().resolve()
    if not target.exists():
        raise ToolError(f"파일이 없습니다: {path}")
    if not target.is_file():
        raise ToolError(f"파일이 아닙니다: {path}")
    if target.stat().st_size > DOCUMENT_MAX_READ_BYTES:
        raise ToolError(f"파일이 너무 큽니다 ({target.stat().st_size:,} bytes).")

    ext = target.suffix.lower()
    if ext not in _SUPPORTED_READ_EXTENSIONS:
        raise ToolError(f"지원하지 않는 형식입니다: {ext}. 지원: {', '.join(_SUPPORTED_READ_EXTENSIONS)}")

    text = ""
    meta: Dict[str, Any] = {"path": str(target), "ext": ext}

    if ext == ".pdf":
        try:
            import pdfplumber
            with pdfplumber.open(str(target)) as pdf:
                meta["pages"] = len(pdf.pages)
                text = "\n\n".join(
                    (p.extract_text() or "") for p in pdf.pages
                ).strip()
        except Exception as exc:
            raise ToolError(f"PDF 읽기 실패: {exc}") from exc

    elif ext == ".docx":
        try:
            from docx import Document
            doc = Document(str(target))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            text = "\n\n".join(paragraphs)
            meta["paragraphs"] = len(paragraphs)
        except Exception as exc:
            raise ToolError(f"DOCX 읽기 실패: {exc}") from exc

    elif ext == ".xlsx":
        try:
            from openpyxl import load_workbook
            wb = load_workbook(str(target), data_only=True)
            rows_all = []
            for ws in wb.worksheets:
                rows_all.append(f"[Sheet: {ws.title}]")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    rows_all.append("\t".join(cells))
            text = "\n".join(rows_all)
            meta["sheets"] = len(wb.worksheets)
        except Exception as exc:
            raise ToolError(f"XLSX 읽기 실패: {exc}") from exc

    elif ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(target))
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                parts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
                slides_text.append(f"[Slide {i}]\n" + "\n".join(parts))
            text = "\n\n".join(slides_text)
            meta["slides"] = len(prs.slides)
        except Exception as exc:
            raise ToolError(f"PPTX 읽기 실패: {exc}") from exc

    elif ext in {".txt", ".md", ".csv"}:
        try:
            text = target.read_text(encoding="utf-8", errors="replace")
        except Exception as exc:
            raise ToolError(f"파일 읽기 실패: {exc}") from exc

    meta["chars"] = len(text)
    meta["preview"] = text[:500]
    meta["content"] = text[:50_000]   # 50K char cap for context
    return meta


def desktop_bridge_status() -> Dict[str, Any]:
    return {
        "status": "requires_desktop_bridge",
        "available_in_codex": True,
        "note": "Chrome and Mac UI control require the Codex desktop Computer Use/Chrome bridge, not a headless FastAPI worker.",
    }


def _safe_brain_folder(folder: str) -> str:
    if folder not in STRUCTURE:
        raise ToolError(f"Unknown knowledge folder: {folder}")
    return folder


def knowledge_save(content: str, folder: str = "00_Raw", title: Optional[str] = None) -> Dict[str, Any]:
    folder = _safe_brain_folder(folder)
    if not content:
        raise ToolError("Knowledge content is required.")
    if len(content.encode("utf-8")) > MAX_FILE_BYTES:
        raise ToolError("Knowledge content is too large.")

    target_dir = BRAIN_DIR / folder
    target_dir.mkdir(parents=True, exist_ok=True)
    safe_title = title or content.strip().splitlines()[0][:60] or "note"
    safe_title = "".join(ch if ch.isalnum() or ch in (" ", "-", "_") else "" for ch in safe_title).strip()
    safe_title = "_".join(safe_title.split()) or "note"
    filename = f"{safe_title}.md"
    target = target_dir / filename
    counter = 2
    while target.exists():
        target = target_dir / f"{safe_title}_{counter}.md"
        counter += 1
    target.write_text(content, encoding="utf-8")
    return {"folder": folder, "filename": target.name, "path": str(target)}


def knowledge_search(query: str, max_results: int = 5) -> Dict[str, Any]:
    if not query:
        raise ToolError("Query is required.")
    max_results = max(1, min(int(max_results), 20))
    query_lower = query.lower()
    results: List[Dict[str, Any]] = []

    for file_path in BRAIN_DIR.rglob("*.md"):
        if len(results) >= max_results:
            break
        try:
            content = file_path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            continue
        if query_lower in content.lower() or query_lower in file_path.name.lower():
            results.append(
                {
                    "path": str(file_path),
                    "relative_path": str(file_path.relative_to(BRAIN_DIR)),
                    "preview": content[:500],
                }
            )

    return {"query": query, "results": results}


def knowledge_tree() -> Dict[str, Any]:
    entries: List[Dict[str, Any]] = []
    for folder in STRUCTURE:
        root = BRAIN_DIR / folder
        root.mkdir(parents=True, exist_ok=True)
        for file_path in sorted(root.rglob("*.md")):
            entries.append(
                {
                    "folder": folder,
                    "relative_path": str(file_path.relative_to(BRAIN_DIR)),
                    "size": file_path.stat().st_size,
                }
            )
    return {"root": str(BRAIN_DIR), "entries": entries}


def obsidian_save(content: str, folder: str = "00_Raw", title: Optional[str] = None) -> Dict[str, Any]:
    result = knowledge_save(content, folder, title)
    result["vault_root"] = str(BRAIN_DIR)
    result["obsidian_uri_hint"] = f"obsidian://open?path={result['path']}"
    return result


def obsidian_search(query: str, max_results: int = 5) -> Dict[str, Any]:
    result = knowledge_search(query, max_results)
    result["vault_root"] = str(BRAIN_DIR)
    return result


def obsidian_tree() -> Dict[str, Any]:
    return knowledge_tree()


def _run_network_command(parts: List[str], timeout: int = 5) -> str:
    try:
        completed = subprocess.run(parts, capture_output=True, text=True, timeout=timeout, check=False)
        if completed.returncode != 0:
            return ""
        return completed.stdout.strip()
    except Exception:
        return ""


def network_status() -> Dict[str, Any]:
    """현재 Mac의 내부 IP, 외부 IP, 주요 네트워크 정보를 반환합니다."""
    local_ips: Dict[str, str] = {}
    for interface in ["en0", "en1", "bridge100"]:
        value = _run_network_command(["ipconfig", "getifaddr", interface])
        if value:
            local_ips[interface] = value

    ifconfig_text = _run_network_command(["ifconfig"])
    current_interface = ""
    for line in ifconfig_text.splitlines():
        if line and not line.startswith(("\t", " ")):
            current_interface = line.split(":", 1)[0]
            continue
        match = re.search(r"\binet\s+(\d+\.\d+\.\d+\.\d+)\b", line)
        if match and current_interface and match.group(1) != "127.0.0.1":
            local_ips.setdefault(current_interface, match.group(1))

    hostname = socket.gethostname()
    guessed_ip = ""
    try:
        with socket.socket(socket.AF_INET, socket.SOCK_DGRAM) as sock:
            sock.connect(("8.8.8.8", 80))
            guessed_ip = sock.getsockname()[0]
    except Exception:
        pass
    if guessed_ip and guessed_ip not in local_ips.values():
        local_ips["default_route"] = guessed_ip

    public_ip = _run_network_command(["curl", "-sS", "--max-time", "3", "https://api.ipify.org"])
    wifi_info = _run_network_command(["networksetup", "-getinfo", "Wi-Fi"])

    primary_local_ip = local_ips.get("en0") or local_ips.get("en1") or guessed_ip or ""
    return {
        "hostname": hostname,
        "local_ip": primary_local_ip,
        "local_ips": local_ips,
        "public_ip": public_ip,
        "wifi_info": wifi_info,
        "ifconfig_available": bool(ifconfig_text),
        "note": "local_ip은 같은 네트워크 안에서 보이는 내부 IP이고, public_ip는 인터넷에서 보이는 외부 IP입니다.",
    }


_BLOCKED_FIND_FLAGS = {"-exec", "-execdir", "-delete", "-ok", "-okdir"}

def run_command(command: str, cwd: Optional[str] = None) -> Dict[str, Any]:
    ensure_agent_root()
    parts = shlex.split(command)
    if not parts:
        raise ToolError("Command is empty.")

    executable = Path(parts[0]).name
    if executable in BLOCKED_COMMANDS or executable not in ALLOWED_COMMANDS:
        raise ToolError(f"Command is not allowed: {executable}")
    if executable == "git":
        raise ToolError("Use the read-only git_status, git_diff, git_log, or git_show tools.")
    if any(token in command for token in ["|", "&&", "||", ";", ">", "<", "$(", "`"]):
        raise ToolError("Shell operators are not allowed.")
    if executable == "find":
        blocked = [f for f in parts[1:] if f in _BLOCKED_FIND_FLAGS]
        if blocked:
            raise ToolError(f"find flags are not allowed: {', '.join(blocked)}")
    abs_args = [a for a in parts[1:] if a.startswith("/") and a not in ("/dev/null",)]
    if abs_args:
        raise ToolError(f"Absolute paths in command arguments are not allowed: {abs_args[0]}")

    workdir = _resolve_path(cwd or ".")
    if not workdir.exists() or not workdir.is_dir():
        raise ToolError("Working directory does not exist.")

    try:
        completed = subprocess.run(
            parts,
            cwd=workdir,
            capture_output=True,
            text=True,
            timeout=MAX_COMMAND_SECONDS,
            check=False,
        )
    except subprocess.TimeoutExpired:
        raise ToolError(f"Command timed out after {MAX_COMMAND_SECONDS} seconds.")

    stdout = completed.stdout[-MAX_COMMAND_OUTPUT:]
    stderr = completed.stderr[-MAX_COMMAND_OUTPUT:]
    return {
        "command": command,
        "cwd": _relative(workdir) if workdir != AGENT_ROOT else ".",
        "returncode": completed.returncode,
        "stdout": stdout,
        "stderr": stderr,
    }


def _load_package_scripts(workdir: Path) -> Dict[str, str]:
    package_json = workdir / "package.json"
    if not package_json.exists():
        return {}
    try:
        import json
        data = json.loads(package_json.read_text(encoding="utf-8"))
    except Exception as exc:
        raise ToolError(f"Could not parse package.json: {exc}") from exc
    scripts = data.get("scripts") or {}
    if not isinstance(scripts, dict):
        return {}
    return {str(key): str(value) for key, value in scripts.items()}


def _run_script(script: str, cwd: Optional[str], allowed: set[str], timeout: int) -> Dict[str, Any]:
    ensure_agent_root()
    if script not in allowed:
        raise ToolError(f"Script is not allowed here: {script}")
    workdir = _resolve_path(cwd or ".")
    if not workdir.exists() or not workdir.is_dir():
        raise ToolError("Working directory does not exist.")

    scripts = _load_package_scripts(workdir)
    if script not in scripts:
        raise ToolError(f"package.json does not define a '{script}' script.")

    try:
        completed = subprocess.run(
            ["npm", "run", script],
            cwd=workdir,
            capture_output=True,
            text=True,
            timeout=timeout,
            check=False,
        )
    except subprocess.TimeoutExpired:
        raise ToolError(f"npm run {script} timed out after {timeout} seconds.")

    return {
        "command": f"npm run {script}",
        "cwd": _relative(workdir) if workdir != AGENT_ROOT else ".",
        "script_body": scripts[script],
        "returncode": completed.returncode,
        "stdout": completed.stdout[-MAX_COMMAND_OUTPUT:],
        "stderr": completed.stderr[-MAX_COMMAND_OUTPUT:],
    }


def build_project(cwd: Optional[str] = None, script: str = "build") -> Dict[str, Any]:
    return _run_script(script, cwd, BUILD_SCRIPT_NAMES, MAX_BUILD_SECONDS)


def deploy_project(cwd: Optional[str] = None, script: str = "deploy") -> Dict[str, Any]:
    return _run_script(script, cwd, DEPLOY_SCRIPT_NAMES, MAX_DEPLOY_SECONDS)


def _run_git(args: List[str], cwd: Optional[str] = None) -> Dict[str, Any]:
    if not args:
        raise ToolError("Git subcommand is required.")
    subcommand = args[0]
    if subcommand not in ALLOWED_GIT_SUBCOMMANDS:
        raise ToolError(f"Git subcommand is not allowed: {subcommand}")
    if any(arg.startswith(("git@", "http://", "https://", "ssh://")) for arg in args):
        raise ToolError("Remote git targets are not allowed.")

    workdir = _resolve_path(cwd or ".")
    if not workdir.exists() or not workdir.is_dir():
        raise ToolError("Working directory does not exist.")

    try:
        completed = subprocess.run(
            ["git", *args],
            cwd=workdir,
            capture_output=True,
            text=True,
            timeout=MAX_COMMAND_SECONDS,
            check=False,
        )
    except subprocess.TimeoutExpired:
        raise ToolError(f"Git command timed out after {MAX_COMMAND_SECONDS} seconds.")

    return {
        "command": "git " + " ".join(args),
        "cwd": _relative(workdir) if workdir != AGENT_ROOT else ".",
        "returncode": completed.returncode,
        "stdout": completed.stdout[-MAX_COMMAND_OUTPUT:],
        "stderr": completed.stderr[-MAX_COMMAND_OUTPUT:],
    }


def git_status(cwd: Optional[str] = None) -> Dict[str, Any]:
    return _run_git(["status", "--short"], cwd)


def git_diff(path: Optional[str] = None, cwd: Optional[str] = None) -> Dict[str, Any]:
    args = ["diff", "--"]
    if path:
        target = _resolve_path(path)
        args.append(_relative(target))
    return _run_git(args, cwd)


def git_log(max_count: int = 5, cwd: Optional[str] = None) -> Dict[str, Any]:
    max_count = max(1, min(int(max_count), 20))
    return _run_git(["log", f"--max-count={max_count}", "--oneline", "--decorate"], cwd)


def git_show(revision: str = "HEAD", cwd: Optional[str] = None) -> Dict[str, Any]:
    if revision.startswith("-") or any(token in revision for token in ["..", ":", "/", "\\"]):
        raise ToolError("Revision is not allowed.")
    return _run_git(["show", "--stat", "--oneline", "--decorate", revision], cwd)


def execute_tool(action: str, args: Dict[str, Any]) -> Dict[str, Any]:
    if action == "list_dir":
        return list_dir(args.get("path", "."))
    if action == "workspace_tree":
        return workspace_tree(args.get("path", "."), args.get("max_depth", 3))
    if action == "read_file":
        return read_file(
            args["path"],
            offset=args.get("offset", 0),
            limit=args.get("limit", 0),
            line_numbers=args.get("line_numbers", True),
        )
    if action == "write_file":
        return write_file(args["path"], args.get("content", ""))
    if action == "edit_file":
        return edit_file(
            args["path"],
            args["old_string"],
            args["new_string"],
            replace_all=bool(args.get("replace_all", False)),
        )
    if action == "grep":
        return grep(
            args["pattern"],
            path=args.get("path", "."),
            glob=args.get("glob"),
            max_results=args.get("max_results", 50),
            case_insensitive=bool(args.get("case_insensitive", False)),
            context_lines=args.get("context_lines", 0),
        )
    if action == "search_files":
        return search_files(args["query"], args.get("path", "."), args.get("max_results", 20))
    if action == "todo_read":
        return todo_read()
    if action == "todo_write":
        return todo_write(args.get("todos") or [])
    if action == "inspect_html":
        return inspect_html(args["path"])
    if action == "preview_url":
        return preview_url(args.get("path", "index.html"))
    if action == "create_docx":
        return create_docx(args.get("title", ""), args.get("body", ""), args.get("filename", "document.docx"))
    if action == "create_xlsx":
        rows = args.get("rows", [])
        if isinstance(rows, str):
            rows = json.loads(rows)
        return create_xlsx(rows, args.get("filename", "spreadsheet.xlsx"), args.get("sheet_name", "Sheet1"))
    if action == "create_pptx":
        slides = args.get("slides", [])
        if isinstance(slides, str):
            slides = json.loads(slides)
        return create_pptx(args.get("title", ""), slides, args.get("filename", "presentation.pptx"))
    if action == "create_pdf":
        return create_pdf(args.get("title", ""), args.get("body", ""), args.get("filename", "document.pdf"))
    if action == "create_web_project":
        return create_web_project(args.get("path", ""), args.get("framework", "react"), args.get("template", "vite"))
    if action == "local_list":
        return local_list(args["path"])
    if action == "local_read":
        return local_read(args["path"])
    if action == "local_write":
        return local_write(args["path"], args.get("content", ""))
    if action == "read_document":
        return read_document(args["path"])
    if action == "network_status":
        return network_status()
    if action == "computer_screenshot":
        return computer_screenshot()
    if action == "computer_open_app":
        return computer_open_app(args.get("app", "Google Chrome"))
    if action == "computer_open_url":
        return computer_open_url(args["url"], args.get("app", "Google Chrome"))
    if action == "computer_click":
        return computer_click(args.get("x", 0), args.get("y", 0), args.get("button", "left"), args.get("double", False))
    if action == "computer_type":
        return computer_type(args["text"], args.get("interval", 0.04))
    if action == "computer_key":
        return computer_key(args["key"])
    if action == "computer_scroll":
        return computer_scroll(args.get("x", 0), args.get("y", 0), args.get("direction", "down"), args.get("clicks", 3))
    if action == "computer_move":
        return computer_move(args.get("x", 0), args.get("y", 0))
    if action == "computer_drag":
        return computer_drag(args.get("x1", 0), args.get("y1", 0), args.get("x2", 0), args.get("y2", 0))
    if action == "computer_status":
        return computer_status()
    if action in {"chrome_status", "computer_use_status"}:
        return desktop_bridge_status()
    if action == "knowledge_save":
        return knowledge_save(args["content"], args.get("folder", "00_Raw"), args.get("title"))
    if action == "knowledge_search":
        return knowledge_search(args["query"], args.get("max_results", 5))
    if action == "knowledge_tree":
        return knowledge_tree()
    if action == "obsidian_save":
        return obsidian_save(args["content"], args.get("folder", "00_Raw"), args.get("title"))
    if action == "obsidian_search":
        return obsidian_search(args["query"], args.get("max_results", 5))
    if action == "obsidian_tree":
        return obsidian_tree()
    if action == "git_status":
        return git_status(args.get("cwd"))
    if action == "git_diff":
        return git_diff(args.get("path"), args.get("cwd"))
    if action == "git_log":
        return git_log(args.get("max_count", 5), args.get("cwd"))
    if action == "git_show":
        return git_show(args.get("revision", "HEAD"), args.get("cwd"))
    if action == "run_command":
        return run_command(args["command"], args.get("cwd"))
    if action == "build_project":
        return build_project(args.get("cwd"), args.get("script", "build"))
    if action == "deploy_project":
        return deploy_project(args.get("cwd"), args.get("script", "deploy"))
    raise ToolError(f"Unknown action: {action}")
