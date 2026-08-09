"""wp13 coverage — ``latticeai.tools.documents``.

Two halves. The creators (docx/xlsx/pptx/pdf) sanitize the caller's filename
into a fixed output directory inside the workspace and then hand the content
to an optional third-party writer; the reader (``read_document``) dispatches
on extension and turns every parser failure into a ``ToolError`` rather than a
traceback.

docx/xlsx/pptx are declared dependencies, so those tests use the real
libraries and read the artifacts back. ``reportlab`` is *not* a declared
dependency — ``create_pdf``'s own error message says to install it separately
— so the PDF path is driven against injected fake modules. That keeps the
font-selection and escaping logic under test on a machine that has no
reportlab at all, which is the machine CI runs on.
"""

from __future__ import annotations

import sys
import zipfile
from pathlib import Path
from types import ModuleType
from typing import Any, Dict, List

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

import latticeai.tools as tools
from latticeai.tools import ToolError
from latticeai.tools import documents as documents_module
from latticeai.tools.documents import (
    _body_to_str,
    create_docx,
    create_pdf,
    create_pptx,
    create_xlsx,
    read_document,
)


@pytest.fixture()
def workspace(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> Path:
    root = tmp_path / "agent_workspace"
    root.mkdir()
    monkeypatch.setattr(tools, "AGENT_ROOT", root)
    tools.ensure_agent_root()
    return root


# ── fake reportlab ───────────────────────────────────────────────────────────


class _FakeStyle:
    def __init__(self, name: str, **kwargs: Any) -> None:
        self.name = name
        self.kwargs = kwargs


class _FakeParagraph:
    def __init__(self, text: str, style: _FakeStyle) -> None:
        self.text = text
        self.style = style


class _FakeSpacer:
    def __init__(self, width: float, height: float) -> None:
        self.size = (width, height)


class _Recorder:
    def __init__(self) -> None:
        self.registered: List[Any] = []
        self.ttfonts: List[Any] = []
        self.doc_kwargs: Dict[str, Any] = {}
        self.doc_path = ""
        self.story: List[Any] = []
        self.register_raises = False

    @property
    def paragraphs(self) -> List[_FakeParagraph]:
        return [item for item in self.story if isinstance(item, _FakeParagraph)]


def _install_fake_reportlab(monkeypatch: pytest.MonkeyPatch) -> _Recorder:
    """Replace reportlab with modules that record what ``create_pdf`` builds."""
    recorder = _Recorder()

    def register_font(font: Any) -> None:
        if recorder.register_raises:
            raise ValueError("cannot parse font file")
        recorder.registered.append(font)

    def ttfont(name: str, path: str) -> Any:
        recorder.ttfonts.append((name, path))
        return ("ttfont", name, path)

    class _FakeDocTemplate:
        def __init__(self, path: str, **kwargs: Any) -> None:
            recorder.doc_path = path
            recorder.doc_kwargs = kwargs

        def build(self, story: List[Any]) -> None:
            recorder.story = list(story)
            Path(recorder.doc_path).write_bytes(b"%PDF-1.4\n% fake document\n")

    modules: Dict[str, ModuleType] = {}
    for name in (
        "reportlab",
        "reportlab.lib",
        "reportlab.lib.pagesizes",
        "reportlab.lib.styles",
        "reportlab.lib.units",
        "reportlab.pdfbase",
        "reportlab.pdfbase.pdfmetrics",
        "reportlab.pdfbase.ttfonts",
        "reportlab.platypus",
    ):
        modules[name] = ModuleType(name)

    modules["reportlab.lib.pagesizes"].A4 = (595.276, 841.890)
    modules["reportlab.lib.styles"].ParagraphStyle = _FakeStyle
    modules["reportlab.lib.units"].mm = 2.834645669291339
    modules["reportlab.pdfbase.pdfmetrics"].registerFont = register_font
    modules["reportlab.pdfbase.ttfonts"].TTFont = ttfont
    modules["reportlab.platypus"].Paragraph = _FakeParagraph
    modules["reportlab.platypus"].SimpleDocTemplate = _FakeDocTemplate
    modules["reportlab.platypus"].Spacer = _FakeSpacer
    modules["reportlab.lib"].pagesizes = modules["reportlab.lib.pagesizes"]
    modules["reportlab.lib"].styles = modules["reportlab.lib.styles"]
    modules["reportlab.lib"].units = modules["reportlab.lib.units"]
    modules["reportlab.pdfbase"].pdfmetrics = modules["reportlab.pdfbase.pdfmetrics"]
    modules["reportlab.pdfbase"].ttfonts = modules["reportlab.pdfbase.ttfonts"]
    modules["reportlab"].lib = modules["reportlab.lib"]
    modules["reportlab"].pdfbase = modules["reportlab.pdfbase"]
    modules["reportlab"].platypus = modules["reportlab.platypus"]

    for name, module in modules.items():
        monkeypatch.setitem(sys.modules, name, module)
    return recorder


# ── _body_to_str ─────────────────────────────────────────────────────────────


@pytest.mark.parametrize(
    "body,expected",
    [
        (["first", "second"], "first\n\nsecond"),
        ([1, None], "1\n\nNone"),
        ([], ""),
        ("plain", "plain"),
        (None, ""),
        (0, ""),
    ],
)
def test_body_to_str_normalises_every_accepted_shape(body: Any, expected: str) -> None:
    assert _body_to_str(body) == expected


# ── missing optional dependencies ────────────────────────────────────────────


@pytest.mark.parametrize(
    "module_names,call,message",
    [
        (["docx"], lambda: create_docx("t", "b"), "python-docx is not installed"),
        (["openpyxl"], lambda: create_xlsx([["a"]]), "openpyxl is not installed"),
        (["pptx"], lambda: create_pptx("t", []), "python-pptx is not installed"),
        (
            [
                "reportlab",
                "reportlab.lib",
                "reportlab.lib.pagesizes",
                "reportlab.lib.styles",
                "reportlab.lib.units",
                "reportlab.pdfbase",
                "reportlab.pdfbase.pdfmetrics",
                "reportlab.pdfbase.ttfonts",
                "reportlab.platypus",
            ],
            lambda: create_pdf("t", "b"),
            "reportlab is not installed",
        ),
    ],
)
def test_a_missing_writer_library_is_a_readable_tool_error(
    workspace: Path, monkeypatch: pytest.MonkeyPatch, module_names, call, message: str
) -> None:
    for name in module_names:
        monkeypatch.setitem(sys.modules, name, None)

    with pytest.raises(ToolError, match=message):
        call()


# ── create_docx ──────────────────────────────────────────────────────────────


def test_create_docx_writes_into_the_document_output_directory(workspace: Path) -> None:
    result = create_docx("Report", ["intro para", "second para"], filename="my report")

    assert result["path"] == "generated_documents/my report.docx"
    written = workspace / "generated_documents" / "my report.docx"
    assert result["bytes"] == written.stat().st_size

    from docx import Document

    texts = [p.text for p in Document(str(written)).paragraphs if p.text.strip()]
    assert texts == ["Report", "intro para", "second para"]


# ── create_xlsx ──────────────────────────────────────────────────────────────


@pytest.mark.parametrize("rows", ["not a list", [["ok"], "not a row"], [1, 2]])
def test_create_xlsx_requires_a_list_of_lists(workspace: Path, rows: Any) -> None:
    with pytest.raises(ToolError, match="Rows must be a list of lists"):
        create_xlsx(rows)


def test_create_xlsx_writes_rows_and_names_the_sheet(workspace: Path) -> None:
    result = create_xlsx([["name", "qty"], ["widget", 3]], filename="stock", sheet_name="Q1")

    assert result["path"] == "generated_spreadsheets/stock.xlsx"
    assert result["rows"] == 2
    assert result["bytes"] > 0

    from openpyxl import load_workbook

    sheet = load_workbook(workspace / "generated_spreadsheets" / "stock.xlsx").active
    assert sheet.title == "Q1"
    assert [list(row) for row in sheet.iter_rows(values_only=True)] == [
        ["name", "qty"],
        ["widget", 3],
    ]


def test_create_xlsx_truncates_an_over_long_sheet_name(workspace: Path) -> None:
    create_xlsx([["x"]], filename="wide.xlsx", sheet_name="s" * 60)

    from openpyxl import load_workbook

    sheet = load_workbook(workspace / "generated_spreadsheets" / "wide.xlsx").active
    assert sheet.title == "s" * 31


def test_create_xlsx_falls_back_to_a_default_sheet_name(workspace: Path) -> None:
    create_xlsx([], filename="empty.xlsx", sheet_name="")

    from openpyxl import load_workbook

    assert load_workbook(workspace / "generated_spreadsheets" / "empty.xlsx").active.title == "Sheet1"


# ── create_pptx ──────────────────────────────────────────────────────────────


def test_create_pptx_builds_a_title_slide_plus_one_slide_per_entry(workspace: Path) -> None:
    slides = [
        {"title": "Agenda", "bullets": ["one", "two", "three"]},
        {"title": "Detail", "bullets": "single bullet as a string"},
        {},
    ]

    result = create_pptx("Quarterly", slides, filename="deck")

    assert result["path"] == "generated_presentations/deck.pptx"
    assert result["slides"] == 4  # title slide + three content slides
    assert result["bytes"] > 0

    from pptx import Presentation

    deck = Presentation(str(workspace / "generated_presentations" / "deck.pptx"))
    titles = [slide.shapes.title.text for slide in deck.slides]
    assert titles == ["Quarterly", "Agenda", "Detail", "Slide"]
    agenda_body = deck.slides[1].placeholders[1].text_frame
    assert [p.text for p in agenda_body.paragraphs] == ["one", "two", "three"]
    assert [p.text for p in deck.slides[2].placeholders[1].text_frame.paragraphs] == [
        "single bullet as a string"
    ]


def test_create_pptx_without_slides_still_produces_a_titled_deck(workspace: Path) -> None:
    result = create_pptx("", None, filename="bare.pptx")

    from pptx import Presentation

    deck = Presentation(str(workspace / "generated_presentations" / "bare.pptx"))
    assert result["slides"] == 1
    assert deck.slides[0].shapes.title.text == "Presentation"


# ── create_pdf ───────────────────────────────────────────────────────────────


def test_create_pdf_registers_a_cjk_font_when_one_is_present(
    workspace: Path, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    recorder = _install_fake_reportlab(monkeypatch)
    present = tmp_path / "NanumGothic.ttf"
    present.write_bytes(b"fake font bytes")
    monkeypatch.setattr(
        documents_module,
        "_CJK_FONT_CANDIDATES",
        [str(tmp_path / "absent.ttf"), str(present), str(tmp_path / "never-reached.ttf")],
    )

    result = create_pdf("보고서", "본문 한 줄", filename="korean")

    assert result["path"] == "generated_pdfs/korean.pdf"
    assert result["bytes"] > 0
    # The first existing candidate wins and the scan stops there.
    assert recorder.ttfonts == [("KoreanFont", str(present))]
    assert [p.style.kwargs["fontName"] for p in recorder.paragraphs] == [
        "KoreanFont",
        "KoreanFont",
    ]


def test_create_pdf_falls_back_to_helvetica_when_registration_fails(
    workspace: Path, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """A font file that reportlab cannot parse must not fail the whole export."""
    recorder = _install_fake_reportlab(monkeypatch)
    recorder.register_raises = True
    broken = tmp_path / "broken.ttf"
    broken.write_bytes(b"not really a font")
    monkeypatch.setattr(documents_module, "_CJK_FONT_CANDIDATES", [str(broken)])

    create_pdf("Title", "Body", filename="fallback.pdf")

    assert recorder.registered == []
    assert {p.style.kwargs["fontName"] for p in recorder.paragraphs} == {"Helvetica"}


def test_create_pdf_uses_helvetica_when_no_cjk_font_is_installed(
    workspace: Path, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    recorder = _install_fake_reportlab(monkeypatch)
    monkeypatch.setattr(
        documents_module,
        "_CJK_FONT_CANDIDATES",
        [str(tmp_path / ("nope-" + str(index))) for index in range(3)],
    )

    create_pdf("", ["only body"], filename="plain.pdf")

    assert recorder.ttfonts == []
    # No title was given, so the story starts with the body paragraph.
    assert [p.text for p in recorder.paragraphs] == ["only body"]
    assert recorder.doc_kwargs["leftMargin"] == pytest.approx(20 * 2.834645669291339)


def test_create_pdf_escapes_markup_in_the_body(
    workspace: Path, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """reportlab treats paragraph text as mini-HTML; raw <> would break the build."""
    recorder = _install_fake_reportlab(monkeypatch)
    monkeypatch.setattr(documents_module, "_CJK_FONT_CANDIDATES", [])

    create_pdf("Title & Co", "<script>alert(1)</script>\n\n   \n\nsecond & last")

    assert [p.text for p in recorder.paragraphs] == [
        "Title & Co",  # the title is not escaped
        "&lt;script&gt;alert(1)&lt;/script&gt;",
        "second &amp; last",
    ]
    assert recorder.doc_path.endswith("generated_pdfs/document.pdf")


# ── read_document ────────────────────────────────────────────────────────────


def test_read_document_refuses_a_missing_path(tmp_path: Path) -> None:
    with pytest.raises(ToolError, match="파일이 없습니다"):
        read_document(str(tmp_path / "ghost.txt"))


def test_read_document_refuses_a_directory(tmp_path: Path) -> None:
    with pytest.raises(ToolError, match="파일이 아닙니다"):
        read_document(str(tmp_path))


def test_read_document_refuses_a_file_over_the_read_cap(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    monkeypatch.setattr(documents_module, "DOCUMENT_MAX_READ_BYTES", 16)
    target = tmp_path / "big.txt"
    target.write_text("x" * 64, encoding="utf-8")

    with pytest.raises(ToolError, match="파일이 너무 큽니다"):
        read_document(str(target))


@pytest.mark.parametrize("suffix", [".exe", ".zip", ".png", ""])
def test_read_document_refuses_an_unsupported_extension(tmp_path: Path, suffix: str) -> None:
    target = tmp_path / ("file" + suffix)
    target.write_bytes(b"payload")

    with pytest.raises(ToolError, match="지원하지 않는 형식입니다"):
        read_document(str(target))


def test_read_document_reads_a_pdf(tmp_path: Path) -> None:
    from PIL import Image

    target = tmp_path / "scan.pdf"
    Image.new("RGB", (40, 20), "white").save(target, "PDF")

    meta = read_document(str(target))

    assert meta["ext"] == ".pdf"
    assert meta["pages"] == 1
    assert meta["chars"] == len(meta["content"])


def test_read_document_reports_a_broken_pdf(tmp_path: Path) -> None:
    target = tmp_path / "broken.pdf"
    target.write_bytes(b"this is not a pdf at all")

    with pytest.raises(ToolError, match="PDF 읽기 실패"):
        read_document(str(target))


def test_read_document_reads_a_docx(workspace: Path, tmp_path: Path) -> None:
    create_docx("Heading", ["alpha", "beta"], filename="doc.docx")
    target = workspace / "generated_documents" / "doc.docx"

    meta = read_document(str(target))

    assert meta["paragraphs"] == 3
    assert "alpha" in meta["content"]
    assert meta["preview"].startswith("Heading")


def test_read_document_reports_a_broken_docx(tmp_path: Path) -> None:
    target = tmp_path / "broken.docx"
    target.write_bytes(b"PK not really an office file")

    with pytest.raises(ToolError, match="DOCX 읽기 실패"):
        read_document(str(target))


def test_read_document_reads_every_sheet_of_an_xlsx(workspace: Path) -> None:
    create_xlsx([["a", 1], [None, "c"]], filename="grid.xlsx", sheet_name="First")
    target = workspace / "generated_spreadsheets" / "grid.xlsx"

    meta = read_document(str(target))

    assert meta["sheets"] == 1
    assert "[Sheet: First]" in meta["content"]
    assert "a\t1" in meta["content"]
    assert "\tc" in meta["content"]  # a None cell becomes an empty column


def test_read_document_reports_a_broken_xlsx(tmp_path: Path) -> None:
    target = tmp_path / "broken.xlsx"
    with zipfile.ZipFile(target, "w") as archive:
        archive.writestr("nothing.txt", "not a workbook")

    with pytest.raises(ToolError, match="XLSX 읽기 실패"):
        read_document(str(target))


def test_read_document_reads_a_pptx_slide_by_slide(workspace: Path) -> None:
    create_pptx("Deck", [{"title": "Agenda", "bullets": ["one", "two"]}], filename="deck.pptx")
    target = workspace / "generated_presentations" / "deck.pptx"

    meta = read_document(str(target))

    assert meta["slides"] == 2
    assert "[Slide 1]" in meta["content"]
    assert "[Slide 2]" in meta["content"]
    assert "Agenda" in meta["content"]
    assert "one\ntwo" in meta["content"]


def test_read_document_reports_a_broken_pptx(tmp_path: Path) -> None:
    target = tmp_path / "broken.pptx"
    target.write_bytes(b"not a presentation")

    with pytest.raises(ToolError, match="PPTX 읽기 실패"):
        read_document(str(target))


@pytest.mark.parametrize("suffix", [".txt", ".md", ".csv"])
def test_read_document_reads_plain_text_formats(tmp_path: Path, suffix: str) -> None:
    target = tmp_path / ("notes" + suffix)
    target.write_text("col1,col2\nvalue,other\n", encoding="utf-8")

    meta = read_document(str(target))

    assert meta["ext"] == suffix
    assert meta["content"] == "col1,col2\nvalue,other\n"
    assert meta["chars"] == len(meta["content"])
    assert meta["path"] == str(target)


def test_read_document_caps_the_preview_and_the_content(tmp_path: Path) -> None:
    target = tmp_path / "long.md"
    target.write_text("y" * 60_000, encoding="utf-8")

    meta = read_document(str(target))

    assert meta["chars"] == 60_000
    assert len(meta["preview"]) == 500
    assert len(meta["content"]) == 50_000


def test_read_document_reports_an_unreadable_text_file(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """A file the OS refuses mid-read must surface as a ToolError, not an OSError."""
    target = tmp_path / "locked.txt"
    target.write_text("secret", encoding="utf-8")
    original = Path.read_text

    def refuse(self: Path, *args: Any, **kwargs: Any):
        if self == target:
            raise PermissionError("permission denied")
        return original(self, *args, **kwargs)

    monkeypatch.setattr(Path, "read_text", refuse)

    with pytest.raises(ToolError, match="파일 읽기 실패"):
        read_document(str(target))
