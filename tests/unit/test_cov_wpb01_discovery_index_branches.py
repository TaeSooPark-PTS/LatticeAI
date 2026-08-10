"""wpb01 branch coverage — ``lattice_brain.graph.discovery_index``.

Three seams:

* ``_extract_local_file_text`` — the "nothing to extract here" directions of
  the per-format parsers: an all-empty .docx table row, .pptx shapes with no
  text frame / no text, a slide that contributes nothing, and a category whose
  extension has no extractor at all.
* the graph maintenance around a file node — linked nodes that are neither
  owned children nor auto-extracted candidates, and an auto-extracted node
  that is still referenced from outside the candidate set (so it must survive
  both re-indexing and deletion), plus deleting a file node that is not there.
* ``index_local_folder`` on a *first* scan, where the three failure paths
  (not indexable / unreadable bytes / parser error) have no previous index row
  to clean up.
"""

from __future__ import annotations

import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

import lattice_brain.graph._kg_common as kg_common  # noqa: E402
from lattice_brain.graph.store import KnowledgeGraphStore  # noqa: E402

CONCEPT_TEXT = (
    "Lattice AI stores the Knowledge Graph locally. "
    "Graph RAG improves retrieval quality for Lattice AI users."
)


@pytest.fixture(autouse=True)
def _rule_based_extraction(monkeypatch: pytest.MonkeyPatch) -> None:
    """Concept extraction must come from the deterministic rule engine.

    A global LLM router left registered by another module would otherwise make
    the extracted concepts (and therefore the graph shape asserted here)
    depend on test ordering.
    """
    monkeypatch.setattr(kg_common, "get_llm_router", lambda: None)


@pytest.fixture()
def store(tmp_path: Path) -> KnowledgeGraphStore:
    return KnowledgeGraphStore(tmp_path / "kg.sqlite", tmp_path / "blobs")


def _png(path: Path) -> Path:
    from PIL import Image

    Image.new("RGB", (10, 10), "red").save(path)
    return path


def _file_node_id(store: KnowledgeGraphStore, source_id: str) -> str:
    with store._connect() as conn:
        row = conn.execute(
            "SELECT graph_node_id FROM local_file_index WHERE source_id=? AND status='indexed'",
            (source_id,),
        ).fetchone()
    assert row is not None and row["graph_node_id"]
    return str(row["graph_node_id"])


def _attach_satellites(store: KnowledgeGraphStore, file_node_id: str) -> None:
    """A plain neighbour and an auto-extracted node with an outside reference.

    Neither may be treated as a child of the file node: the plain one is not
    auto-extracted, and the auto-extracted one is still cited by a node the
    file does not own.
    """
    with store._connect() as conn:
        store._upsert_node(
            conn, "plain:other", "Concept", "Other", summary="", metadata={}
        )
        store._upsert_node(
            conn,
            "auto:widget",
            "Concept",
            "Widget",
            summary="",
            metadata={"auto_extracted": True, "source": "local_folder"},
        )
        store._upsert_edge(conn, file_node_id, "plain:other", "언급함", weight=0.5, metadata={})
        store._upsert_edge(conn, file_node_id, "auto:widget", "언급함", weight=0.5, metadata={})
        store._upsert_edge(conn, "auto:widget", "plain:other", "언급함", weight=0.5, metadata={})


def _node_ids(store: KnowledgeGraphStore) -> set:
    with store._connect() as conn:
        return {row["id"] for row in conn.execute("SELECT id FROM nodes")}


# ── _extract_local_file_text ────────────────────────────────────────────────


def test_docx_table_row_with_only_empty_cells_is_skipped(
    store: KnowledgeGraphStore, tmp_path: Path
) -> None:
    from docx import Document

    doc = Document()
    doc.add_paragraph("Intro paragraph")
    table = doc.add_table(rows=2, cols=2)
    table.rows[1].cells[0].text = "value"
    path = tmp_path / "tables.docx"
    doc.save(str(path))

    text, meta = store._extract_local_file_text(path, "document", include_ocr=False)

    assert meta["tables"] == 1
    # Two rows in the table, but only the one carrying content became a line.
    assert meta["table_rows"] == 1
    assert "Intro paragraph" in text
    assert "value" in text


def test_pptx_shapes_without_text_and_slides_without_content_are_skipped(
    store: KnowledgeGraphStore, tmp_path: Path
) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # blank: contributes nothing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(
        str(_png(tmp_path / "pic.png")), Inches(1), Inches(1), Inches(1), Inches(1)
    )
    slide.shapes.add_textbox(Inches(1), Inches(3), Inches(2), Inches(1))  # empty text
    filled = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(2), Inches(1))
    filled.text_frame.text = "Real slide content"
    path = tmp_path / "deck.pptx"
    prs.save(str(path))

    text, meta = store._extract_local_file_text(path, "slide_deck", include_ocr=False)

    assert meta["slides"] == 2
    assert meta["text_slides"] == 1
    assert text == "[Slide 2]\nReal slide content"


def test_a_category_with_no_extractor_yields_parser_metadata_only(
    store: KnowledgeGraphStore, tmp_path: Path
) -> None:
    path = tmp_path / "payload.bin"
    path.write_bytes(b"\x00\x01\x02binary")

    text, meta = store._extract_local_file_text(path, "unsupported", include_ocr=False)

    assert text == ""
    assert meta == {"parser": "bin"}


# ── graph maintenance around a file node ────────────────────────────────────


def test_reindex_keeps_neighbours_the_file_does_not_own(
    store: KnowledgeGraphStore, tmp_path: Path
) -> None:
    root = tmp_path / "folder"
    root.mkdir()
    note = root / "notes.txt"
    note.write_text(CONCEPT_TEXT, encoding="utf-8")
    first = store.index_local_folder(root)
    source_id = first["source"]["id"]
    file_node_id = _file_node_id(store, source_id)
    _attach_satellites(store, file_node_id)

    note.write_text(CONCEPT_TEXT + " A second revision adds more text.", encoding="utf-8")
    second = store.index_local_folder(root)

    assert second["counts"]["indexed"] == 1
    remaining = _node_ids(store)
    # The plain neighbour is not a child, and the auto-extracted node is still
    # cited from outside the candidate set, so neither is garbage-collected.
    assert "plain:other" in remaining
    assert "auto:widget" in remaining


def test_deleting_a_file_node_keeps_externally_cited_neighbours(
    store: KnowledgeGraphStore, tmp_path: Path
) -> None:
    root = tmp_path / "folder"
    root.mkdir()
    (root / "notes.txt").write_text(CONCEPT_TEXT, encoding="utf-8")
    result = store.index_local_folder(root)
    file_node_id = _file_node_id(store, result["source"]["id"])
    _attach_satellites(store, file_node_id)

    with store._connect() as conn:
        store._delete_local_file_graph(conn, file_node_id)

    remaining = _node_ids(store)
    assert file_node_id not in remaining
    assert "plain:other" in remaining
    assert "auto:widget" in remaining
    with store._connect() as conn:
        chunks = conn.execute(
            "SELECT COUNT(*) AS c FROM chunks WHERE source_node=?", (file_node_id,)
        ).fetchone()["c"]
    assert chunks == 0


def test_two_spellings_of_one_concept_never_produce_a_self_edge(
    store: KnowledgeGraphStore, tmp_path: Path
) -> None:
    """Both spellings slug to the same node id, so the triple between them
    would be an edge from a node to itself — it must be dropped."""
    root = tmp_path / "folder"
    root.mkdir()
    (root / "notes.txt").write_text(
        "The `mlx vlm` runtime and the \"mlx-vlm\" package are the same thing.",
        encoding="utf-8",
    )

    store.index_local_folder(root)

    with store._connect() as conn:
        assert conn.execute(
            "SELECT COUNT(*) AS c FROM nodes WHERE id='concept:mlx-vlm'"
        ).fetchone()["c"] == 1
        assert conn.execute(
            "SELECT COUNT(*) AS c FROM edges WHERE from_node=to_node"
        ).fetchone()["c"] == 0


def test_deleting_an_unknown_file_node_is_a_no_op(
    store: KnowledgeGraphStore, tmp_path: Path
) -> None:
    root = tmp_path / "folder"
    root.mkdir()
    (root / "notes.txt").write_text(CONCEPT_TEXT, encoding="utf-8")
    store.index_local_folder(root)
    before = _node_ids(store)

    with store._connect() as conn:
        store._delete_local_file_graph(conn, "local-file:never-existed")

    assert _node_ids(store) == before


# ── index_local_folder: first scan, nothing to clean up ─────────────────────


def test_first_scan_records_each_failure_without_a_previous_row(
    store: KnowledgeGraphStore, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    root = tmp_path / "folder"
    root.mkdir()
    (root / "good.txt").write_text(CONCEPT_TEXT, encoding="utf-8")
    (root / "notes.xyz").write_text("unsupported extension", encoding="utf-8")
    (root / "broken.pdf").write_bytes(b"not a pdf at all")
    (root / "unreadable.txt").write_text("never read", encoding="utf-8")

    real_read_bytes = Path.read_bytes

    def fake_read_bytes(self, *args, **kwargs):
        if self.name == "unreadable.txt":
            raise OSError(5, "simulated I/O error")
        return real_read_bytes(self, *args, **kwargs)

    monkeypatch.setattr(Path, "read_bytes", fake_read_bytes)

    result = store.index_local_folder(root)

    counts = result["counts"]
    assert counts["indexed"] == 1
    assert counts["unsupported"] == 1
    assert counts["failed"] == 2
    failed_paths = {Path(err["path"]).name for err in result["errors"]}
    assert failed_paths == {"unreadable.txt", "broken.pdf"}
    with store._connect() as conn:
        statuses = {
            row["relative_path"]: row["status"]
            for row in conn.execute(
                "SELECT relative_path, status FROM local_file_index WHERE source_id=?",
                (result["source"]["id"],),
            )
        }
    assert statuses == {
        "good.txt": "indexed",
        "notes.xyz": "unsupported",
        "broken.pdf": "failed",
        "unreadable.txt": "failed",
    }
