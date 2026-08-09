"""wp13 coverage — the ``latticeai.tools`` package seam itself.

``TOOL_HANDLERS`` is the single name → invocation table, and three of its
entries are more than a lambda: the xlsx and pptx handlers accept their
structured argument as either real data or a JSON string (models emit both),
and every knowledge tool goes through ``_knowledge_scope``, which refuses to
run unless the caller injected an authenticated workspace *and* user. Those
three, plus the public path resolver, are what this file drives — through
``execute_tool`` rather than by calling the private helpers, so the table
wiring is under test too.
"""

from __future__ import annotations

import json
import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

import latticeai.tools as tools
from latticeai.tools import (
    ToolError,
    execute_tool,
    registered_tools,
    resolve_workspace_path,
)
from latticeai.tools import knowledge as knowledge_module


@pytest.fixture()
def workspace(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> Path:
    root = tmp_path / "agent_workspace"
    root.mkdir()
    monkeypatch.setattr(tools, "AGENT_ROOT", root)
    tools.ensure_agent_root()
    return root


# ── resolve_workspace_path ───────────────────────────────────────────────────


def test_resolve_workspace_path_is_the_sandboxed_resolver(workspace: Path) -> None:
    assert resolve_workspace_path() == workspace
    assert resolve_workspace_path("sub/file.txt") == workspace / "sub" / "file.txt"


def test_resolve_workspace_path_refuses_an_escape(workspace: Path) -> None:
    with pytest.raises(ToolError, match="escapes the agent workspace"):
        resolve_workspace_path("../../etc/passwd")


def test_resolve_workspace_path_creates_the_workspace_on_demand(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    root = tmp_path / "not-created-yet"
    monkeypatch.setattr(tools, "AGENT_ROOT", root)

    assert resolve_workspace_path() == root
    assert root.is_dir()


# ── the xlsx / pptx argument adapters ────────────────────────────────────────


def test_create_xlsx_handler_accepts_rows_as_data(workspace: Path) -> None:
    result = execute_tool(
        "create_xlsx",
        {"rows": [["a", "b"], ["c", "d"]], "filename": "data.xlsx", "sheet_name": "Tab"},
    )

    from openpyxl import load_workbook

    sheet = load_workbook(workspace / "generated_spreadsheets" / "data.xlsx").active
    assert result["rows"] == 2
    assert sheet.title == "Tab"
    assert [list(row) for row in sheet.iter_rows(values_only=True)] == [["a", "b"], ["c", "d"]]


def test_create_xlsx_handler_accepts_rows_as_a_json_string(workspace: Path) -> None:
    """Models routinely emit tool arguments as JSON text rather than as data."""
    result = execute_tool("create_xlsx", {"rows": json.dumps([["x", 1]])})

    assert result["path"] == "generated_spreadsheets/spreadsheet.xlsx"
    assert result["rows"] == 1


def test_create_xlsx_handler_defaults_every_argument(workspace: Path) -> None:
    result = execute_tool("create_xlsx", {})

    from openpyxl import load_workbook

    assert result["rows"] == 0
    assert load_workbook(workspace / "generated_spreadsheets" / "spreadsheet.xlsx").active.title == (
        "Sheet1"
    )


def test_create_pptx_handler_accepts_slides_as_a_json_string(workspace: Path) -> None:
    result = execute_tool(
        "create_pptx",
        {
            "title": "Roadmap",
            "slides": json.dumps([{"title": "Now", "bullets": ["ship"]}]),
            "filename": "plan.pptx",
        },
    )

    from pptx import Presentation

    deck = Presentation(str(workspace / "generated_presentations" / "plan.pptx"))
    assert result["slides"] == 2
    assert [slide.shapes.title.text for slide in deck.slides] == ["Roadmap", "Now"]


def test_create_pptx_handler_defaults_to_an_empty_deck(workspace: Path) -> None:
    result = execute_tool("create_pptx", {})

    assert result["path"] == "generated_presentations/presentation.pptx"
    assert result["slides"] == 1


def test_a_malformed_json_argument_surfaces_as_a_json_error(workspace: Path) -> None:
    with pytest.raises(json.JSONDecodeError):
        execute_tool("create_xlsx", {"rows": "[[unquoted]]"})


# ── _knowledge_scope ─────────────────────────────────────────────────────────


@pytest.fixture()
def vault(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> Path:
    root = tmp_path / "brain"
    root.mkdir()
    monkeypatch.setattr(knowledge_module, "BRAIN_DIR", root)
    return root


_SCOPED_TOOLS = [
    ("knowledge_save", {"content": "body"}),
    ("knowledge_search", {"query": "body"}),
    ("knowledge_tree", {}),
    ("obsidian_save", {"content": "body"}),
    ("obsidian_search", {"query": "body"}),
    ("obsidian_tree", {}),
]


@pytest.mark.parametrize("action,args", _SCOPED_TOOLS)
def test_every_knowledge_tool_refuses_an_unauthenticated_call(vault: Path, action, args) -> None:
    with pytest.raises(ToolError, match="authenticated workspace and user scope"):
        execute_tool(action, dict(args))


@pytest.mark.parametrize(
    "scope",
    [
        {"workspace_id": "ws-1"},
        {"user_email": "owner@example.com"},
        {"workspace_id": "", "user_email": "owner@example.com"},
        {"workspace_id": "ws-1", "user_email": "   "},
    ],
)
def test_a_half_scoped_knowledge_call_is_refused_at_the_registry(vault: Path, scope) -> None:
    with pytest.raises(ToolError, match="authenticated workspace and user scope"):
        execute_tool("knowledge_tree", dict(scope))


def test_a_scoped_knowledge_call_writes_into_its_own_partition(vault: Path) -> None:
    scope = {"workspace_id": "ws-1", "user_email": "Owner@Example.com"}

    saved = execute_tool("knowledge_save", {"content": "scoped note", "title": "n", **scope})
    tree = execute_tool("knowledge_tree", dict(scope))
    found = execute_tool("knowledge_search", {"query": "scoped", **scope})

    assert Path(saved["path"]).is_relative_to(Path(tree["root"]))
    assert [entry["relative_path"] for entry in tree["entries"]] == ["00_Raw/n.md"]
    assert len(found["results"]) == 1
    # The email is normalised before it becomes part of the partition digest.
    upper = execute_tool(
        "knowledge_tree", {"workspace_id": "ws-1", "user_email": "OWNER@EXAMPLE.COM"}
    )
    assert upper["root"] == tree["root"]


def test_scopes_cannot_read_each_others_notes(vault: Path) -> None:
    execute_tool(
        "knowledge_save",
        {"content": "first tenant", "title": "a", "workspace_id": "ws-1", "user_email": "a@x.com"},
    )

    other = execute_tool(
        "knowledge_search", {"query": "tenant", "workspace_id": "ws-2", "user_email": "a@x.com"}
    )

    assert other["results"] == []


# ── the table itself ─────────────────────────────────────────────────────────


def test_unknown_actions_are_rejected_by_name(workspace: Path) -> None:
    with pytest.raises(ToolError, match="Unknown action: not_a_tool"):
        execute_tool("not_a_tool", {})


def test_every_scoped_knowledge_tool_is_actually_registered() -> None:
    assert {name for name, _ in _SCOPED_TOOLS} <= registered_tools()
