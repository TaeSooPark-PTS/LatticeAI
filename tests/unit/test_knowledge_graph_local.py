"""Local folder Graph RAG tests."""
import os
import sys
import time
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).parent.parent.parent))

from fastapi import FastAPI
from fastapi.testclient import TestClient

from lattice_brain.graph.store import KnowledgeGraphStore
from latticeai.services.local_knowledge import LocalKnowledgeWatcher, create_local_knowledge_router


def _store(tmp_path: Path) -> KnowledgeGraphStore:
    return KnowledgeGraphStore(tmp_path / "kg.sqlite", tmp_path / "blobs")


def test_local_folder_index_emits_scoped_brain_ingestion_event():
    events = []
    indexed = {}

    class _Graph:
        def index_local_folder(self, *_args, **kwargs):
            indexed.update(kwargs)
            return {
                "status": "ok",
                "source": {"id": "source-local"},
                "counts": {"indexed": 2, "deleted": 0},
            }

    class _Hooks:
        @staticmethod
        def fire_hook(kind, event, **kwargs):
            events.append((kind, event, kwargs))
            return {"blocked": False}

    class _WorkspaceService:
        @staticmethod
        def resolve_write_scope(requested, user):
            assert requested == "org:acme"
            assert user == "owner@example.com"
            return "org:acme"

    app = FastAPI()
    app.include_router(create_local_knowledge_router(
        get_graph=lambda: _Graph(),
        require_graph=lambda: None,
        require_user=lambda _request: "owner@example.com",
        require_local_user=lambda _request: "owner@example.com",
        local_permission_response=lambda *_args, **_kwargs: {},
        require_local_approval=lambda **_kwargs: None,
        hooks=_Hooks(),
        workspace_service=_WorkspaceService(),
    ))

    response = TestClient(app).post(
        "/knowledge-graph/local/index",
        headers={"X-Workspace-Id": "org:acme"},
        json={"path": "/approved/source", "approved": True},
    )

    assert response.status_code == 200
    brain_event = next(item for item in events if item[1] == "tool.kg_ingest.local_folder")
    assert brain_event[0] == "post_tool"
    assert brain_event[2]["payload"]["source_type"] == "local_folder"
    assert brain_event[2]["user_email"] == "owner@example.com"
    assert brain_event[2]["workspace_id"] == "org:acme"
    assert indexed["workspace_id"] == "org:acme"
    assert indexed["consent"]["approved_by"] == "owner@example.com"
    assert indexed["consent"]["workspace_id"] == "org:acme"


def test_local_folder_index_rejects_unauthorized_workspace_before_scanning():
    calls = []

    class _Graph:
        def index_local_folder(self, *_args, **_kwargs):
            calls.append("indexed")
            return {"status": "ok"}

    class _WorkspaceService:
        @staticmethod
        def resolve_write_scope(_requested, _user):
            raise PermissionError("workspace write denied")

    app = FastAPI()
    app.include_router(create_local_knowledge_router(
        get_graph=lambda: _Graph(),
        require_graph=lambda: None,
        require_user=lambda _request: "owner@example.com",
        require_local_user=lambda _request: "owner@example.com",
        local_permission_response=lambda *_args, **_kwargs: {},
        require_local_approval=lambda **_kwargs: None,
        workspace_service=_WorkspaceService(),
    ))

    response = TestClient(app).post(
        "/knowledge-graph/local/index",
        headers={"X-Workspace-Id": "org:forbidden"},
        json={"path": "/approved/source", "approved": True},
    )

    assert response.status_code == 403
    assert calls == []


def test_watcher_restore_recovers_workspace_scope_from_persisted_consent():
    captured = []

    class _Graph:
        @staticmethod
        def local_sources():
            return {"sources": [{
                "id": "source:one",
                "root_path": "/approved/source",
                "watch_enabled": True,
                "consent": {
                    "approved_by": "owner@example.com",
                    "workspace_id": "org:acme",
                },
            }]}

    watcher = LocalKnowledgeWatcher(lambda: _Graph())
    watcher.start_source = lambda source: captured.append(source) or {"watching": True}

    result = watcher.restore_enabled_sources()

    assert result["restored"] == 1
    assert captured[0]["workspace_id"] == "org:acme"


def test_audit_local_folder_classifies_supported_sensitive_and_unsupported(tmp_path):
    root = tmp_path / "source"
    root.mkdir()
    (root / "notes.md").write_text("Graph RAG local folder scan", encoding="utf-8")
    (root / "app.py").write_text("def run():\n    return 'Lattice AI'\n", encoding="utf-8")
    (root / ".env").write_text("TOKEN=secret", encoding="utf-8")
    (root / "payload.bin").write_bytes(b"\x00\x01")

    audit = _store(tmp_path).audit_local_folder(root)

    assert audit["summary"]["total_files"] == 4
    assert audit["summary"]["readable_files"] == 2
    assert audit["summary"]["sensitive_files"] == 1
    assert audit["summary"]["unsupported_files"] == 1
    assert audit["by_category"]["text"] == 1
    assert audit["by_category"]["code"] == 1


def test_index_local_folder_creates_graph_and_skips_unchanged_files(tmp_path):
    root = tmp_path / "source"
    src = root / "src"
    src.mkdir(parents=True)
    (root / "notes.md").write_text(
        "Lattice AI uses Graph RAG. TODO implement local folder scan.",
        encoding="utf-8",
    )
    (src / "app.py").write_text("def graph_rag():\n    return 'local knowledge'\n", encoding="utf-8")
    excluded = root / "node_modules" / "pkg"
    excluded.mkdir(parents=True)
    (excluded / "index.js").write_text("console.log('skip')", encoding="utf-8")

    store = _store(tmp_path)
    first = store.index_local_folder(root, user_email="user@example.com")
    second = store.index_local_folder(root, user_email="user@example.com")

    assert first["counts"]["indexed"] == 2
    assert second["counts"]["skipped_unchanged"] == 2

    stats = store.stats()
    assert stats["local_sources"] == 1
    assert stats["local_file_status"]["indexed"] == 2
    assert stats["nodes"]["Folder"] >= 2
    assert stats["nodes"]["Document"] >= 1
    assert stats["nodes"]["CodeFile"] >= 1

    matches = store.search("Graph RAG")["matches"]
    assert any(match["title"] == "notes.md" for match in matches)


def test_local_folder_nodes_are_scoped_and_cannot_be_reassigned(tmp_path):
    root = tmp_path / "shared-source"
    root.mkdir()
    (root / "roadmap.md").write_text(
        "Project Atlas roadmap decision and private workspace context.",
        encoding="utf-8",
    )

    store = _store(tmp_path)
    first = store.index_local_folder(
        root,
        user_email="owner@example.com",
        workspace_id="org:alpha",
    )
    first_node = first["indexed_nodes"][0]
    assert store.workspaces_of([first_node]) == {first_node: "org:alpha"}
    assert store.filter_scoped_nodes([{"id": first_node}], {"org:alpha"})
    assert store.filter_scoped_nodes([{"id": first_node}], {"org:beta"}) == []
    with store._connect() as conn:
        scopes = {row["workspace_id"] for row in conn.execute("SELECT workspace_id FROM nodes_v2")}
    assert scopes == {"org:alpha"}, "file, hierarchy, chunk, concept, and semantic nodes must share scope"
    assert store.local_sources()["sources"][0]["consent"]["workspace_id"] == "org:alpha"

    with pytest.raises(ValueError, match="another workspace"):
        store.index_local_folder(
            root,
            user_email="owner@example.com",
            workspace_id="org:beta",
        )
    with pytest.raises(ValueError, match="another workspace"):
        store.index_local_folder(
            root,
            user_email="owner@example.com",
            workspace_id="org:beta",
            source_id_override=first["source"]["id"],
        )


def test_legacy_local_folder_nodes_reproject_to_personal_without_id_rewrite(tmp_path):
    root = tmp_path / "legacy-source"
    root.mkdir()
    (root / "notes.md").write_text("Legacy local Brain knowledge.", encoding="utf-8")

    store = _store(tmp_path)
    legacy = store.index_local_folder(root, user_email="owner@example.com")
    source_id = legacy["source"]["id"]
    node_id = legacy["indexed_nodes"][0]
    assert store.workspaces_of([node_id])[node_id] is None

    reprojected = store.index_local_folder(
        root,
        user_email="owner@example.com",
        workspace_id="personal",
        source_id_override=source_id,
    )

    assert reprojected["source"]["id"] == source_id
    assert reprojected["indexed_nodes"] == [node_id]
    assert store.workspaces_of([node_id])[node_id] == "personal"
    assert store.filter_scoped_nodes([{"id": node_id}], {"org:other"}) == []
    with store._connect() as conn:
        scopes = {row["workspace_id"] for row in conn.execute("SELECT workspace_id FROM nodes_v2")}
    assert scopes == {"personal"}, "legacy local-folder artifacts must not remain globally visible"


def test_index_local_folder_skips_empty_spreadsheet_until_text_is_extracted(tmp_path):
    from openpyxl import Workbook

    root = tmp_path / "source"
    root.mkdir()
    workbook_path = root / "report.xlsx"
    Workbook().save(workbook_path)

    store = _store(tmp_path)
    result = store.index_local_folder(root, user_email="user@example.com")

    assert result["counts"]["skipped_empty_text"] == 1
    assert store.stats()["local_file_status"]["skipped_empty_text"] == 1
    assert all(node["title"] != "report.xlsx" for node in store.graph()["nodes"])


def test_index_local_folder_removes_stale_graph_node_when_spreadsheet_loses_text(tmp_path):
    from openpyxl import Workbook

    root = tmp_path / "source"
    root.mkdir()
    workbook_path = root / "quarterly.xlsx"
    workbook = Workbook()
    workbook.active["A1"] = "Lattice revenue forecast graph signal"
    workbook.save(workbook_path)

    store = _store(tmp_path)
    first = store.index_local_folder(root, user_email="user@example.com")
    assert first["counts"]["indexed"] == 1
    assert any(node["title"] == "quarterly.xlsx" for node in store.graph()["nodes"])

    Workbook().save(workbook_path)
    future = time.time() + 2
    os.utime(workbook_path, (future, future))
    second = store.index_local_folder(root, user_email="user@example.com")

    assert second["counts"]["skipped_empty_text"] == 1
    assert all(node["title"] != "quarterly.xlsx" for node in store.graph()["nodes"])
    with store._connect() as conn:
        row = conn.execute(
            "SELECT status, graph_node_id FROM local_file_index WHERE file_name=?",
            ("quarterly.xlsx",),
        ).fetchone()
    assert row["status"] == "skipped_empty_text"
    assert row["graph_node_id"] is None


def test_index_local_folder_skips_pdf_when_parser_extracts_no_text(tmp_path):
    root = tmp_path / "source"
    root.mkdir()
    pdf_path = root / "scan.pdf"
    pdf_path.write_bytes(b"%PDF-1.4\n% empty image-only fixture\n")

    store = _store(tmp_path)

    def empty_pdf_text(path, category, include_ocr=False):
        assert path == pdf_path
        assert category == "pdf"
        return "", {"parser": "pdfplumber"}

    store._extract_local_file_text = empty_pdf_text
    result = store.index_local_folder(root, user_email="user@example.com")

    assert result["counts"]["skipped_empty_text"] == 1
    assert all(node["title"] != "scan.pdf" for node in store.graph()["nodes"])


def test_index_local_folder_skips_empty_word_and_powerpoint_files(tmp_path):
    from docx import Document
    from pptx import Presentation

    root = tmp_path / "source"
    root.mkdir()
    Document().save(root / "blank.docx")
    Presentation().save(root / "blank.pptx")

    store = _store(tmp_path)
    result = store.index_local_folder(root, user_email="user@example.com")

    assert result["counts"]["skipped_empty_text"] == 2
    graph_titles = {node["title"] for node in store.graph()["nodes"]}
    assert "blank.docx" not in graph_titles
    assert "blank.pptx" not in graph_titles


def test_index_local_folder_extracts_word_tables_and_powerpoint_text(tmp_path):
    from docx import Document
    from pptx import Presentation
    from pptx.util import Inches

    root = tmp_path / "source"
    root.mkdir()

    doc = Document()
    table = doc.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Owner"
    table.cell(0, 1).text = "Lattice contract renewal"
    doc.save(root / "contract.docx")

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(1))
    box.text = "Roadmap decision: ship graph text extraction"
    deck.save(root / "roadmap.pptx")

    store = _store(tmp_path)
    result = store.index_local_folder(root, user_email="user@example.com")

    assert result["counts"]["indexed"] == 2
    nodes = {node["title"]: node for node in store.graph(limit=20)["nodes"]}
    assert "Lattice contract renewal" in nodes["contract.docx"]["summary"]
    assert "Roadmap decision" in nodes["roadmap.pptx"]["summary"]


def test_index_local_folder_marks_missing_files_deleted(tmp_path):
    root = tmp_path / "source"
    root.mkdir()
    target = root / "notes.md"
    target.write_text("Lattice AI Graph RAG", encoding="utf-8")

    store = _store(tmp_path)
    store.index_local_folder(root, user_email="user@example.com")
    target.unlink()
    result = store.index_local_folder(root, user_email="user@example.com")

    assert result["counts"]["deleted"] == 1
    assert store.stats()["local_file_status"]["deleted"] == 1
    rep = store._v2_sync_report()
    assert rep.get("in_sync", False), f"v2 drift after reindex-delete: {rep}"


def test_set_local_source_watch_updates_source(tmp_path):
    root = tmp_path / "source"
    root.mkdir()
    (root / "notes.md").write_text("Lattice AI Graph RAG", encoding="utf-8")

    store = _store(tmp_path)
    source_id = store.index_local_folder(root, user_email="user@example.com")["source"]["id"]
    store.set_local_source_watch(source_id, True)

    source = store.local_sources()["sources"][0]
    assert source["id"] == source_id
    assert source["watch_enabled"] is True


def test_remove_local_source_removes_only_derived_index(tmp_path):
    root = tmp_path / "source"
    root.mkdir()
    local_file = root / "notes.md"
    local_file.write_text("Lattice AI Graph RAG", encoding="utf-8")

    store = _store(tmp_path)
    source_id = store.index_local_folder(root, user_email="user@example.com")["source"]["id"]
    result = store.remove_local_source(source_id)

    assert result["source_id"] == source_id
    assert local_file.exists()
    assert store.local_sources()["sources"] == []
    assert store.stats()["local_sources"] == 0

    # explicit drift assertion for local-folder delete/reindex/remove paths
    rep = store._v2_sync_report()
    assert rep.get("in_sync", False), f"v2 drift after remove_local_source: {rep}"
