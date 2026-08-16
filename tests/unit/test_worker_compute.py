"""The v11.6.0 pure-compute seams: data in, data out, and no store touched.

Parity against the write-path call, render bytes re-opened by the same library,
honesty when a port is absent. Optional libs are stubbed via
``sys.modules[name] = None``; the PDF font hunt uses reportlab's ``Vera.ttf``.
"""

from __future__ import annotations

import base64
import io
import os
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional

import pytest
import reportlab
from fastapi import FastAPI, HTTPException, Request
from fastapi.testclient import TestClient

from latticeai.api import worker_compute
from latticeai.api.agent_worker_seam import SEAM_ENV_VAR, SEAM_RATE_BUCKET
from latticeai.api.worker_compute import (
    EMBED_KINDS,
    EXTRACT_KINDS,
    EXTRACT_LIMITS,
    PASSAGE_MAX_CHARS,
    WORKER_COMPUTE_MESSAGES,
    build_docx_bytes,
    build_extract_reply,
    build_pdf_bytes,
    build_pptx_bytes,
    build_xlsx_bytes,
    create_worker_compute_router,
)
from latticeai.core.embedding_providers import resolve_embedder
from latticeai.core.messages import LANGUAGE_HEADER, MESSAGES
from latticeai.services.voice_capture import MAX_AUDIO_BYTES

USER = "worker@local"
VERA_TTF = str(Path(reportlab.__file__).parent / "fonts" / "Vera.ttf")

COMPUTE_PATHS = [
    ("/worker/embed", {"texts": ["x"]}),
    ("/worker/parse", {"filename": "a.txt", "content_b64": ""}),
    ("/worker/render/docx", {}),
    ("/worker/render/xlsx", {}),
    ("/worker/render/pptx", {}),
    ("/worker/render/pdf", {}),
    ("/worker/asr", {"audio_b64": ""}),
    ("/worker/extract", {"text": "x"}),
]


class Recorder:
    """Every admission this router performed, in order."""

    def __init__(self) -> None:
        self.rate: List[tuple] = []

    def enforce(self, user: str, bucket: str) -> None:
        self.rate.append((user, bucket))


def _b64(data: bytes) -> str:
    return base64.b64encode(data).decode("ascii")


def _client(
    *,
    embedder: Any = "default",
    transcriber: Any = None,
    require_user: Any = None,
    recorder: Optional[Recorder] = None,
) -> TestClient:
    """A bare app carrying only the compute router — no product wiring at all."""
    if embedder == "default":
        embedder = resolve_embedder("")
    recorder = recorder or Recorder()
    app = FastAPI()
    app.include_router(
        create_worker_compute_router(
            embedder=embedder,
            transcriber=transcriber,
            require_user=require_user or (lambda request: USER),
            enforce_rate_limit=recorder.enforce,
        )
    )
    client = TestClient(app, raise_server_exceptions=False)
    client.recorder = recorder  # type: ignore[attr-defined]
    return client


@pytest.fixture(autouse=True)
def _seam_open(monkeypatch):
    """The gate is open for every test that does not deliberately close it."""
    monkeypatch.setenv(SEAM_ENV_VAR, "1")


@pytest.mark.parametrize("path,body", COMPUTE_PATHS)
def test_every_seam_is_404_while_the_host_has_not_opened_it(monkeypatch, path, body):
    monkeypatch.delenv(SEAM_ENV_VAR, raising=False)
    client = _client()

    response = client.post(path, json=body)

    assert response.status_code == 404
    assert response.json()["detail"] == MESSAGES["agent_seam.disabled"]["ko"]


@pytest.mark.parametrize("path,body", COMPUTE_PATHS)
def test_every_seam_authenticates_before_it_computes(path, body):
    def _refuse(request: Request) -> str:
        raise HTTPException(status_code=401, detail="nope")

    client = _client(require_user=_refuse)

    assert client.post(path, json=body).status_code == 401
    assert client.recorder.rate == []


def test_the_admitted_call_is_charged_to_the_per_step_seam_bucket():
    client = _client()

    client.post("/worker/embed", json={"texts": ["hello"]})

    assert client.recorder.rate == [(USER, SEAM_RATE_BUCKET)]


def test_the_messages_are_bilingual_and_land_in_the_one_catalog():
    for key, entry in WORKER_COMPUTE_MESSAGES.items():
        assert MESSAGES[key] is entry or MESSAGES[key] == entry
        assert entry["ko"] and entry["en"] and entry["ko"] != entry["en"]
        assert any("가" <= ch <= "힣" for ch in entry["ko"])
        assert not any("가" <= ch <= "힣" for ch in entry["en"])




def test_the_vectors_are_the_ones_the_write_path_would_have_stored():
    """Parity with the literal call ``write_master._upsert_vector_item`` makes."""
    from lattice_brain.embeddings import LocalEmbeddingModel

    embedder = resolve_embedder("")
    client = _client(embedder=embedder)
    text = "Lattice AI 로컬 브레인 vector parity check"

    body = client.post("/worker/embed", json={"texts": [text], "kind": "passage"}).json()

    inline = embedder.provider.embed(text[:PASSAGE_MAX_CHARS])
    assert body["vectors"] == [inline]
    assert body["vectors"][0] == LocalEmbeddingModel().embed(text)
    assert body["dim"] == embedder.provider.dim == len(inline)
    assert body["provider"] == "hash"
    assert body["model_id"] == embedder.provider.model_id
    assert body["kind"] == "passage"


def test_a_passage_is_clamped_where_the_write_path_clamps_it():
    embedder = resolve_embedder("")
    client = _client(embedder=embedder)
    # The tail token sits past the clamp, so the two kinds cannot coincide the
    # way a repeated word would (feature hashing normalises that away).
    long_text = ("lattice " * 9_000) + "zetatoken"
    assert len(long_text) > PASSAGE_MAX_CHARS

    passage = client.post("/worker/embed", json={"texts": [long_text], "kind": "passage"}).json()
    query = client.post("/worker/embed", json={"texts": [long_text], "kind": "query"}).json()

    assert passage["vectors"] == [embedder.provider.embed(long_text[:PASSAGE_MAX_CHARS])]
    assert query["vectors"] == [embedder.provider.embed(long_text)]
    assert passage["vectors"] != query["vectors"]


def test_many_texts_come_back_in_the_order_they_were_sent():
    client = _client()

    body = client.post("/worker/embed", json={"texts": ["alpha", "beta", "gamma"]}).json()

    provider = resolve_embedder("").provider
    assert body["vectors"] == [provider.embed(t) for t in ("alpha", "beta", "gamma")]


def test_an_unknown_embedding_kind_is_refused_with_the_two_that_exist():
    client = _client()

    response = client.post(
        "/worker/embed",
        json={"texts": ["x"], "kind": "Document"},
        headers={LANGUAGE_HEADER: "en"},
    )

    assert response.status_code == 422
    assert "Document" in response.json()["detail"]
    for kind in EMBED_KINDS:
        assert kind in response.json()["detail"]


def test_a_worker_without_an_embedder_says_so_instead_of_crashing():
    client = _client(embedder=None)

    response = client.post("/worker/embed", json={"texts": ["x"]})

    assert response.status_code == 503
    assert response.json()["detail"] == MESSAGES["worker_compute.embedder_unavailable"]["ko"]


def test_the_reported_identity_is_read_after_the_call_not_before():
    """A network provider locks its width to what the model returned."""

    class WidensOnFirstCall:
        provider = "ollama"
        dim = 384
        model_id = "ollama:bge-m3:384"

        def embed_batch(self, texts):
            self.dim = 1024
            self.model_id = "ollama:bge-m3:1024"
            return [[0.5] * 1024 for _ in texts]

    resolved = type("Resolved", (), {"provider": WidensOnFirstCall(), "active": "ollama"})()
    client = _client(embedder=resolved)

    body = client.post("/worker/embed", json={"texts": ["x"]}).json()

    assert body["dim"] == 1024
    assert body["model_id"] == "ollama:bge-m3:1024"
    assert body["provider"] == "ollama"




def test_a_text_document_comes_back_as_the_parser_matrix_reports_it():
    client = _client()
    payload = "첫 문단입니다.\n\n두 번째 문단."

    body = client.post(
        "/worker/parse",
        json={"filename": "note.md", "content_b64": _b64(payload.encode("utf-8"))},
    ).json()

    assert body["ext"] == ".md"
    assert body["content"] == payload
    assert body["preview"] == payload[:500]
    assert body["chars"] == len(payload)
    assert body["filename"] == "note.md"
    # The temp file this seam wrote is not a path Rust may act on.
    assert "path" not in body


def test_a_rendered_docx_parses_back_through_the_same_worker():
    client = _client()
    rendered = client.post(
        "/worker/render/docx",
        json={"title": "보고서", "body": "한 문단.\n\n두 문단.", "filename": "r.docx"},
    ).json()

    # The render reply is bytes and cost only — the name belongs to the caller
    # that asked for it, and to the Rust writer that resolved the target.
    assert "filename" not in rendered

    body = client.post(
        "/worker/parse",
        json={"filename": "r.docx", "content_b64": rendered["content_b64"]},
    ).json()

    assert body["ext"] == ".docx"
    assert body["paragraphs"] == 3
    assert "한 문단." in body["content"]


def test_a_rendered_xlsx_and_pptx_report_their_sheet_and_slide_counts():
    client = _client()
    xlsx = client.post("/worker/render/xlsx", json={"rows": [["a", 1]]}).json()
    pptx = client.post(
        "/worker/render/pptx", json={"title": "T", "slides": [{"title": "S"}]}
    ).json()

    sheet = client.post(
        "/worker/parse", json={"filename": "s.xlsx", "content_b64": xlsx["content_b64"]}
    ).json()
    deck = client.post(
        "/worker/parse", json={"filename": "d.pptx", "content_b64": pptx["content_b64"]}
    ).json()

    assert sheet["sheets"] == 1
    assert "a" in sheet["content"]
    assert deck["slides"] == 2


def test_a_rendered_pdf_parses_back_with_its_page_count():
    client = _client()
    pdf = client.post("/worker/render/pdf", json={"title": "T", "body": "hello pdf"}).json()

    body = client.post(
        "/worker/parse", json={"filename": "d.pdf", "content_b64": pdf["content_b64"]}
    ).json()

    assert body["pages"] == 1
    assert "hello pdf" in body["content"]


def test_an_unparseable_document_is_a_400_carrying_the_parser_s_reason():
    client = _client()

    response = client.post(
        "/worker/parse", json={"filename": "thing.exe", "content_b64": _b64(b"MZ")}
    )

    assert response.status_code == 400
    assert ".exe" in response.json()["detail"]


def test_a_body_that_is_not_base64_is_a_422_not_a_500():
    client = _client()

    response = client.post(
        "/worker/parse",
        json={"filename": "a.txt", "content_b64": "not base64!!"},
        headers={LANGUAGE_HEADER: "en"},
    )

    assert response.status_code == 422
    assert "not valid base64" in response.json()["detail"]




def test_the_docx_bytes_open_as_a_document_with_the_blocks_that_were_sent():
    from docx import Document

    client = _client()

    body = client.post(
        "/worker/render/docx",
        json={"title": "Title", "body": "one\n\n \n\ntwo", "filename": "my report.docx"},
    ).json()

    document = Document(io.BytesIO(base64.b64decode(body["content_b64"])))
    texts = [p.text for p in document.paragraphs]
    assert texts == ["Title", "one", "two"]
    assert body["bytes"] == len(base64.b64decode(body["content_b64"]))
    assert "filename" not in body


def test_a_docx_without_a_title_has_no_heading():
    from docx import Document

    document = Document(io.BytesIO(build_docx_bytes("", "body")))

    assert [p.text for p in document.paragraphs] == ["body"]


def test_the_xlsx_bytes_open_as_a_workbook_with_the_named_sheet():
    from openpyxl import load_workbook

    client = _client()

    body = client.post(
        "/worker/render/xlsx",
        json={"rows": [["a", 1], ["b", 2]], "sheet_name": "결과", "filename": "s"},
    ).json()

    workbook = load_workbook(io.BytesIO(base64.b64decode(body["content_b64"])))
    sheet = workbook.active
    assert sheet.title == "결과"
    assert [[c.value for c in row] for row in sheet.iter_rows()] == [["a", 1], ["b", 2]]
    assert body["rows"] == 2
    assert "filename" not in body


def test_the_pptx_bytes_open_as_a_deck_whose_first_slide_is_the_title():
    from pptx import Presentation

    client = _client()

    body = client.post(
        "/worker/render/pptx",
        json={
            "title": "Deck",
            "slides": [
                {"title": "One", "bullets": ["a", "b"]},
                {"bullets": "single"},
            ],
            "filename": "deck",
        },
    ).json()

    deck = Presentation(io.BytesIO(base64.b64decode(body["content_b64"])))
    assert len(deck.slides) == 3
    assert deck.slides[0].shapes.title.text == "Deck"
    assert deck.slides[1].shapes.title.text == "One"
    assert deck.slides[2].shapes.title.text == "Slide"
    assert body["slides"] == 3
    assert "filename" not in body


def test_a_deck_with_no_title_still_gets_the_default_one():
    from pptx import Presentation

    deck = Presentation(io.BytesIO(build_pptx_bytes("", [])))

    assert deck.slides[0].shapes.title.text == "Presentation"


def test_the_pdf_bytes_open_as_a_pdf_with_the_escaped_text():
    import pdfplumber

    client = _client()

    body = client.post(
        "/worker/render/pdf",
        json={"title": "Report", "body": "a & b <c>\n\n \n\nsecond", "filename": "r"},
    ).json()

    with pdfplumber.open(io.BytesIO(base64.b64decode(body["content_b64"]))) as pdf:
        text = "\n".join((page.extract_text() or "") for page in pdf.pages)
    assert "Report" in text
    assert "a & b <c>" in text
    assert "second" in text
    assert "filename" not in body


def test_a_pdf_without_a_title_is_still_a_pdf():
    assert build_pdf_bytes("", "body").startswith(b"%PDF")


def test_the_cjk_font_is_registered_when_one_of_the_candidates_exists(monkeypatch):
    monkeypatch.setattr(worker_compute, "CJK_FONT_CANDIDATES", ["/nope/missing.ttf", VERA_TTF])

    assert build_pdf_bytes("제목", "본문").startswith(b"%PDF")


def test_a_candidate_that_is_not_a_font_falls_back_instead_of_failing(tmp_path, monkeypatch):
    impostor = tmp_path / "not-a-font.ttf"
    impostor.write_bytes(b"this is not a font")
    monkeypatch.setattr(worker_compute, "CJK_FONT_CANDIDATES", [str(impostor)])

    assert build_pdf_bytes("제목", "본문").startswith(b"%PDF")


def test_no_cjk_font_anywhere_still_renders(monkeypatch):
    """Every candidate absent — the loop runs out and Helvetica is used."""
    monkeypatch.setattr(worker_compute, "CJK_FONT_CANDIDATES", ["/definitely/not/here.ttf"])

    assert build_pdf_bytes("title", "body").startswith(b"%PDF")


def test_rows_that_are_not_a_list_of_lists_are_refused_by_the_builder():
    with pytest.raises(worker_compute.ToolError, match="list of lists"):
        build_xlsx_bytes("not rows")
    with pytest.raises(worker_compute.ToolError, match="list of lists"):
        build_xlsx_bytes([{"a": 1}])


@pytest.mark.parametrize(
    "kind,module,payload",
    [
        ("docx", "docx", {}),
        ("xlsx", "openpyxl", {"rows": []}),
        ("pptx", "pptx", {}),
        ("pdf", "reportlab.lib.pagesizes", {}),
    ],
)
def test_a_missing_document_library_is_a_503_naming_the_kind(monkeypatch, kind, module, payload):
    monkeypatch.setitem(sys.modules, module, None)
    client = _client()

    response = client.post(
        f"/worker/render/{kind}", json=payload, headers={LANGUAGE_HEADER: "en"}
    )

    assert response.status_code == 503
    assert kind in response.json()["detail"]
    assert "not installed" in response.json()["detail"]


def test_a_builder_that_breaks_is_a_500_naming_the_kind_and_the_reason(monkeypatch):
    def _boom(*args, **kwargs):
        raise RuntimeError("the renderer exploded")

    monkeypatch.setattr(worker_compute, "build_pdf_bytes", _boom)
    client = _client()

    response = client.post("/worker/render/pdf", json={}, headers={LANGUAGE_HEADER: "en"})

    assert response.status_code == 500
    assert "the renderer exploded" in response.json()["detail"]
    assert "pdf" in response.json()["detail"]


def test_the_render_reply_carries_no_target_of_its_own():
    """The one place a document's *name* is decided is ``lattice-agent``.

    ``documents.document_output_target`` (Rust) sanitises and resolves it
    before the render call; this seam answers with bytes and cost only, so
    there is no second name for the two sides to disagree about.
    """
    client = _client()

    body = client.post("/worker/render/docx", json={"filename": "../../etc/passwd"}).json()

    assert set(body) == {"content_b64", "bytes"}




def test_a_transcriber_that_hears_words_returns_them_with_its_identity():
    def transcribe(path: str) -> str:
        assert Path(path).read_bytes() == b"audio"
        return "  안녕하세요  "

    client = _client(transcriber=transcribe)

    body = client.post(
        "/worker/asr", json={"audio_b64": _b64(b"audio"), "mime": "audio/m4a"}
    ).json()

    assert body["text"] == "안녕하세요"
    assert body["status"] == "ok"
    assert body["segments"] is None
    assert body["provider"].endswith("transcribe")
    assert body["detail"] == ""


def test_a_worker_with_no_transcriber_reports_the_absence_rather_than_inventing():
    client = _client(transcriber=None)

    body = client.post("/worker/asr", json={"audio_b64": _b64(b"audio")}).json()

    assert body == {
        "text": "",
        "segments": None,
        "provider": "",
        "status": "unavailable",
        "detail": "no local transcriber is configured",
    }


def test_an_empty_transcript_is_not_reported_as_text():
    client = _client(transcriber=lambda path: "   ")

    body = client.post("/worker/asr", json={"audio_b64": _b64(b"a"), "mime": "audio/wav"}).json()

    assert body["status"] == "unavailable"
    assert body["detail"] == "the transcriber returned no text"
    assert body["text"] == ""


def test_a_broken_transcriber_is_a_reported_state_not_a_500():
    def _boom(path: str) -> str:
        raise RuntimeError("model crashed")

    client = _client(transcriber=_boom)

    response = client.post("/worker/asr", json={"audio_b64": _b64(b"a"), "filename": "m.mp3"})

    assert response.status_code == 200
    assert response.json()["status"] == "failed"
    assert response.json()["detail"] == "model crashed"


def test_a_container_no_transcriber_can_open_is_refused_before_any_work():
    called: List[str] = []
    client = _client(transcriber=lambda path: called.append(path) or "x")

    response = client.post(
        "/worker/asr",
        json={"audio_b64": _b64(b"a"), "filename": "memo.txt"},
        headers={LANGUAGE_HEADER: "en"},
    )

    assert response.status_code == 400
    assert ".txt" in response.json()["detail"]
    assert called == []


def test_a_recording_over_the_memo_ceiling_is_refused_with_both_numbers():
    client = _client(transcriber=lambda path: "x")
    oversized = b"0" * (MAX_AUDIO_BYTES + 1)

    response = client.post(
        "/worker/asr",
        json={"audio_b64": _b64(oversized), "mime": "audio/m4a"},
        headers={LANGUAGE_HEADER: "en"},
    )

    assert response.status_code == 413
    assert str(MAX_AUDIO_BYTES) in response.json()["detail"]


def test_audio_that_is_not_base64_is_a_422():
    client = _client()

    assert client.post("/worker/asr", json={"audio_b64": "!!!"}).status_code == 422


@pytest.mark.parametrize(
    "filename,mime,expected",
    [
        ("memo.WAV", "audio/m4a", ".wav"),
        (None, "audio/ogg; codecs=opus", ".ogg"),
        (None, None, ".m4a"),
        ("", "", ".m4a"),
    ],
)
def test_the_temp_suffix_prefers_the_filename_then_the_mime_then_the_default(
    filename, mime, expected
):
    assert worker_compute._suffix_for(filename, mime, ".m4a") == expected




def test_the_temp_file_is_gone_once_the_call_returns():
    seen: List[str] = []
    with worker_compute._temp_payload(b"payload", ".txt") as path:
        seen.append(path)
        assert Path(path).read_bytes() == b"payload"

    assert not Path(seen[0]).exists()


def test_a_temp_file_that_cannot_be_removed_does_not_fail_the_call(monkeypatch):
    """A read-only tmpdir is an operator's problem, not a failed transcription."""

    def _refuse(self, missing_ok=False):
        raise OSError("read-only file system")

    monkeypatch.setattr(Path, "unlink", _refuse)

    with worker_compute._temp_payload(b"x", ".txt") as path:
        leftover = path

    assert os.path.exists(leftover)
    monkeypatch.undo()
    os.remove(leftover)


def test_the_callable_identity_falls_back_to_the_type_for_an_object_port():
    class Port:
        def __call__(self, path: str) -> str:
            return "heard"

    client = _client(transcriber=Port())

    body = client.post("/worker/asr", json={"audio_b64": _b64(b"a"), "mime": "audio/wav"}).json()

    assert body["provider"].endswith(".Port")
    assert body["text"] == "heard"


def test_a_port_with_no_module_is_named_by_its_qualname_alone():
    def port(path: str) -> str:
        return ""

    port.__module__ = ""
    assert worker_compute._callable_identity(port) == port.__qualname__




def test_the_router_mounts_exactly_the_eight_compute_paths():
    from latticeai.runtime.build_phases.worker_profile import WORKER_COMPUTE_ROUTES

    # Read the router's own, always-flat route list rather than an
    # application's. FastAPI >= 0.140 stopped flattening `include_router` into
    # `app.routes` and appends a wrapper instead, where the usual
    # `isinstance(route, APIRoute)` walk finds nothing at all — so the app-side
    # spelling of this assertion compares an empty set and passes vacuously.
    router = create_worker_compute_router(
        embedder=None,
        require_user=lambda request: USER,
        enforce_rate_limit=lambda user, bucket: None,
    )

    keys = {(method, route.path) for route in router.routes for method in route.methods}
    assert keys == set(WORKER_COMPUTE_ROUTES)


def test_nothing_in_this_module_reaches_a_store_or_the_filesystem_it_keeps():
    """The compute seam's whole claim, checked against its own source."""
    source = Path(worker_compute.__file__).read_text(encoding="utf-8")
    for forbidden in ("sqlite3", "KnowledgeGraph", "ingestion_pipeline", "CONVERSATIONS"):
        assert forbidden not in source, f"{forbidden} has no business in a compute seam"


def test_the_request_bodies_are_the_ones_the_create_handlers_accepted():
    """The four bodies ``POST /tools/create_*`` took, pinned as literals.

    They used to be checked against ``latticeai.api.tools.Tool{Docx,Xlsx,Pptx,
    Pdf}Request`` — the request models of the Python routes these seams were
    factored out of. WP-W3b made those routes native and WP-P1 deleted the
    Python side, so the anchor is gone and the shape has to be written down.

    It is written down twice on purpose: the other copy is
    ``rust/lattice-agent/src/tools/render.rs``, which composes exactly these
    fields (``filename`` defaulted per creator, ``sheet_name`` defaulted to
    ``Sheet1``, ``rows``/``slides`` coerced from a JSON string). A field added
    here and not there is a render call that silently drops it.
    """
    expected = {
        worker_compute.RenderDocxRequest: {
            "title": "", "body": "", "filename": "document.docx",
        },
        worker_compute.RenderXlsxRequest: {
            "rows": [], "filename": "spreadsheet.xlsx", "sheet_name": "Sheet1",
        },
        worker_compute.RenderPptxRequest: {
            "title": "", "slides": [], "filename": "presentation.pptx",
        },
        worker_compute.RenderPdfRequest: {
            "title": "", "body": "", "filename": "document.pdf",
        },
    }
    for model, defaults in expected.items():
        assert set(model.model_fields) == set(defaults), model.__name__
        assert _defaults(model) == defaults, model.__name__


def _defaults(model) -> Dict[str, Any]:
    """Each field's effective default, factory-built ones included."""
    values: Dict[str, Any] = {}
    for name, field in model.model_fields.items():
        factory = field.default_factory
        values[name] = factory() if factory is not None else field.default
    return values


def test_the_seams_answer_json_and_never_a_file_path():
    """A response that named a path would be a write this seam does not own."""
    client = _client()
    docx = client.post("/worker/render/docx", json={"title": "t"}).json()
    xlsx = client.post("/worker/render/xlsx", json={"rows": []}).json()

    for body in (docx, xlsx):
        assert set(body) >= {"content_b64", "bytes"}
        assert "path" not in body
        # Not even a name: naming a target is the writer's job, and this seam
        # is not the writer.
        assert "filename" not in body




@pytest.fixture
def _rules_only(monkeypatch):
    """Pin extraction to the rule fallback so two calls in one test agree."""
    monkeypatch.setattr(
        "lattice_brain.graph._kg_common.extraction.ENABLE_LLM_EXTRACTION",
        False,
    )


def _python_extract(text: str, kind: str) -> Dict[str, Any]:
    """The literal composition ingest_* makes after calling the extractors."""
    from lattice_brain.graph._kg_common import (
        _classify_node_type,
        _extract_concepts,
        _extract_triples,
        _semantic_items,
    )

    limit = EXTRACT_LIMITS[kind]
    raw = _extract_concepts(text, limit=limit)
    return {
        "concepts": [
            {"text": concept, "node_type": _classify_node_type(concept, text)}
            for concept in raw
        ],
        "triples": [
            {
                "subject": triple["subject"],
                "object": triple["object"],
                "relation": triple["relation"],
                "weight": float(triple.get("weight") or 1.0),
                "context": str(triple.get("context") or ""),
                "evidence": str(triple.get("evidence") or ""),
                "confidence": triple.get("confidence"),
            }
            for triple in _extract_triples(text, raw)
        ],
        "semantic": [
            {
                "item_type": item["type"],
                "title": item["title"],
                "summary": item["summary"],
                "raw": dict(item),
            }
            for item in _semantic_items(text)
        ],
    }


def test_extract_returns_the_structures_ingestion_consumes(_rules_only):
    """Parity with the call sites ingest_message / ingest_document use."""
    client = _client()
    text = (
        "Lattice AI uses Graph RAG.\n"
        "We decided to keep the write path native.\n"
        "TODO: implement the extract seam."
    )
    body = client.post("/worker/extract", json={"text": text, "kind": "message"}).json()
    expected = _python_extract(text, "message")
    assert body == expected
    assert set(body) == {"concepts", "triples", "semantic"}
    assert any(item["raw"]["type"] == item["item_type"] for item in body["semantic"])


def test_extract_document_kind_uses_the_ingest_document_limit(monkeypatch):
    seen: List[int] = []

    def _fake_concepts(text: str, limit: int = 12) -> List[str]:
        seen.append(limit)
        return ["Lattice AI"]

    import lattice_brain.graph._kg_common as kg

    monkeypatch.setattr(kg, "_extract_concepts", _fake_concepts)
    client = _client()
    body = client.post(
        "/worker/extract", json={"text": "Lattice AI", "kind": "document"}
    ).json()
    assert seen == [EXTRACT_LIMITS["document"]]
    assert body["concepts"] == [{"text": "Lattice AI", "node_type": "Concept"}]


def test_extract_default_kind_is_the_chat_turn_door(_rules_only):
    client = _client()
    text = "FastAPI and Pydantic share a request body."
    omitted = client.post("/worker/extract", json={"text": text}).json()
    named = client.post("/worker/extract", json={"text": text, "kind": "message"}).json()
    assert omitted == named == _python_extract(text, "message")


def test_extract_empty_text_is_empty_structures_not_an_error():
    client = _client()
    body = client.post("/worker/extract", json={"text": ""}).json()
    assert body == {"concepts": [], "triples": [], "semantic": []}
    assert build_extract_reply("", "message") == body


def test_extract_rejects_an_unknown_kind():
    client = _client()
    response = client.post("/worker/extract", json={"text": "x", "kind": "query"})
    assert response.status_code == 422
    assert response.json()["detail"] == MESSAGES["worker_compute.extract_kind_invalid"][
        "ko"
    ].format(kind="query", allowed=", ".join(EXTRACT_KINDS))


def test_extract_kind_is_case_and_space_insensitive(_rules_only):
    client = _client()
    text = "Claude Sonnet writes the note."
    body = client.post(
        "/worker/extract", json={"text": text, "kind": " MESSAGE "}
    ).json()
    assert body == _python_extract(text, "message")


def test_extract_maps_a_thin_triple_the_way_the_write_path_does(monkeypatch):
    """``weight or 1.0`` and empty context/evidence are the ingest mapping."""
    import lattice_brain.graph._kg_common as kg

    monkeypatch.setattr(kg, "_extract_concepts", lambda text, limit=12: ["A", "B"])
    monkeypatch.setattr(
        kg,
        "_extract_triples",
        lambda text, concepts, limit=20: [
            {"subject": "A", "object": "B", "relation": "관련됨"}
        ],
    )
    monkeypatch.setattr(kg, "_semantic_items", lambda text: [])
    monkeypatch.setattr(kg, "_classify_node_type", lambda concept, text: "Concept")
    reply = build_extract_reply("A and B", "message")
    assert reply["triples"] == [
        {
            "subject": "A",
            "object": "B",
            "relation": "관련됨",
            "weight": 1.0,
            "context": "",
            "evidence": "",
            "confidence": None,
        }
    ]


def test_extract_kinds_are_the_two_ingest_doors():
    assert EXTRACT_KINDS == ("message", "document")
    assert EXTRACT_LIMITS == {"message": 12, "document": 15}
