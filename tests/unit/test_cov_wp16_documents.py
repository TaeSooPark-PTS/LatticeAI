"""wp16 coverage — ``lattice_brain.graph.documents``.

Two jobs live in this mixin: turning a parsed document's *structure* (slides,
pages, sheets, images) into graph nodes, and reading that structure out of the
real file formats. The tests build genuine .pptx/.pdf/.docx/.xlsx files in
``tmp_path`` with the same libraries the product parses them with, so a parser
change shows up here rather than in production.
"""

from __future__ import annotations

import sys
import zipfile
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from lattice_brain.graph.store import KnowledgeGraphStore


@pytest.fixture()
def store(tmp_path: Path) -> KnowledgeGraphStore:
    return KnowledgeGraphStore(tmp_path / "kg.sqlite", tmp_path / "blobs")


def _png(path: Path) -> Path:
    from PIL import Image

    Image.new("RGB", (12, 8), "blue").save(path)
    return path


def _pdf(path: Path) -> Path:
    from PIL import Image

    Image.new("RGB", (40, 20), "white").save(path, "PDF")
    return path


def _pptx(path: Path, image: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Lattice AI Roadmap"
    slide.shapes.add_picture(str(image), Inches(1), Inches(2), width=Inches(1))
    prs.save(str(path))
    return path


def _docx(path: Path) -> Path:
    from docx import Document

    doc = Document()
    doc.add_heading("Lattice AI Roadmap", level=1)
    doc.add_paragraph("Graph RAG keeps citations honest.")
    doc.add_paragraph("   ")  # blank paragraph: skipped, not counted
    doc.add_table(rows=1, cols=2)
    doc.save(str(path))
    return path


def _xlsx(path: Path) -> Path:
    from openpyxl import Workbook

    wb = Workbook()
    wb.active.title = "Q3"
    wb.active["A1"] = "revenue"
    wb.create_sheet("Q4")
    wb.save(str(path))
    return path


def _add_zip_entry(path: Path, name: str, data: bytes) -> None:
    kept = []
    with zipfile.ZipFile(path) as source:
        for info in source.infolist():
            kept.append((info, source.read(info.filename)))
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as target:
        for info, payload in kept:
            target.writestr(info, payload)
        target.writestr(name, data)


# ── find_documents_by_uri_prefix ─────────────────────────────────────────────


def test_find_documents_by_uri_prefix_requires_a_prefix(store) -> None:
    assert store.find_documents_by_uri_prefix("") == []
    assert store.find_documents_by_uri_prefix("   ") == []
    assert store.find_documents_by_uri_prefix(None) == []


# ── delete_document_tree ─────────────────────────────────────────────────────


def test_delete_document_tree_skips_blank_ids(store) -> None:
    assert store.delete_document_tree("") == {"status": "skipped", "removed_nodes": 0}


def test_delete_document_tree_reports_missing_nodes(store) -> None:
    assert store.delete_document_tree("doc:ghost") == {
        "status": "not_found",
        "node_id": "doc:ghost",
        "removed_nodes": 0,
    }


def test_delete_document_tree_keeps_a_source_shared_with_another_document(
    store,
) -> None:
    with store._connect() as conn:
        store._upsert_node(conn, "doc:a", "Document", "A")
        store._upsert_node(conn, "doc:b", "Document", "B")
        store._upsert_node(conn, "source:shared", "Source", "shared source")
        store._upsert_edge(conn, "doc:a", "source:shared", "indexed_from")
        store._upsert_edge(conn, "doc:b", "source:shared", "indexed_from")

    result = store.delete_document_tree("doc:a")

    assert result == {"status": "ok", "node_id": "doc:a", "removed_nodes": 1}
    with store._connect() as conn:
        remaining = {
            row["id"] for row in conn.execute("SELECT id FROM nodes").fetchall()
        }
    assert "source:shared" in remaining
    assert "doc:a" not in remaining

    # Once the last document goes, the now-orphaned Source goes with it.
    assert store.delete_document_tree("doc:b")["removed_nodes"] == 2
    with store._connect() as conn:
        assert conn.execute("SELECT COUNT(*) AS n FROM nodes").fetchone()["n"] == 0


# ── _ingest_structure_nodes ──────────────────────────────────────────────────


STRUCTURE = {
    "slides": [{"index": 1, "texts": ["Lattice AI Roadmap covers Graph RAG."]}],
    "pages": [{"index": 1, "preview": "Lattice AI Roadmap covers Graph RAG."}],
    "sheets": [{"title": "Q3", "max_row": 4, "max_column": 2}],
    "images": [
        {"sha256": "abc123def456", "page": 1, "name": "ppt/media/image1.png"},
        {"bytes": 42},
    ],
}


@pytest.mark.parametrize("workspace_id", [None, "org:acme"])
def test_ingest_structure_nodes_projects_every_structural_part(
    store, workspace_id
) -> None:
    with store._connect() as conn:
        store._upsert_node(conn, "file:deck", "SlideDeck", "deck.pptx")
        store._ingest_structure_nodes(
            conn,
            "file:deck",
            "deck.pptx",
            STRUCTURE,
            owner="owner@example.com",
            workspace_id=workspace_id,
        )
        types = {
            row["type"]: row["n"]
            for row in conn.execute(
                "SELECT type, COUNT(*) AS n FROM nodes GROUP BY type"
            ).fetchall()
        }
        edges = {
            row["type"]
            for row in conn.execute("SELECT DISTINCT type FROM edges").fetchall()
        }
        titles = [
            row["title"]
            for row in conn.execute(
                "SELECT title FROM nodes WHERE type='Image' ORDER BY id"
            ).fetchall()
        ]

    assert types["Slide"] == 1
    assert types["Page"] == 1
    assert types["Sheet"] == 1
    assert types["Image"] == 2
    assert types.get("Topic", 0) >= 1
    # Verb-form edges are the graph's vocabulary; the labels are normalized.
    assert {"HAS_SLIDE", "HAS_PAGE", "HAS_SHEET", "CONTAINS_IMAGE", "DISCUSSES"} <= edges
    # A named image carries page + filename in its title; an anonymous one does not.
    assert "deck.pptx / image / page 1 / image1.png" in titles
    assert "deck.pptx / image" in titles


# ── _document_structure dispatch ─────────────────────────────────────────────


def test_document_structure_dispatches_by_extension(store, tmp_path: Path) -> None:
    image = _png(tmp_path / "pic.png")

    slides = store._document_structure(_pptx(tmp_path / "deck.pptx", image), ".pptx")
    assert [slide["index"] for slide in slides["slides"]] == [1]

    pages = store._document_structure(_pdf(tmp_path / "doc.pdf"), ".pdf")
    assert [page["index"] for page in pages["pages"]] == [1]

    words = store._document_structure(_docx(tmp_path / "doc.docx"), ".docx")
    assert words["paragraphs"] == 2

    sheets = store._document_structure(_xlsx(tmp_path / "book.xlsx"), ".xlsx")
    assert [sheet["title"] for sheet in sheets["sheets"]] == ["Q3", "Q4"]

    assert store._document_structure(tmp_path / "pic.png", ".png") == {}


def test_document_structure_reports_parser_failures(store, tmp_path: Path) -> None:
    broken = tmp_path / "not-really.docx"
    broken.write_text("plain text pretending to be a document", encoding="utf-8")

    result = store._document_structure(broken, ".docx")

    assert "error" in result and result["error"]


# ── format-specific structure readers ────────────────────────────────────────


def test_pptx_structure_reads_shapes_text_and_media(store, tmp_path: Path) -> None:
    deck = _pptx(tmp_path / "deck.pptx", _png(tmp_path / "pic.png"))
    _add_zip_entry(deck, "ppt/media/image99.png", b"not a real png")

    result = store._pptx_structure(deck)

    assert "error" not in result
    slide = result["slides"][0]
    assert slide["index"] == 1
    assert "Lattice AI Roadmap" in slide["texts"]
    text_shapes = [shape for shape in slide["shapes"] if "text" in shape]
    assert text_shapes[0]["bbox"]["width"] > 0
    assert len(slide["shapes"]) == 2  # title + picture
    images = {image["name"]: image for image in result["images"]}
    assert images["ppt/media/image1.png"]["width"] == 12
    assert images["ppt/media/image1.png"]["format"] == "PNG"
    assert images["ppt/media/image1.png"]["sha256"]
    # The unreadable media part is still reported, just without dimensions.
    assert "width" not in images["ppt/media/image99.png"]


def test_pptx_structure_reports_failure_for_non_pptx(store, tmp_path: Path) -> None:
    broken = tmp_path / "deck.pptx"
    broken.write_text("nope", encoding="utf-8")

    result = store._pptx_structure(broken)

    assert result["error"]
    assert result["slides"] == []


def test_pdf_structure_reads_pages_and_images(store, tmp_path: Path) -> None:
    result = store._pdf_structure(_pdf(tmp_path / "doc.pdf"))

    assert "error" not in result
    page = result["pages"][0]
    assert page["index"] == 1
    assert page["width"] == 40.0
    assert page["height"] == 20.0
    assert page["chars"] == 0
    assert page["image_count"] == 1
    assert result["images"][0]["page"] == 1
    assert result["images"][0]["index"] == 1
    assert set(result["images"][0]["bbox"]) == {"x0", "top", "x1", "bottom"}
    assert isinstance(result["metadata"], dict)


def test_pdf_structure_reports_failure_for_non_pdf(store, tmp_path: Path) -> None:
    broken = tmp_path / "doc.pdf"
    broken.write_text("nope", encoding="utf-8")

    result = store._pdf_structure(broken)

    assert result["error"]
    assert result["pages"] == []


def test_docx_structure_counts_paragraphs_headings_and_tables(
    store, tmp_path: Path
) -> None:
    result = store._docx_structure(_docx(tmp_path / "doc.docx"))

    assert result["paragraphs"] == 2
    assert result["tables"] == 1
    assert [heading["text"] for heading in result["headings"]] == [
        "Lattice AI Roadmap"
    ]
    assert result["headings"][0]["style"].lower().startswith("heading")


def test_xlsx_structure_lists_every_worksheet(store, tmp_path: Path) -> None:
    result = store._xlsx_structure(_xlsx(tmp_path / "book.xlsx"))

    assert [sheet["title"] for sheet in result["sheets"]] == ["Q3", "Q4"]
    assert result["sheets"][0]["max_row"] == 1
    assert result["sheets"][0]["max_column"] == 1
