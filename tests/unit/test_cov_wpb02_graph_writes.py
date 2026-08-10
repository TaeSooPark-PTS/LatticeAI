"""wpb02 branch coverage — the Knowledge Graph write side.

The guards driven here are the ones that only fire on data the happy path never
produces: a relation triple naming a concept that was never promoted to a node,
a PDF whose page table starts after the text does, an edge whose caller already
supplied the canonical label as its legacy label, a build with the v2 schema
module absent, an import artifact carrying an edge with no type at all, and a
projection whose freshly written row is gone by the time it is read back.

Everything runs against a real :class:`KnowledgeGraphStore` on ``tmp_path``;
only pure module-level helpers (``_extract_triples``, ``pdf_page_offsets``,
``EdgeType``, ``KGStoreV2``) are swapped, and always through ``monkeypatch``.
"""

from __future__ import annotations

import sqlite3
from pathlib import Path
from typing import Any, Dict, List

import pytest

from lattice_brain.graph import ingest as ingest_mod
from lattice_brain.graph import projection as projection_mod
from lattice_brain.graph import provenance as provenance_mod
from lattice_brain.graph import write_master as write_master_mod
from lattice_brain.graph.store import KnowledgeGraphStore

ORPHAN_TRIPLE = [
    {"subject": "ghost-subject", "object": "ghost-object", "relation": "관련됨", "weight": 0.5}
]


@pytest.fixture()
def store(tmp_path: Path) -> KnowledgeGraphStore:
    return KnowledgeGraphStore(tmp_path / "kg.sqlite", tmp_path / "blobs")


def _edge_types(store: KnowledgeGraphStore) -> List[str]:
    with store._connect() as conn:
        return [row["type"] for row in conn.execute("SELECT type FROM edges ORDER BY id ASC")]


# ── ingest.py: triples whose endpoints never became concept nodes ───────────


def test_a_message_triple_without_concept_nodes_creates_no_edge(
    store: KnowledgeGraphStore, monkeypatch
):
    monkeypatch.setattr(ingest_mod, "_extract_triples", lambda *_a, **_k: list(ORPHAN_TRIPLE))

    store.ingest_message("user", "Vector search and graph recall work together.")

    assert "관련됨" not in _edge_types(store)


def test_a_document_triple_without_concept_nodes_creates_no_edge(
    store: KnowledgeGraphStore, monkeypatch, tmp_path: Path
):
    monkeypatch.setattr(ingest_mod, "_extract_triples", lambda *_a, **_k: list(ORPHAN_TRIPLE))
    source = tmp_path / "notes.txt"
    source.write_text("Vector search and graph recall work together.", encoding="utf-8")

    store.ingest_document(
        source,
        extracted={"content": "Vector search and graph recall work together."},
    )

    assert "관련됨" not in _edge_types(store)


def test_a_source_triple_without_concept_nodes_creates_no_edge(
    store: KnowledgeGraphStore, monkeypatch
):
    monkeypatch.setattr(ingest_mod, "_extract_triples", lambda *_a, **_k: list(ORPHAN_TRIPLE))

    store.ingest_source(
        source_type="note",
        title="Retrieval notes",
        text="Vector search and graph recall work together.",
    )

    assert "관련됨" not in _edge_types(store)


def test_a_chunk_that_starts_before_the_first_page_offset_gets_no_page_label(
    store: KnowledgeGraphStore, monkeypatch, tmp_path: Path
):
    # A page table whose first page starts after byte 0 cannot label the
    # opening chunk; the ingest records no page rather than a wrong one.
    monkeypatch.setattr(ingest_mod, "pdf_page_offsets", lambda _meta: [40, 120])
    source = tmp_path / "report.pdf"
    source.write_bytes(b"%PDF-1.4 not a real pdf")
    text = "Durable recall matters. " * 20

    store.ingest_document(source, extracted={"content": text})

    with store._connect() as conn:
        rows = conn.execute(
            "SELECT metadata_json FROM nodes WHERE type='Chunk' ORDER BY id ASC"
        ).fetchall()
    assert rows
    assert all('"page"' not in row["metadata_json"] for row in rows)


# ── write_master.py: the single write door ──────────────────────────────────


def test_an_edge_whose_legacy_label_is_already_canonical_records_no_legacy_label(
    store: KnowledgeGraphStore,
):
    with store._connect() as conn:
        store._upsert_node(conn, "n:a", "Concept", "A")
        store._upsert_node(conn, "n:b", "Concept", "B")
        store._upsert_edge(conn, "n:a", "n:b", "포함함", legacy_label="CONTAINS")
        row = conn.execute(
            "SELECT type, metadata_json FROM edges WHERE from_node='n:a'"
        ).fetchone()

    assert row["type"] == "CONTAINS"
    assert "legacy_label" not in row["metadata_json"]


def test_without_the_v2_schema_module_an_edge_keeps_its_free_string_type(
    store: KnowledgeGraphStore, monkeypatch
):
    monkeypatch.setattr(write_master_mod, "EdgeType", None)
    with store._connect() as conn:
        store._upsert_node(conn, "n:a", "Concept", "A")
        store._upsert_node(conn, "n:b", "Concept", "B")
        store._upsert_edge(conn, "n:a", "n:b", "포함함")
        row = conn.execute(
            "SELECT type FROM edges WHERE from_node='n:a'"
        ).fetchone()

    assert row["type"] == "포함함"


# ── provenance.py: import ───────────────────────────────────────────────────


def _artifact(edges: List[Dict[str, Any]]) -> Dict[str, Any]:
    return {
        "header": {"graph_schema_version": 1},
        "nodes": [
            {"id": "n:a", "type": "Concept", "title": "A", "metadata_json": "{}", "raw_json": "{}"},
            {"id": "n:b", "type": "Concept", "title": "B", "metadata_json": "{}", "raw_json": "{}"},
        ],
        "edges": edges,
        "chunks": [],
        "knowledge_sources": [],
        "provenance": [],
    }


def test_a_replace_import_without_the_v2_schema_module_still_rebuilds_the_graph(
    store: KnowledgeGraphStore, monkeypatch
):
    monkeypatch.setattr(provenance_mod, "KGStoreV2", None)

    result = store.import_graph_data(_artifact([]), mode="replace")

    assert result["nodes"] == 2
    with store._connect() as conn:
        assert conn.execute("SELECT COUNT(*) AS c FROM nodes").fetchone()["c"] == 2


def test_an_imported_edge_with_no_type_invents_no_legacy_label(store: KnowledgeGraphStore):
    artifact = _artifact(
        [{"from_node": "n:a", "to_node": "n:b", "type": "", "weight": 1.0, "metadata_json": "{}"}]
    )

    result = store.import_graph_data(artifact, mode="merge")

    assert result["edges"] == 1
    with store._connect() as conn:
        row = conn.execute(
            "SELECT type, metadata_json FROM edges WHERE from_node='n:a'"
        ).fetchone()
    # The importer passes no legacy label of its own; the write door normalizes
    # the empty type to MENTIONS and records the empty original verbatim.
    assert row["type"] == "MENTIONS"
    assert row["metadata_json"] == '{"legacy_label": ""}'


# ── projection.py: the row is gone before it can be read back ───────────────


class _PurgingConnection(sqlite3.Connection):
    """Another writer purges ``edges_v2`` between the upsert and the read-back."""

    def execute(self, sql, *args):  # type: ignore[override]
        if sql.lstrip().startswith("SELECT id FROM edges_v2"):
            sqlite3.Connection.execute(self, "DELETE FROM edges_v2")
        return sqlite3.Connection.execute(self, sql, *args)


def test_a_projected_edge_that_vanishes_records_no_occurrence(store: KnowledgeGraphStore):
    conn = sqlite3.connect(str(store.db_path), factory=_PurgingConnection)
    conn.row_factory = sqlite3.Row
    try:
        store._v2_project_edge(conn, "n:a", "n:b", "CONTAINS", 1.0, "{}", strict=True)
        occurrences = conn.execute(
            "SELECT COUNT(*) AS c FROM edge_occurrences"
        ).fetchone()["c"]
    finally:
        conn.close()

    assert occurrences == 0


# ── documents.py: a slide placeholder with an empty text frame ──────────────


def test_an_empty_slide_placeholder_contributes_no_text(
    store: KnowledgeGraphStore, tmp_path: Path
):
    from pptx import Presentation

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])  # "Title Only", title left blank
    deck = tmp_path / "blank.pptx"
    prs.save(str(deck))

    structure = store._document_structure(deck, ".pptx")

    assert structure["slides"][0]["texts"] == []
    assert structure["slides"][0]["shapes"]
    assert "text" not in structure["slides"][0]["shapes"][0]


def test_projection_module_exposes_the_v2_edge_projector():
    assert hasattr(projection_mod.KnowledgeGraphProjectionMixin, "_v2_project_edge")
