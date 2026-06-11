"""Document generation (docx/xlsx/pptx/pdf) and extraction (read_document)."""

from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List

from tools import (
    ToolError,
    _resolve_path,
    _relative,
    DOCUMENT_OUTPUT_DIR,
    PRESENTATION_OUTPUT_DIR,
    SPREADSHEET_OUTPUT_DIR,
    PDF_OUTPUT_DIR,
    DOCUMENT_MAX_READ_BYTES,
    _CJK_FONT_CANDIDATES,
    _SUPPORTED_READ_EXTENSIONS,
)


def _safe_filename(name: str, suffix: str) -> str:
    base = Path(name or f"artifact{suffix}").name
    if not base.lower().endswith(suffix):
        base += suffix
    safe = "".join(ch if ch.isalnum() or ch in ("-", "_", ".", " ") else "_" for ch in base).strip()
    return safe or f"artifact{suffix}"


def _body_to_str(body) -> str:
    if isinstance(body, list):
        return "\n\n".join(str(item) for item in body)
    return str(body or "")


def create_docx(title: str, body, filename: str = "document.docx") -> Dict[str, Any]:
    try:
        from docx import Document
    except Exception as exc:
        raise ToolError("python-docx is not installed. Run `pip install -r requirements.txt`.") from exc

    output_dir = _resolve_path(DOCUMENT_OUTPUT_DIR)
    output_dir.mkdir(parents=True, exist_ok=True)
    target = output_dir / _safe_filename(filename, ".docx")

    document = Document()
    if title:
        document.add_heading(str(title), level=1)
    for block in _body_to_str(body).split("\n\n"):
        text = block.strip()
        if text:
            document.add_paragraph(text)
    document.save(target)
    return {"path": _relative(target), "bytes": target.stat().st_size}


def create_xlsx(rows: List[List[Any]], filename: str = "spreadsheet.xlsx", sheet_name: str = "Sheet1") -> Dict[str, Any]:
    try:
        from openpyxl import Workbook
    except Exception as exc:
        raise ToolError("openpyxl is not installed. Run `pip install -r requirements.txt`.") from exc

    if not isinstance(rows, list) or not all(isinstance(row, list) for row in rows):
        raise ToolError("Rows must be a list of lists.")

    output_dir = _resolve_path(SPREADSHEET_OUTPUT_DIR)
    output_dir.mkdir(parents=True, exist_ok=True)
    target = output_dir / _safe_filename(filename, ".xlsx")

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = (sheet_name or "Sheet1")[:31]
    for row in rows:
        sheet.append(row)
    workbook.save(target)
    return {"path": _relative(target), "rows": len(rows), "bytes": target.stat().st_size}


def create_pptx(title: str, slides: List[Dict[str, Any]], filename: str = "presentation.pptx") -> Dict[str, Any]:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise ToolError("python-pptx is not installed. Run `pip install -r requirements.txt`.") from exc

    output_dir = _resolve_path(PRESENTATION_OUTPUT_DIR)
    output_dir.mkdir(parents=True, exist_ok=True)
    target = output_dir / _safe_filename(filename, ".pptx")

    presentation = Presentation()
    first_layout = presentation.slide_layouts[0]
    first = presentation.slides.add_slide(first_layout)
    first.shapes.title.text = title or "Presentation"
    first.placeholders[1].text = ""

    content_layout = presentation.slide_layouts[1]
    for slide_data in slides or []:
        slide = presentation.slides.add_slide(content_layout)
        slide.shapes.title.text = str(slide_data.get("title") or "Slide")
        body = slide.placeholders[1].text_frame
        body.clear()
        bullets = slide_data.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        for index, bullet in enumerate(bullets):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = str(bullet)
            paragraph.level = 0

    presentation.save(target)
    return {"path": _relative(target), "slides": len(presentation.slides), "bytes": target.stat().st_size}




def create_pdf(title: str, body, filename: str = "document.pdf") -> Dict[str, Any]:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle
        from reportlab.lib.units import mm
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
    except Exception as exc:
        raise ToolError("reportlab is not installed. Run `pip install reportlab`.") from exc

    output_dir = _resolve_path(PDF_OUTPUT_DIR)
    output_dir.mkdir(parents=True, exist_ok=True)
    target = output_dir / _safe_filename(filename, ".pdf")

    # CJK 폰트 등록
    font_name = "Helvetica"
    for font_path in _CJK_FONT_CANDIDATES:
        if Path(font_path).exists():
            try:
                pdfmetrics.registerFont(TTFont("KoreanFont", font_path))
                font_name = "KoreanFont"
            except Exception:
                pass
            break

    title_style = ParagraphStyle("Title", fontName=font_name, fontSize=18, spaceAfter=8, leading=24)
    body_style  = ParagraphStyle("Body",  fontName=font_name, fontSize=11, spaceAfter=6, leading=16)

    story = []
    if title:
        story.append(Paragraph(str(title), title_style))
        story.append(Spacer(1, 4 * mm))

    for block in _body_to_str(body).split("\n\n"):
        text = block.strip()
        if text:
            safe_text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe_text, body_style))
            story.append(Spacer(1, 2 * mm))

    doc = SimpleDocTemplate(str(target), pagesize=A4,
                             leftMargin=20*mm, rightMargin=20*mm,
                             topMargin=20*mm, bottomMargin=20*mm)
    doc.build(story)
    return {"path": _relative(target), "bytes": target.stat().st_size}


def read_document(path: str) -> Dict[str, Any]:
    """Extract text from PDF, DOCX, XLSX, PPTX, TXT, MD, CSV files."""
    target = Path(path).expanduser().resolve()
    if not target.exists():
        raise ToolError(f"파일이 없습니다: {path}")
    if not target.is_file():
        raise ToolError(f"파일이 아닙니다: {path}")
    if target.stat().st_size > DOCUMENT_MAX_READ_BYTES:
        raise ToolError(f"파일이 너무 큽니다 ({target.stat().st_size:,} bytes).")

    ext = target.suffix.lower()
    if ext not in _SUPPORTED_READ_EXTENSIONS:
        raise ToolError(f"지원하지 않는 형식입니다: {ext}. 지원: {', '.join(_SUPPORTED_READ_EXTENSIONS)}")

    text = ""
    meta: Dict[str, Any] = {"path": str(target), "ext": ext}

    if ext == ".pdf":
        try:
            import pdfplumber
            with pdfplumber.open(str(target)) as pdf:
                meta["pages"] = len(pdf.pages)
                text = "\n\n".join(
                    (p.extract_text() or "") for p in pdf.pages
                ).strip()
        except Exception as exc:
            raise ToolError(f"PDF 읽기 실패: {exc}") from exc

    elif ext == ".docx":
        try:
            from docx import Document
            doc = Document(str(target))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            text = "\n\n".join(paragraphs)
            meta["paragraphs"] = len(paragraphs)
        except Exception as exc:
            raise ToolError(f"DOCX 읽기 실패: {exc}") from exc

    elif ext == ".xlsx":
        try:
            from openpyxl import load_workbook
            wb = load_workbook(str(target), data_only=True)
            rows_all = []
            for ws in wb.worksheets:
                rows_all.append(f"[Sheet: {ws.title}]")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    rows_all.append("\t".join(cells))
            text = "\n".join(rows_all)
            meta["sheets"] = len(wb.worksheets)
        except Exception as exc:
            raise ToolError(f"XLSX 읽기 실패: {exc}") from exc

    elif ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(target))
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                parts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
                slides_text.append(f"[Slide {i}]\n" + "\n".join(parts))
            text = "\n\n".join(slides_text)
            meta["slides"] = len(prs.slides)
        except Exception as exc:
            raise ToolError(f"PPTX 읽기 실패: {exc}") from exc

    elif ext in {".txt", ".md", ".csv"}:
        try:
            text = target.read_text(encoding="utf-8", errors="replace")
        except Exception as exc:
            raise ToolError(f"파일 읽기 실패: {exc}") from exc

    meta["chars"] = len(text)
    meta["preview"] = text[:500]
    meta["content"] = text[:50_000]   # 50K char cap for context
    return meta
